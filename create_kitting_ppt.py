#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
키팅 자동화 로직 구현방법 PPT 생성기
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ───────────────────────────────────────────
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
BLUE       = RGBColor(0x29, 0x52, 0xa3)
GREEN      = RGBColor(0x1a, 0x8a, 0x4a)
GREEN_LT   = RGBColor(0xd4, 0xed, 0xda)
RED        = RGBColor(0xc0, 0x39, 0x2b)
RED_LT     = RGBColor(0xfb, 0xea, 0xe8)
ORANGE     = RGBColor(0xe6, 0x7e, 0x22)
ORANGE_LT  = RGBColor(0xfd, 0xf3, 0xe3)
PURPLE     = RGBColor(0x6c, 0x35, 0x9e)
PURPLE_LT  = RGBColor(0xed, 0xe3, 0xf7)
TEAL       = RGBColor(0x0d, 0x86, 0x8c)
TEAL_LT    = RGBColor(0xd1, 0xf0, 0xf2)
BLUE_LT    = RGBColor(0xd6, 0xe4, 0xf7)
GRAY_BG    = RGBColor(0xf4, 0xf6, 0xf9)
WHITE      = RGBColor(0xff, 0xff, 0xff)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
BORDER     = RGBColor(0xcc, 0xd6, 0xe5)
CODE_BG    = RGBColor(0x1e, 0x1e, 0x2e)
CODE_FG    = RGBColor(0xa6, 0xe2, 0x2e)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=Pt(12), bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = size
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return txb

def code_box(slide, l, t, w, h, lines):
    box(slide, l, t, w, h, fill_color=CODE_BG)
    txb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.1),
                                    w - Inches(0.3), h - Inches(0.2))
    txb.word_wrap = False
    tf = txb.text_frame; tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line; run.font.size = Pt(9.5)
        run.font.bold = False
        if line.startswith('#'):
            run.font.color.rgb = RGBColor(0x75, 0x71, 0x5e)
        elif any(k in line for k in ['def ', 'import ', 'from ', 'class ']):
            run.font.color.rgb = RGBColor(0x66, 0xd9, 0xef)
        elif line.strip().startswith(('log(', 'time.', 'raise', 'return')):
            run.font.color.rgb = CODE_FG
        else:
            run.font.color.rgb = RGBColor(0xf8, 0xf8, 0xf2)
        run.font.name = 'Consolas'

def header_bar(slide, title, subtitle=""):
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    txt(slide, title, Inches(0.4), Inches(0.1), Inches(10), Inches(0.6),
        size=Pt(25), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.42), Inches(0.68), Inches(11), Inches(0.35),
            size=Pt(11.5), color=RGBColor(0xb0, 0xc4, 0xde))
    txt(slide, "키팅 자동화 시스템",
        Inches(9.5), Inches(0.15), Inches(3.5), Inches(0.4),
        size=Pt(10), color=RGBColor(0x80, 0xa0, 0xc8), align=PP_ALIGN.RIGHT)

def tag(slide, text, l, t, color):
    w = Inches(len(text) * 0.09 + 0.3)
    box(slide, l, t, w, Inches(0.3), fill_color=color)
    txt(slide, text, l + Inches(0.08), t + Inches(0.03),
        w - Inches(0.1), Inches(0.26), size=Pt(9.5), bold=True, color=WHITE)
    return w


# ════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════
def slide_cover(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    for ox, oy, ow, col in [
        (Inches(10.2), Inches(-1.0), Inches(6), RGBColor(0x25, 0x4e, 0x7e)),
        (Inches(-1.5), Inches(4.8),  Inches(5), RGBColor(0x18, 0x30, 0x50)),
        (Inches(5.5),  Inches(5.8),  Inches(3), RGBColor(0x1e, 0x45, 0x72)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    box(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H,
        fill_color=BLUE)

    txt(s, "sMES 키팅 자동화",
        Inches(0.55), Inches(1.5), Inches(12), Inches(0.85),
        size=Pt(44), bold=True, color=WHITE)
    txt(s, "로직 구현 방법 설명서",
        Inches(0.6), Inches(2.38), Inches(10), Inches(0.65),
        size=Pt(28), color=RGBColor(0x90, 0xb8, 0xe8))
    txt(s, "pywinauto · pyautogui · Playwright · Win32 API 연동",
        Inches(0.6), Inches(3.1), Inches(10), Inches(0.45),
        size=Pt(14), color=RGBColor(0xb0, 0xc8, 0xe0), italic=True)

    box(s, Inches(0.6), Inches(3.65), Inches(5.5), Pt(1.5), fill_color=BLUE)

    items = [
        "Step 1  sMES.exe 자동 실행 & 관리자 권한 획득",
        "Step 2  pywinauto 창 연결 (win32 / uia 백엔드)",
        "Step 3  로그인 자동화 (Edit 컨트롤 탐지 + 클립보드 입력)",
        "Step 4  메뉴 탐색 — 생산관리 > 조립 자재 kitting",
        "Step 5  날짜 설정 + 조회 + 품목별 Excel 다운로드",
        "Step 6  Playwright — 자재부족현황 웹앱 자동 업로드",
        "Step 7  오류 처리 & 폴백 전략",
    ]
    ty = Inches(3.85)
    for item in items:
        txt(s, item, Inches(0.7), ty, Inches(9), Inches(0.34),
            size=Pt(11.5), color=RGBColor(0xc0, 0xd8, 0xf0))
        ty += Inches(0.36)


# ════════════════════════════════════════════════════════════
# 슬라이드 2: 전체 흐름 개요
# ════════════════════════════════════════════════════════════
def slide_overview(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "전체 자동화 흐름 개요",
               "main() 진입점부터 완료 팝업까지 — 7단계 순차 실행")

    stages = [
        ("🔐", "관리자\n권한 확인",  RED,    RED_LT,   "is_admin()\nUAC 권한 상승"),
        ("🚀", "sMES\n실행",         BLUE,   BLUE_LT,  "subprocess.Popen\n(sMES.exe)"),
        ("🔗", "창 연결",            TEAL,   TEAL_LT,  "pywinauto\nwin32/uia"),
        ("🔑", "로그인",             PURPLE, PURPLE_LT,"Edit 컨트롤 탐지\n클립보드 입력"),
        ("📋", "메뉴+\n다운로드",   ORANGE, ORANGE_LT,"생산관리>\nkitting>조회>Excel"),
        ("🌐", "웹 업로드",          GREEN,  GREEN_LT, "Playwright\n자동 파일 업로드"),
        ("✅", "완료 팝업",          NAVY,   BLUE_LT,  "page.evaluate()\n팝업 표시"),
    ]

    bw = Inches(1.73)
    bh = Inches(3.2)
    sx = Inches(0.22)
    sy = Inches(1.3)

    for i, (icon, title, c, bg_c, desc) in enumerate(stages):
        lx = sx + i * (bw + Inches(0.07))
        box(s, lx, sy, bw, bh, fill_color=bg_c, line_color=c, line_width=Pt(1.5))
        box(s, lx, sy, bw, Inches(0.55), fill_color=c)
        txt(s, icon, lx, sy + Inches(0.06), bw, Inches(0.38),
            size=Pt(20), color=WHITE, align=PP_ALIGN.CENTER)
        for j, line in enumerate(title.split('\n')):
            txt(s, line, lx, sy + Inches(0.58 + j*0.28), bw, Inches(0.28),
                size=Pt(12), bold=True, color=c, align=PP_ALIGN.CENTER)
        box(s, lx+Inches(0.12), sy+Inches(1.15), bw-Inches(0.24), Pt(0.8),
            fill_color=BORDER)
        for j, line in enumerate(desc.split('\n')):
            txt(s, line, lx+Inches(0.1), sy+Inches(1.3+j*0.32), bw-Inches(0.2), Inches(0.3),
                size=Pt(10), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        # 단계 번호
        txt(s, f"Step {i+1 if i>0 else '0'}", lx, sy+Inches(2.85), bw, Inches(0.25),
            size=Pt(9), color=c, align=PP_ALIGN.CENTER, bold=True)
        if i < len(stages)-1:
            txt(s, "→", lx+bw, sy+Inches(1.35), Inches(0.1), Inches(0.38),
                size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 하단 라이브러리 태그
    box(s, Inches(0.22), Inches(4.7), Inches(12.9), Inches(1.5),
        fill_color=WHITE, line_color=BORDER, line_width=Pt(1))
    txt(s, "사용 라이브러리", Inches(0.4), Inches(4.78),
        Inches(2), Inches(0.32), size=Pt(12), bold=True, color=NAVY)
    libs = [
        ("pywinauto",    BLUE,   "Windows GUI 컨트롤 조작"),
        ("pyautogui",    TEAL,   "마우스·키보드 시뮬레이션"),
        ("Playwright",   GREEN,  "Chromium 브라우저 자동화"),
        ("psutil",       PURPLE, "프로세스 PID 탐색"),
        ("win32clipboard", ORANGE, "클립보드 텍스트 입력"),
        ("subprocess",   RED,    "외부 프로세스 실행"),
    ]
    lx = Inches(0.4)
    for lib, col, desc in libs:
        w = tag(s, lib, lx, Inches(5.18), col)
        txt(s, desc, lx, Inches(5.52), w + Inches(0.1), Inches(0.28),
            size=Pt(9), color=MID_GRAY, align=PP_ALIGN.CENTER)
        lx += w + Inches(0.25)


# ════════════════════════════════════════════════════════════
# 슬라이드 3: Step 1~2 — sMES 실행 & 창 연결
# ════════════════════════════════════════════════════════════
def slide_launch(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "Step 1~2  sMES 실행 & 창 연결",
               "관리자 권한 확인 → subprocess 실행 → pywinauto 창 연결 (백엔드 이중화)")

    # 좌: 관리자 권한
    box(s, Inches(0.3), Inches(1.25), Inches(4.0), Inches(5.3),
        fill_color=WHITE, line_color=RED, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(4.0), Inches(0.42), fill_color=RED)
    txt(s, "🔐  관리자 권한 확인", Inches(0.45), Inches(1.27),
        Inches(3.7), Inches(0.38), size=Pt(13), bold=True, color=WHITE)
    code_box(s, Inches(0.38), Inches(1.78), Inches(3.84), Inches(1.9), [
        "def is_admin():",
        "    return ctypes.windll",
        "        .shell32",
        "        .IsUserAnAdmin()",
        "",
        "if not is_admin():",
        "    elevate()  # UAC 재실행",
    ])
    txt(s, "• ctypes Win32 API로 관리자 여부 확인",
        Inches(0.45), Inches(3.8), Inches(3.7), Inches(0.3),
        size=Pt(10.5), color=DARK_GRAY)
    txt(s, "• 권한 없으면 ShellExecuteW('runas') 로\n  관리자 권한으로 자신을 재실행",
        Inches(0.45), Inches(4.1), Inches(3.7), Inches(0.55),
        size=Pt(10.5), color=DARK_GRAY)
    txt(s, "• pywinauto 일부 기능은 관리자 권한 필수",
        Inches(0.45), Inches(4.7), Inches(3.7), Inches(0.3),
        size=Pt(10.5), color=DARK_GRAY)

    # 중: sMES 실행
    box(s, Inches(4.55), Inches(1.25), Inches(4.0), Inches(5.3),
        fill_color=WHITE, line_color=BLUE, line_width=Pt(1.5))
    box(s, Inches(4.55), Inches(1.25), Inches(4.0), Inches(0.42), fill_color=BLUE)
    txt(s, "🚀  sMES.exe 실행", Inches(4.7), Inches(1.27),
        Inches(3.7), Inches(0.38), size=Pt(13), bold=True, color=WHITE)
    code_box(s, Inches(4.63), Inches(1.78), Inches(3.84), Inches(1.5), [
        "SMES_EXE = Path(",
        r"  r'C:\Program Files\I2R\sMES'",
        ")",
        "",
        "subprocess.Popen([str(SMES_EXE)])",
        "time.sleep(4.0)  # 로딩 대기",
    ])
    txt(s, "• subprocess.Popen으로 비동기 실행\n  (프로세스가 뜰 때까지 4초 대기)",
        Inches(4.7), Inches(3.42), Inches(3.7), Inches(0.55),
        size=Pt(10.5), color=DARK_GRAY)
    txt(s, "• 이미 실행 중이면 psutil로 PID를 찾아 재활용",
        Inches(4.7), Inches(4.0), Inches(3.7), Inches(0.38),
        size=Pt(10.5), color=DARK_GRAY)
    txt(s, "• exe 파일 부재 시 FileNotFoundError 발생",
        Inches(4.7), Inches(4.42), Inches(3.7), Inches(0.3),
        size=Pt(10.5), color=DARK_GRAY)

    # 우: 창 연결
    box(s, Inches(8.8), Inches(1.25), Inches(4.2), Inches(5.3),
        fill_color=WHITE, line_color=TEAL, line_width=Pt(1.5))
    box(s, Inches(8.8), Inches(1.25), Inches(4.2), Inches(0.42), fill_color=TEAL)
    txt(s, "🔗  pywinauto 창 연결", Inches(8.95), Inches(1.27),
        Inches(3.9), Inches(0.38), size=Pt(13), bold=True, color=WHITE)
    code_box(s, Inches(8.88), Inches(1.78), Inches(4.04), Inches(2.0), [
        "for backend in ('win32','uia'):",
        "  app = Application(backend)",
        "      .connect(process=pid)",
        "  wins = app.windows()",
        "  # 가장 큰 창 선택",
        "  wins.sort(key=area)",
        "  return wins[0]",
    ])
    txt(s, "• win32 백엔드 우선 시도 → 실패 시 uia 재시도\n  (구형/신형 sMES 모두 호환)",
        Inches(8.95), Inches(3.92), Inches(3.9), Inches(0.55),
        size=Pt(10.5), color=DARK_GRAY)
    txt(s, "• psutil로 실행 중 PID를 먼저 탐색\n  → 프로세스 ID로 직접 연결",
        Inches(8.95), Inches(4.52), Inches(3.9), Inches(0.55),
        size=Pt(10.5), color=DARK_GRAY)
    txt(s, "• 창 면적(width×height)이 가장 큰 창을 메인 창으로 선택",
        Inches(8.95), Inches(5.12), Inches(3.9), Inches(0.38),
        size=Pt(10.5), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════
# 슬라이드 4: Step 3 — sMES 로그인
# ════════════════════════════════════════════════════════════
def slide_login(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "Step 3  sMES 로그인 자동화",
               "타이틀바 없는 borderless 로그인 폼 탐지 → Edit 컨트롤 Y좌표 정렬 → 클립보드 안전 입력")

    # 로그인 폼 탐지 (좌)
    box(s, Inches(0.3), Inches(1.25), Inches(6.1), Inches(5.3),
        fill_color=WHITE, line_color=PURPLE, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(6.1), Inches(0.42), fill_color=PURPLE)
    txt(s, "🔍  로그인 폼 탐지 로직", Inches(0.45), Inches(1.27),
        Inches(5.8), Inches(0.38), size=Pt(13), bold=True, color=WHITE)

    code_box(s, Inches(0.38), Inches(1.78), Inches(5.94), Inches(2.5), [
        "# Edit 컨트롤 2개 이상 있는 창 = 로그인 폼",
        "for w in app.windows():",
        "    edits = w.children(",
        "        class_name_re='WindowsForms10.EDIT.*'",
        "    )",
        "    if len(edits) >= 2:",
        "        return app, w  # 로그인 폼 발견!",
        "",
        "# 폴백: 가장 작은 창 (로그인 다이얼로그)",
        "smallest = min(wins, key=area)",
    ])
    steps = [
        "① sMES 프로세스 PID로 연결",
        "② 모든 하위 창 순회",
        "③ WindowsForms10.EDIT 컨트롤 2개+ 창 탐지",
        "④ 타임아웃(20초) 내 반복 탐색",
        "⑤ 실패 시 → 가장 작은 창을 로그인 폼으로 가정",
    ]
    ty = Inches(4.42)
    for st in steps:
        txt(s, st, Inches(0.45), ty, Inches(5.8), Inches(0.3),
            size=Pt(10.5), color=DARK_GRAY)
        ty += Inches(0.3)

    # 입력 로직 (우)
    box(s, Inches(6.7), Inches(1.25), Inches(6.3), Inches(5.3),
        fill_color=WHITE, line_color=ORANGE, line_width=Pt(1.5))
    box(s, Inches(6.7), Inches(1.25), Inches(6.3), Inches(0.42), fill_color=ORANGE)
    txt(s, "⌨  ID / PW 입력 전략", Inches(6.85), Inches(1.27),
        Inches(6.0), Inches(0.38), size=Pt(13), bold=True, color=WHITE)

    code_box(s, Inches(6.78), Inches(1.78), Inches(6.14), Inches(2.4), [
        "# Edit 컨트롤 Y좌표 정렬 → 위=ID, 아래=PW",
        "edits.sort(key=lambda c: c.rectangle().top)",
        "id_field, pw_field = edits[0], edits[1]",
        "",
        "# ID: type_keys 직접 입력",
        "id_field.type_keys(SMES_ID)",
        "",
        "# PW: 클립보드 붙여넣기 (특수문자 안전)",
        "win32clipboard.SetClipboardText(SMES_PW)",
        "pyautogui.hotkey('ctrl', 'v')",
        "pw_field.type_keys('{ENTER}')",
    ])

    reasons = [
        ("ID",  BLUE,   "type_keys()로 직접 입력\n일반 영문+숫자 조합"),
        ("PW",  ORANGE, "클립보드 우회 입력\n!@#$ 등 특수문자 오입력 방지"),
    ]
    ry = Inches(4.35)
    for label, col, desc in reasons:
        box(s, Inches(6.85), ry, Inches(0.55), Inches(0.55), fill_color=col)
        txt(s, label, Inches(6.85), ry + Inches(0.1), Inches(0.55), Inches(0.35),
            size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(7.5), ry + Inches(0.06), Inches(5.3), Inches(0.5),
            size=Pt(10.5), color=DARK_GRAY)
        ry += Inches(0.72)

    txt(s, "* 로그인 성공 검증: 로그인 폼(Edit 2개 창)이 닫히면 성공 판정",
        Inches(6.85), Inches(5.82), Inches(6.0), Inches(0.3),
        size=Pt(10), color=MID_GRAY, italic=True)


# ════════════════════════════════════════════════════════════
# 슬라이드 5: Step 4~5 — 메뉴 이동 & Excel 다운로드
# ════════════════════════════════════════════════════════════
def slide_navigate(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "Step 4~5  메뉴 이동 & Excel 다운로드",
               "생산관리 > 조립 자재 kitting > 날짜 설정 > 조회 > 품목별 Excel 다운로드")

    # 메뉴 클릭 흐름 (상단)
    flow_items = [
        (BLUE,   "생산관리 클릭",      "_try_click(['생산관리'])"),
        (TEAL,   "kitting 메뉴 클릭",  "_try_click(['조립 자재 kitting'])"),
        (PURPLE, "생산일자 설정",       "set_date(main_win)\nToday 자동 입력"),
        (ORANGE, "조회 버튼 클릭",      "_try_click(['조회', '검색'])"),
        (GREEN,  "Excel 다운로드",      "download_all_items(main_win)"),
    ]
    fw = Inches(2.42)
    fx = Inches(0.3)
    for i, (col, title, desc) in enumerate(flow_items):
        box(s, fx, Inches(1.25), fw, Inches(1.6),
            fill_color=WHITE, line_color=col, line_width=Pt(1.5))
        box(s, fx, Inches(1.25), fw, Inches(0.4), fill_color=col)
        txt(s, title, fx+Inches(0.12), Inches(1.28), fw-Inches(0.24), Inches(0.35),
            size=Pt(11.5), bold=True, color=WHITE)
        for j, line in enumerate(desc.split('\n')):
            txt(s, line, fx+Inches(0.15), Inches(1.75)+j*Inches(0.3),
                fw-Inches(0.3), Inches(0.3), size=Pt(10), color=DARK_GRAY)
        if i < len(flow_items)-1:
            txt(s, "→", fx+fw, Inches(1.7), Inches(0.12), Inches(0.4),
                size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        fx += fw + Inches(0.12)

    # _try_click 설명 (좌하)
    box(s, Inches(0.3), Inches(3.1), Inches(5.8), Inches(3.45),
        fill_color=WHITE, line_color=TEAL, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(3.1), Inches(5.8), Inches(0.42), fill_color=TEAL)
    txt(s, "🔍  _try_click() 다중 탐색 전략",
        Inches(0.45), Inches(3.12), Inches(5.5), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)
    code_box(s, Inches(0.38), Inches(3.63), Inches(5.64), Inches(1.65), [
        "for name in candidates:",
        "  for ct in [Button, MenuItem,",
        "    TreeItem, ListItem, Text]:",
        "    ctrl = win.child_window(",
        "        title_re=f'.*{name}.*',",
        "        control_type=ct)",
        "    if ctrl.exists(): ctrl.click_input()",
    ])
    txt(s, "• 후보 텍스트 여러 개 & 컨트롤 타입 전수 시도",
        Inches(0.45), Inches(5.38), Inches(5.5), Inches(0.28),
        size=Pt(10.5), color=DARK_GRAY)
    txt(s, "• 클릭 성공 시 True, 전부 실패 시 수동 입력 요청",
        Inches(0.45), Inches(5.66), Inches(5.5), Inches(0.28),
        size=Pt(10.5), color=DARK_GRAY)

    # 다운로드 방식 (우하)
    box(s, Inches(6.4), Inches(3.1), Inches(6.6), Inches(3.45),
        fill_color=WHITE, line_color=ORANGE, line_width=Pt(1.5))
    box(s, Inches(6.4), Inches(3.1), Inches(6.6), Inches(0.42), fill_color=ORANGE)
    txt(s, "📥  품목별 Excel 다운로드 2가지 방식",
        Inches(6.55), Inches(3.12), Inches(6.3), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)

    for i, (col, method, desc1, desc2) in enumerate([
        (BLUE, "방식 A: 그리드 행 탐지",
         "DataItem/ListItem 컨트롤 목록 감지",
         "행 클릭 → Excel 버튼 → 저장 다이얼로그 처리"),
        (PURPLE, "방식 B: 키보드 Down 방식 (폴백)",
         "그리드 행 감지 실패 시 자동 전환",
         "Ctrl+Home → Down키 반복 → 최대 200행 처리"),
    ]):
        ry = Inches(3.65) + i * Inches(1.45)
        box(s, Inches(6.55), ry, Inches(6.25), Inches(1.28),
            fill_color=GRAY_BG, line_color=col, line_width=Pt(1))
        txt(s, method, Inches(6.7), ry+Inches(0.06), Inches(5.9), Inches(0.32),
            size=Pt(11.5), bold=True, color=col)
        txt(s, desc1, Inches(6.7), ry+Inches(0.38), Inches(5.9), Inches(0.28),
            size=Pt(10), color=DARK_GRAY)
        txt(s, desc2, Inches(6.7), ry+Inches(0.65), Inches(5.9), Inches(0.28),
            size=Pt(10), color=MID_GRAY)
        txt(s, "→ 파일명: {품목명}_{YYYYMMDD}.xlsx  /  Downloads 폴더에서 kitting 폴더로 이동",
            Inches(6.7), ry+Inches(0.95), Inches(5.9), Inches(0.26),
            size=Pt(9.5), color=TEAL)


# ════════════════════════════════════════════════════════════
# 슬라이드 6: Step 6 — Playwright 웹 자동 업로드
# ════════════════════════════════════════════════════════════
def slide_playwright(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "Step 6  Playwright — 자재부족현황 웹 자동 업로드",
               "Chromium 브라우저 자동 실행 → 로그인 → 키팅 초기화 → 파일 일괄 업로드 → 완료 팝업")

    # 업로드 흐름 5단계
    steps = [
        (BLUE,   "① 브라우저 실행",
                 "sync_playwright()\np.chromium.launch()\nheadless=False"),
        (TEAL,   "② HTML 열기",
                 "page.goto(HTML_FILE\n.as_uri())\n로컬 파일 직접 열기"),
        (PURPLE, "③ 로그인",
                 "#login-email 입력\n#login-pw 입력\n#btn-login 클릭"),
        (ORANGE, "④ 키팅 초기화\n+ 파일 업로드",
                 ".btn-kit-clear 클릭\n#file-kit.set_input_files\n(xlsx 파일 목록)"),
        (GREEN,  "⑤ 완료 팝업",
                 "page.evaluate()\nJS로 팝업 DOM 삽입\n60초 대기 후 종료"),
    ]
    bw = Inches(2.42)
    bx = Inches(0.3)
    for i, (col, title, desc) in enumerate(steps):
        box(s, bx, Inches(1.25), bw, Inches(2.6),
            fill_color=WHITE, line_color=col, line_width=Pt(1.5))
        box(s, bx, Inches(1.25), bw, Inches(0.42), fill_color=col)
        for j, line in enumerate(title.split('\n')):
            txt(s, line, bx+Inches(0.12), Inches(1.28)+j*Inches(0.28),
                bw-Inches(0.24), Inches(0.28), size=Pt(12), bold=True, color=WHITE)
        box(s, bx+Inches(0.12), Inches(1.75), bw-Inches(0.24), Pt(0.8),
            fill_color=BORDER)
        for j, line in enumerate(desc.split('\n')):
            txt(s, line, bx+Inches(0.15), Inches(1.88)+j*Inches(0.36),
                bw-Inches(0.3), Inches(0.35), size=Pt(10.5), color=DARK_GRAY)
        if i < len(steps)-1:
            txt(s, "→", bx+bw, Inches(2.1), Inches(0.12), Inches(0.4),
                size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        bx += bw + Inches(0.12)

    # 하단: 코드 + 특이사항
    box(s, Inches(0.3), Inches(4.05), Inches(6.2), Inches(3.1),
        fill_color=WHITE, line_color=BLUE, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(4.05), Inches(6.2), Inches(0.42), fill_color=BLUE)
    txt(s, "핵심 코드", Inches(0.45), Inches(4.07),
        Inches(5.9), Inches(0.38), size=Pt(13), bold=True, color=WHITE)
    code_box(s, Inches(0.38), Inches(4.58), Inches(6.04), Inches(2.42), [
        "with sync_playwright() as p:",
        "  browser = p.chromium.launch(",
        "      headless=False)",
        "  page = browser.new_context()",
        "          .new_page()",
        "  page.goto(HTML_FILE.as_uri())",
        "  page.locator('#file-kit')",
        "      .set_input_files(valid_files)",
    ])

    box(s, Inches(6.8), Inches(4.05), Inches(6.2), Inches(3.1),
        fill_color=WHITE, line_color=GREEN, line_width=Pt(1.5))
    box(s, Inches(6.8), Inches(4.05), Inches(6.2), Inches(0.42), fill_color=GREEN)
    txt(s, "주요 특이사항", Inches(6.95), Inches(4.07),
        Inches(5.9), Inches(0.38), size=Pt(13), bold=True, color=WHITE)
    notes = [
        ("headless=False",    "브라우저 창이 화면에 보이도록 실행\n  (숨김 모드 비활성)"),
        ("set_input_files()",  "복수 파일 동시 업로드 지원\n  (다운로드된 모든 xlsx 한 번에)"),
        ("page.evaluate()",    "JS 직접 실행으로 완료 팝업 DOM 삽입\n  (alert 대신 커스텀 UI)"),
        ("60초 대기",           "page.wait_for_timeout(60000)\n  사용자 확인 후 브라우저 종료"),
    ]
    ny = Inches(4.58)
    for label, desc in notes:
        txt(s, "• " + label + "  →  " + desc.replace('\n  ', ' '),
            Inches(6.95), ny, Inches(5.9), Inches(0.45),
            size=Pt(10.5), color=DARK_GRAY)
        ny += Inches(0.52)


# ════════════════════════════════════════════════════════════
# 슬라이드 7: 오류 처리 & 폴백 전략
# ════════════════════════════════════════════════════════════
def slide_error(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "오류 처리 & 폴백 전략",
               "자동화 실패 시 사용자 개입(input) 요청 또는 대체 방식으로 계속 진행")

    errors = [
        (RED,    "sMES.exe 없음",
                 "FileNotFoundError 발생",
                 "실행 경로 SMES_EXE 확인 후 수정"),
        (ORANGE, "창 연결 실패",
                 "win32 실패 → uia 재시도\n두 백엔드 모두 실패 시",
                 "\"수동으로 sMES를 열어주세요\"\ninput() 대기 후 재탐색"),
        (PURPLE, "로그인 창 미탐지",
                 "20초 타임아웃 내 Edit 2개 창\n발견 못할 경우",
                 "가장 작은 창을 로그인 폼으로\n가정 후 진행"),
        (TEAL,   "로그인 실패",
                 "10초 후에도 로그인 창이\n닫히지 않은 경우",
                 "RuntimeError → 사용자에게\n수동 로그인 안내"),
        (BLUE,   "그리드 행 미탐지",
                 "DataItem/ListItem 등\n모든 컨트롤 타입 탐색 실패",
                 "키보드 Down 방식으로 자동 전환\n(최대 200회 반복)"),
        (ORANGE, "Excel 버튼 미탐지",
                 "Excel 다운로드 버튼을\n찾지 못한 경우",
                 "\"수동 저장 후 Enter\" 안내\nDownloads 폴더에서 자동 이동"),
        (GREEN,  "저장 다이얼로그 없음",
                 "자동 저장된 경우\n(다이얼로그 미표시)",
                 "Downloads 폴더 최근 xlsx 파일\n(30초 이내) 자동 이동"),
        (RED,    "Playwright 로그인 실패",
                 "locator 탐색 실패\n(HTML 구조 변경 등)",
                 "\"수동 로그인 후 Enter\" 대기\n이후 파일 업로드 진행"),
    ]

    cols = 4
    ew = Inches(3.1)
    eh = Inches(2.45)
    sx = Inches(0.25)
    sy = Inches(1.25)

    for i, (col, title, cause, solution) in enumerate(errors):
        row = i // cols
        ci  = i % cols
        lx = sx + ci * (ew + Inches(0.1))
        ty = sy + row * (eh + Inches(0.12))
        box(s, lx, ty, ew, eh, fill_color=WHITE,
            line_color=col, line_width=Pt(1.2))
        box(s, lx, ty, ew, Inches(0.36), fill_color=col)
        txt(s, title, lx+Inches(0.12), ty+Inches(0.04),
            ew-Inches(0.24), Inches(0.3), size=Pt(11), bold=True, color=WHITE)
        # 원인
        box(s, lx+Inches(0.12), ty+Inches(0.42), ew-Inches(0.24), Inches(0.22),
            fill_color=RGBColor(0xf0, 0xf0, 0xf0))
        txt(s, "원인", lx+Inches(0.15), ty+Inches(0.43),
            Inches(0.5), Inches(0.2), size=Pt(8.5), bold=True, color=MID_GRAY)
        for j, line in enumerate(cause.split('\n')):
            txt(s, line, lx+Inches(0.15), ty+Inches(0.67+j*0.27),
                ew-Inches(0.3), Inches(0.26), size=Pt(9.5), color=DARK_GRAY)
        # 처리
        box(s, lx+Inches(0.12), ty+Inches(1.3), ew-Inches(0.24), Inches(0.22),
            fill_color=col)
        txt(s, "처리", lx+Inches(0.15), ty+Inches(1.31),
            Inches(0.5), Inches(0.2), size=Pt(8.5), bold=True, color=WHITE)
        for j, line in enumerate(solution.split('\n')):
            txt(s, line, lx+Inches(0.15), ty+Inches(1.58+j*0.27),
                ew-Inches(0.3), Inches(0.26), size=Pt(9.5), color=col)


# ════════════════════════════════════════════════════════════
# 슬라이드 8: 설치 & 실행 방법
# ════════════════════════════════════════════════════════════
def slide_setup(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    for ox, oy, ow, col in [
        (Inches(10.5), Inches(-0.5), Inches(5), RGBColor(0x25, 0x4e, 0x7e)),
        (Inches(-1),   Inches(5),    Inches(4), RGBColor(0x18, 0x30, 0x50)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    header_bar(s, "사전 설치 & 실행 방법", "")

    txt(s, "1  라이브러리 설치 (pip)",
        Inches(0.4), Inches(1.25), Inches(12), Inches(0.42),
        size=Pt(16), bold=True, color=WHITE)
    code_box(s, Inches(0.4), Inches(1.75), Inches(12.5), Inches(0.82), [
        "pip install pywinauto pyautogui pillow playwright psutil keyboard pywin32",
    ])

    txt(s, "2  Playwright 브라우저 설치",
        Inches(0.4), Inches(2.72), Inches(12), Inches(0.42),
        size=Pt(16), bold=True, color=WHITE)
    code_box(s, Inches(0.4), Inches(3.22), Inches(12.5), Inches(0.62), [
        "playwright install chromium",
    ])

    txt(s, "3  실행 (관리자 권한 자동 획득)",
        Inches(0.4), Inches(4.02), Inches(12), Inches(0.42),
        size=Pt(16), bold=True, color=WHITE)
    code_box(s, Inches(0.4), Inches(4.52), Inches(12.5), Inches(0.62), [
        "python kitting_automation.py     # 또는  run.bat 더블클릭",
    ])

    points = [
        ("SMES_EXE",    "sMES 설치 경로 확인 필수  (C:\\Program Files (x86)\\I2R\\sMES\\sMES.exe)"),
        ("SMES_ID/PW",  "sMES 로그인 계정 설정  (SSAT045 / rlatndus1!)"),
        ("WEB_EMAIL/PW","자재부족현황 웹앱 계정 설정"),
        ("DOWNLOAD_DIR","Excel 저장 경로  (kitting 자재 폴더)"),
    ]
    ty = Inches(5.3)
    for key, desc in points:
        box(s, Inches(0.4), ty, Inches(2.0), Inches(0.35), fill_color=BLUE)
        txt(s, key, Inches(0.45), ty + Inches(0.04),
            Inches(1.9), Inches(0.28), size=Pt(10.5), bold=True, color=WHITE)
        txt(s, desc, Inches(2.55), ty + Inches(0.05),
            Inches(10.0), Inches(0.28), size=Pt(11), color=RGBColor(0xc0, 0xd8, 0xf0))
        ty += Inches(0.45)


# ════════════════════════════════════════════════════════════
# 메인
# ════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_launch(prs)
    slide_login(prs)
    slide_navigate(prs)
    slide_playwright(prs)
    slide_error(prs)
    slide_setup(prs)

    out_path = (r"C:\Users\조립\Desktop\claude"
                r"\Material Shortage Status vs. Production Plan"
                r"\키팅_자동화_구현방법.pptx")
    prs.save(out_path)
    import sys
    sys.stdout.buffer.write(("저장 완료: " + out_path + "\n").encode("utf-8"))

if __name__ == '__main__':
    main()
