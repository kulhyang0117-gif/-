#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
생산계획 대비 자재부족현황 대시보드 구현 방법 PPT 생성기
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys

# ── 색상 팔레트 ──────────────────────────────────────────
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)   # 헤더 배경
BLUE       = RGBColor(0x29, 0x52, 0xa3)   # 서브 강조
GREEN      = RGBColor(0x27, 0xae, 0x60)   # 정상
RED        = RGBColor(0xe7, 0x4c, 0x3c)   # 위험
ORANGE     = RGBColor(0xe6, 0x7e, 0x22)   # 경고
PURPLE     = RGBColor(0x8e, 0x44, 0xad)   # 키팅
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf2, 0xf5)
DARK_GRAY  = RGBColor(0x44, 0x44, 0x44)
MID_GRAY   = RGBColor(0x77, 0x77, 0x77)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]   # 완전 빈 레이아웃
    return prs.slides.add_slide(blank)


def bg(slide, color):
    """슬라이드 배경색 설정"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """텍스트 박스 추가"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=""):
    """슬라이드 공통 상단 헤더 바"""
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), NAVY)
    txt(slide, title,
        Inches(0.4), Inches(0.12), Inches(10), Inches(0.6),
        size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.42), Inches(0.72), Inches(10), Inches(0.35),
            size=Pt(13), color=RGBColor(0xb0, 0xc4, 0xde))
    # 슬라이드 번호 표시 위치 (오른쪽)
    txt(slide, "자재부족현황 대시보드",
        Inches(9.5), Inches(0.18), Inches(3.5), Inches(0.35),
        size=Pt(11), color=RGBColor(0x80, 0xa0, 0xc8), align=PP_ALIGN.RIGHT)


def section_label(slide, text, l, t, w=Inches(3), color=NAVY):
    """섹션 레이블 (작은 박스)"""
    box(slide, l, t, w, Inches(0.32), fill_color=color)
    txt(slide, text, l + Inches(0.1), t + Inches(0.03),
        w - Inches(0.2), Inches(0.28),
        size=Pt(11), bold=True, color=WHITE)


def card(slide, l, t, w, h, title, lines, title_color=NAVY,
         bg_color=RGBColor(0xff,0xff,0xff),
         border_color=RGBColor(0xcc,0xd6,0xe5)):
    """정보 카드"""
    box(slide, l, t, w, h, fill_color=bg_color, line_color=border_color, line_width=Pt(1))
    txt(slide, title,
        l + Inches(0.15), t + Inches(0.12), w - Inches(0.3), Inches(0.35),
        size=Pt(13), bold=True, color=title_color)
    # 구분선
    box(slide, l + Inches(0.1), t + Inches(0.48), w - Inches(0.2), Pt(1),
        fill_color=border_color)
    body_t = t + Inches(0.55)
    for line in lines:
        txt(slide, line,
            l + Inches(0.18), body_t, w - Inches(0.36), Inches(0.33),
            size=Pt(11), color=DARK_GRAY)
        body_t += Inches(0.3)


def arrow(slide, x1, y1, x2, y2, color=NAVY):
    """간단한 수평/수직 화살표 (line+삼각형 대신 텍스트 화살표)"""
    pass   # python-pptx connector 대신 텍스트로 표현


# ═══════════════════════════════════════════════════════════
# 슬라이드 1 : 표지
# ═══════════════════════════════════════════════════════════
def slide_cover(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    # 배경 장식 - 큰 반투명 원
    circ = s.shapes.add_shape(9, Inches(9), Inches(-1), Inches(6), Inches(6))   # 9=타원
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x25, 0x4e, 0x7e)
    circ.line.fill.background()

    circ2 = s.shapes.add_shape(9, Inches(-1.5), Inches(4.5), Inches(5), Inches(5))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(0x18, 0x30, 0x50)
    circ2.line.fill.background()

    # 메인 타이틀
    txt(s, "📦 생산계획 대비", Inches(1), Inches(1.2), Inches(11), Inches(0.9),
        size=Pt(38), bold=True, color=RGBColor(0xb0, 0xc8, 0xff))
    txt(s, "자재부족현황 대시보드",
        Inches(1), Inches(2.05), Inches(11), Inches(1.1),
        size=Pt(46), bold=True, color=WHITE)
    txt(s, "구현 방법 및 기술 설명서",
        Inches(1), Inches(3.15), Inches(11), Inches(0.55),
        size=Pt(22), color=RGBColor(0x8a, 0xb4, 0xf8))

    # 구분선
    box(s, Inches(1), Inches(3.85), Inches(5.5), Pt(2), fill_color=BLUE)

    # 메타 정보
    metas = [
        ("기술 스택",    "HTML · JavaScript · SheetJS · Supabase · Python"),
        ("배포 환경",    "Vercel (웹앱) + 로컬 실행 (자동화)"),
        ("주요 기능",    "BOM × 생산계획 분석 · 실시간 파일 공유 · 키팅 자동화"),
    ]
    ty = Inches(4.05)
    for label, val in metas:
        txt(s, f"• {label}",
            Inches(1), ty, Inches(2.3), Inches(0.38),
            size=Pt(12), bold=True, color=RGBColor(0x8a, 0xb4, 0xf8))
        txt(s, val,
            Inches(3.2), ty, Inches(8), Inches(0.38),
            size=Pt(12), color=RGBColor(0xcc, 0xdd, 0xf5))
        ty += Inches(0.42)

    txt(s, "2026 · 조립팀 내부 도구",
        Inches(1), Inches(6.7), Inches(6), Inches(0.4),
        size=Pt(11), color=RGBColor(0x55, 0x75, 0xa0), italic=True)


# ═══════════════════════════════════════════════════════════
# 슬라이드 2 : 프로젝트 개요
# ═══════════════════════════════════════════════════════════
def slide_overview(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "01. 프로젝트 개요", "왜 만들었는가, 무엇을 해결하는가")

    # 문제 정의
    section_label(s, "■ 문제 상황", Inches(0.4), Inches(1.35), Inches(2.2), RED)
    problems = [
        "① 생산 착수 직전에야 자재 부족이 발견됨 → 라인 정지",
        "② BOM × 생산계획 수작업 대조로 수십 분 소요",
        "③ 재고·키팅 현황이 담당자마다 달라 정보 불일치",
        "④ sMES 키팅 파일 수동 다운로드 → 업로드 반복 작업",
    ]
    ty = Inches(1.8)
    for p in problems:
        txt(s, p, Inches(0.5), ty, Inches(5.8), Inches(0.38),
            size=Pt(12.5), color=DARK_GRAY)
        ty += Inches(0.4)

    # 솔루션
    section_label(s, "■ 솔루션", Inches(6.8), Inches(1.35), Inches(2.2), GREEN)
    solutions = [
        "✅ BOM + 생산계획 + 재고를 자동 대조 · 시각화",
        "✅ 부족 자재를 위험도(🚨/🔴/🟡) 순으로 즉시 표시",
        "✅ Supabase 클라우드로 팀 전체 실시간 파일 공유",
        "✅ Python 자동화로 sMES → 웹 1-클릭 업로드",
    ]
    ty = Inches(1.8)
    for sol in solutions:
        txt(s, sol, Inches(6.9), ty, Inches(5.9), Inches(0.38),
            size=Pt(12.5), color=DARK_GRAY)
        ty += Inches(0.4)

    # 가운데 화살표
    txt(s, "→", Inches(6.2), Inches(3.2), Inches(0.5), Inches(0.5),
        size=Pt(32), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 효과 요약 카드
    box(s, Inches(0.4), Inches(5.2), Inches(12.4), Inches(1.8),
        fill_color=NAVY, line_color=BLUE, line_width=Pt(1))
    txt(s, "📊 기대 효과",
        Inches(0.6), Inches(5.3), Inches(3), Inches(0.4),
        size=Pt(14), bold=True, color=WHITE)
    effects = [
        ("⏱ 분석 시간", "수십 분 → 수 초"),
        ("🎯 정확도",    "수작업 오류 제거"),
        ("👥 공유",      "팀 전원 동일 정보"),
        ("🤖 자동화",    "sMES 연동 1-클릭"),
    ]
    ex = Inches(0.6)
    for label, val in effects:
        txt(s, label, ex, Inches(5.78), Inches(1.5), Inches(0.32),
            size=Pt(11), bold=True, color=RGBColor(0xb0, 0xc8, 0xff))
        txt(s, val, ex, Inches(6.12), Inches(1.5), Inches(0.32),
            size=Pt(13), bold=True, color=WHITE)
        ex += Inches(3.1)


# ═══════════════════════════════════════════════════════════
# 슬라이드 3 : 시스템 아키텍처
# ═══════════════════════════════════════════════════════════
def slide_architecture(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "02. 시스템 아키텍처", "전체 구성 요소와 데이터 흐름")

    # 레이어 박스들
    layers = [
        # (l, t, w, h, bg, label, desc)
        (Inches(0.3),  Inches(1.3),  Inches(2.7), Inches(4.8),
         RGBColor(0xe8,0xf0,0xfa), "📁 데이터 소스 (Excel)",
         ["BOM 정보\n(제품번호·자품번·소요량)",
          "생산계획\n(일자·P/N·수량)",
          "재고현황\n(품목코드·재고수량)",
          "포장사양\n(포장 기준)",
          "키팅된 자재\n(복수 파일 가능)"]),

        (Inches(3.4),  Inches(1.3),  Inches(3.2), Inches(4.8),
         RGBColor(0xf0,0xf5,0xff), "🌐 웹 대시보드 (HTML+JS)",
         ["SheetJS: Excel 파싱",
          "BOM × 계획 → 소요량 계산",
          "재고 + 키팅 대조 분석",
          "위험도 분류 & 시각화",
          "탭별 결과 테이블 렌더링"]),

        (Inches(6.95), Inches(1.3),  Inches(3.0), Inches(4.8),
         RGBColor(0xf5,0xf0,0xff), "☁️ Supabase (BaaS)",
         ["Auth: 회원가입·로그인",
          "관리자 승인 / 권한 관리",
          "profiles 테이블 (RLS)",
          "upload_logs 공유 동기화",
          "Storage: ms-files 버킷"]),

        (Inches(10.25), Inches(1.3),  Inches(2.75), Inches(4.8),
         RGBColor(0xf5,0xff,0xf0), "🤖 Python 자동화",
         ["pywinauto: sMES 제어",
          "키팅 Excel 자동 다운로드",
          "Playwright: 웹 자동화",
          "파일 자동 업로드",
          "1-클릭 실행 (run.bat)"]),
    ]

    for (l, t, w, h, bg_c, label, items) in layers:
        box(s, l, t, w, h, fill_color=bg_c,
            line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1.5))
        txt(s, label, l+Inches(0.1), t+Inches(0.1),
            w-Inches(0.2), Inches(0.45),
            size=Pt(12), bold=True, color=NAVY)
        box(s, l+Inches(0.05), t+Inches(0.53), w-Inches(0.1), Pt(1),
            fill_color=RGBColor(0xbb,0xcc,0xdd))
        iy = t + Inches(0.65)
        for item in items:
            txt(s, "• "+item, l+Inches(0.15), iy,
                w-Inches(0.3), Inches(0.6),
                size=Pt(10.5), color=DARK_GRAY)
            iy += Inches(0.72)

    # 화살표들
    arrows = [
        (Inches(3.08), Inches(3.7), "→"),
        (Inches(6.63), Inches(3.7), "→"),
        (Inches(9.95), Inches(3.7), "→"),
    ]
    for (al, at, ar) in arrows:
        txt(s, ar, al, at, Inches(0.38), Inches(0.45),
            size=Pt(22), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 하단 설명
    box(s, Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.75),
        fill_color=NAVY)
    txt(s, "💡 서버 없는 아키텍처: 모든 연산은 브라우저에서 처리 (SheetJS) · 클라우드는 파일 공유·인증 전용 (Supabase) · Vercel 정적 배포",
        Inches(0.5), Inches(6.5), Inches(12.2), Inches(0.55),
        size=Pt(11.5), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 4 : 기술 스택
# ═══════════════════════════════════════════════════════════
def slide_techstack(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "03. 기술 스택", "선택 이유와 역할")

    techs = [
        ("🌐 HTML + CSS + JS",
         NAVY,
         ["• 별도 프레임워크 없음 — 순수 Vanilla JS",
          "• 단일 파일(자재부족현황.html)로 배포",
          "• CDN으로 외부 라이브러리 로드",
          "• 설치 없이 브라우저만으로 실행 가능"]),
        ("📊 SheetJS (xlsx@0.18.5)",
         BLUE,
         ["• 브라우저에서 Excel 직접 파싱",
          "• .xlsx / .xls 모두 지원",
          "• ArrayBuffer → JSON 변환",
          "• 결과물 Excel 내보내기 (Export)"]),
        ("☁️ Supabase",
         RGBColor(0x1a, 0x8c, 0x5b),
         ["• PostgreSQL 기반 BaaS",
          "• 이메일/비밀번호 인증 (Auth)",
          "• Row Level Security (RLS) 정책",
          "• Storage 버킷으로 파일 공유"]),
        ("🤖 Python (자동화)",
         PURPLE,
         ["• pywinauto: Win32 GUI 자동화",
          "• pyautogui: 마우스/키보드 제어",
          "• playwright: 브라우저 자동화",
          "• psutil: 프로세스 감지"]),
        ("🚀 Vercel",
         RGBColor(0x00, 0x00, 0x00),
         ["• 정적 파일 무료 배포",
          "• vercel.json 라우팅 설정",
          "• HTTPS 자동 제공",
          "• 팀 URL 공유로 접속"]),
        ("💾 IndexedDB + LocalStorage",
         DARK_GRAY,
         ["• 업로드 파일 브라우저 캐시",
          "• 세션 재시작 후 복원",
          "• 업로드 시각·담당자 저장",
          "• Supabase 미연결 시 오프라인 지원"]),
    ]

    positions = [
        (Inches(0.3), Inches(1.35)),
        (Inches(4.55), Inches(1.35)),
        (Inches(8.8),  Inches(1.35)),
        (Inches(0.3),  Inches(4.0)),
        (Inches(4.55), Inches(4.0)),
        (Inches(8.8),  Inches(4.0)),
    ]

    W = Inches(4.0)
    H = Inches(2.4)

    for i, ((l, t), (title, tc, lines)) in enumerate(zip(positions, techs)):
        box(s, l, t, W, H,
            fill_color=WHITE,
            line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1.2))
        # 컬러 탑바
        box(s, l, t, W, Inches(0.38), fill_color=tc)
        txt(s, title, l+Inches(0.12), t+Inches(0.05),
            W-Inches(0.2), Inches(0.3),
            size=Pt(12.5), bold=True, color=WHITE)
        iy = t + Inches(0.5)
        for line in lines:
            txt(s, line, l+Inches(0.15), iy,
                W-Inches(0.3), Inches(0.38),
                size=Pt(10.5), color=DARK_GRAY)
            iy += Inches(0.42)


# ═══════════════════════════════════════════════════════════
# 슬라이드 5 : 핵심 알고리즘 — 분석 로직
# ═══════════════════════════════════════════════════════════
def slide_algorithm(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "04. 핵심 분석 알고리즘", "BOM × 생산계획 → 소요량 계산 → 재고 대조")

    # 플로우 스텝
    steps = [
        ("①", "BOM 파싱",
         "wb.Sheets['BOM정보']\n→ 제품번호별 자품번+소요량 Map",
         NAVY),
        ("②", "생산계획 파싱",
         "시트별 일자+P/N+수량\n→ planRows 배열",
         BLUE),
        ("③", "소요량 계산",
         "planRow.qty × comp.소요량\n→ reqMap[자품번].totalReq",
         RGBColor(0x16, 0x7a, 0x5e)),
        ("④", "재고·키팅 대조",
         "available = 재고 + 키팅수량\nshortage = available − totalReq",
         RGBColor(0xc0, 0x5a, 0x10)),
        ("⑤", "위험도 분류",
         "danger: shortage < 0\nwarning: ratio < 1.2\nok: 여유 충분",
         RED),
        ("⑥", "긴급도 판정",
         "1영업일 이내 → 🚨 urgent1\n3영업일 이내 → 🔴 urgent3\n이후 → later",
         PURPLE),
    ]

    lx = Inches(0.3)
    for i, (num, title, desc, color) in enumerate(steps):
        if i < 3:
            l = Inches(0.3) + i * Inches(4.3)
            t = Inches(1.35)
        else:
            l = Inches(0.3) + (i-3) * Inches(4.3)
            t = Inches(3.9)

        box(s, l, t, Inches(3.9), Inches(2.2),
            fill_color=WHITE,
            line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1))
        box(s, l, t, Inches(3.9), Inches(0.5), fill_color=color)
        txt(s, num + " " + title,
            l+Inches(0.12), t+Inches(0.08),
            Inches(3.7), Inches(0.36),
            size=Pt(13), bold=True, color=WHITE)
        txt(s, desc,
            l+Inches(0.15), t+Inches(0.62),
            Inches(3.62), Inches(1.48),
            size=Pt(11), color=DARK_GRAY)

    # 수식 요약
    box(s, Inches(0.3), Inches(6.38), Inches(12.7), Inches(0.78), fill_color=NAVY)
    txt(s, "핵심 수식:  [ available = (재고수량) + (키팅수량) ]   →   [ shortage = available − (BOM 소요량 × 계획수량) ]   →   [ ratio = available / totalReq ]",
        Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55),
        size=Pt(11.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# 슬라이드 6 : 데이터 입력 & 파싱
# ═══════════════════════════════════════════════════════════
def slide_parsing(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "05. 데이터 입력 & Excel 파싱", "유연한 헤더 자동 감지 + 다중 파일 지원")

    # 입력 파일 테이블
    section_label(s, "■ 입력 파일 목록", Inches(0.4), Inches(1.35), Inches(2.5))
    headers = ["구분", "파일", "필수", "주요 컬럼", "특이사항"]
    rows_data = [
        ["BOM",       "BOM.xlsx",        "필수", "제품번호 · 자품번 · 소요량 · 레벨", "BOM정보 > BOM > 첫 시트 순 탐색"],
        ["생산계획",  "조립_D.xlsx 등",  "필수", "일자 · P/N · 수량",              "다중 시트 지원, 텍스트 날짜 파싱"],
        ["재고현황",  "재고.xlsx",        "필수", "품목코드 · 재고수량",             "동일 품번 수량 합산"],
        ["포장사양",  "포장사양.xlsx",    "선택", "P/N · 포장수량",                 "포장재 소요량 계산"],
        ["키팅된 자재", "kit_*.xlsx",     "선택", "자재코드 · 수량",                "복수 파일, 모델별 합산"],
    ]
    col_ws = [Inches(1.1), Inches(1.7), Inches(0.7), Inches(3.8), Inches(4.1)]
    col_xs = [Inches(0.35)]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w)

    # 헤더 행
    ty = Inches(1.75)
    for i, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
        box(s, cx, ty, cw, Inches(0.35), fill_color=NAVY)
        txt(s, hdr, cx+Inches(0.05), ty+Inches(0.04), cw-Inches(0.1), Inches(0.28),
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ty += Inches(0.35)

    for ri, row in enumerate(rows_data):
        bg_c = WHITE if ri % 2 == 0 else RGBColor(0xf5, 0xf7, 0xfa)
        for ci, (cell, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
            box(s, cx, ty, cw, Inches(0.38), fill_color=bg_c,
                line_color=RGBColor(0xdd,0xe4,0xee), line_width=Pt(0.5))
            cell_color = RED if cell == "필수" else (MID_GRAY if cell == "선택" else DARK_GRAY)
            txt(s, cell, cx+Inches(0.06), ty+Inches(0.04), cw-Inches(0.12), Inches(0.32),
                size=Pt(10.5), color=cell_color, bold=(cell in ["필수","선택"]))
        ty += Inches(0.38)

    # 파싱 전략 설명
    section_label(s, "■ 헤더 자동 감지 전략", Inches(0.4), Inches(4.42), Inches(2.8))
    strategies = [
        "• 상위 5행을 스캔하여 키워드로 헤더 행 자동 감지 (제품번호, 일자, 수량 등)",
        "• 컬럼 인덱스를 동적으로 찾아 서식 변경에도 자동 적응",
        "• Excel 날짜 시리얼 번호 ↔ 텍스트 날짜 ('3월 14일 (토)') 모두 처리",
        "• 키팅 파일: 코드KW / 수량KW 키워드 배열로 다양한 양식 지원",
    ]
    ty2 = Inches(4.82)
    for st in strategies:
        txt(s, st, Inches(0.5), ty2, Inches(12.2), Inches(0.36),
            size=Pt(11.5), color=DARK_GRAY)
        ty2 += Inches(0.38)

    # 파일 공유 설명
    box(s, Inches(0.35), Inches(6.4), Inches(12.6), Inches(0.76), fill_color=NAVY)
    txt(s, "📤 파일 업로드 시 자동 처리:  IndexedDB 로컬 캐시 저장 → Supabase Storage 업로드 → upload_logs 테이블 동기화 → 다른 기기에서 자동 복원",
        Inches(0.55), Inches(6.52), Inches(12.2), Inches(0.52),
        size=Pt(11), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 7 : 인증 & 권한 시스템
# ═══════════════════════════════════════════════════════════
def slide_auth(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "06. 인증 & 권한 관리 시스템", "Supabase Auth + Row Level Security + 관리자 승인")

    # 인증 플로우
    section_label(s, "■ 인증 플로우", Inches(0.4), Inches(1.35), Inches(2.0))
    flow_steps = [
        ("회원가입 신청", "이름·이메일·비밀번호\n입력 후 Supabase Auth 등록", NAVY),
        ("관리자 승인 대기", "status='pending'\n⏳ 승인 전 접근 불가", MID_GRAY),
        ("관리자 승인", "Admin 패널에서 승인\npermissions JSONB 설정", GREEN),
        ("로그인 & 권한 적용", "read / upload / modify\n권한별 기능 활성화", BLUE),
    ]
    fx = Inches(0.4)
    for i, (title, desc, color) in enumerate(flow_steps):
        box(s, fx, Inches(1.75), Inches(2.7), Inches(1.85),
            fill_color=WHITE, line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1))
        box(s, fx, Inches(1.75), Inches(2.7), Inches(0.4), fill_color=color)
        txt(s, title, fx+Inches(0.1), Inches(1.8),
            Inches(2.5), Inches(0.32), size=Pt(11.5), bold=True, color=WHITE)
        txt(s, desc, fx+Inches(0.12), Inches(2.2),
            Inches(2.5), Inches(1.25), size=Pt(10.5), color=DARK_GRAY)
        if i < 3:
            txt(s, "→", fx+Inches(2.72), Inches(2.55),
                Inches(0.4), Inches(0.38),
                size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        fx += Inches(3.0)

    # DB 스키마
    section_label(s, "■ 데이터베이스 스키마", Inches(0.4), Inches(3.85), Inches(2.8))
    schema_items = [
        ("profiles 테이블",
         ["id (UUID, FK → auth.users)",
          "email, display_name",
          "status: pending / approved / rejected",
          "is_admin (BOOLEAN)",
          "permissions: {read, upload, modify} (JSONB)"],
         NAVY),
        ("upload_logs 테이블",
         ["file_type (PK): bom/plan/inv/pkg/kit",
          "uploaded_at (Unix ms)",
          "uploader_name, file_name",
          "file_names[] (키팅 복수 파일)"],
         BLUE),
        ("RLS 정책",
         ["profiles: 인증된 사용자 읽기/쓰기",
          "upload_logs: 인증된 사용자 ALL",
          "storage/ms-files: 인증 SELECT/INSERT/UPDATE",
          "→ 미인증 사용자 완전 차단"],
         RGBColor(0x16,0x7a,0x5e)),
    ]
    sx = Inches(0.4)
    for (title, items, color) in schema_items:
        box(s, sx, Inches(4.25), Inches(4.1), Inches(2.4),
            fill_color=WHITE, line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1))
        box(s, sx, Inches(4.25), Inches(4.1), Inches(0.38), fill_color=color)
        txt(s, title, sx+Inches(0.12), Inches(4.3),
            Inches(3.9), Inches(0.3), size=Pt(12), bold=True, color=WHITE)
        iy = Inches(4.72)
        for item in items:
            txt(s, "• " + item, sx+Inches(0.15), iy,
                Inches(3.85), Inches(0.35), size=Pt(10.5), color=DARK_GRAY)
            iy += Inches(0.36)
        sx += Inches(4.35)

    # 권한 설명
    box(s, Inches(0.4), Inches(6.4), Inches(12.6), Inches(0.76), fill_color=NAVY)
    txt(s, "🔐 보안:  이메일 확인 OFF (내부 앱) · 관리자 승인 필수 · RLS로 DB 직접 접근 차단 · Storage 버킷도 인증 필수",
        Inches(0.6), Inches(6.52), Inches(12.1), Inches(0.52),
        size=Pt(11.5), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 8 : 결과 화면 & UI 구성
# ═══════════════════════════════════════════════════════════
def slide_ui(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "07. 결과 화면 & UI 구성", "탭 구조 · 요약 카드 · 위험도 시각화")

    # 요약 카드
    section_label(s, "■ 상단 요약 카드 (7개)", Inches(0.4), Inches(1.35), Inches(2.8))
    summary_cards = [
        ("전체 자재", "분석 대상\n자재 수", DARK_GRAY),
        ("🚨 긴급", "1영업일내\n부족 자재", RED),
        ("🔴 위험", "3영업일내\n부족 자재", RGBColor(0xc0,0x39,0x2b)),
        ("🟡 경고", "재고비율\n< 120%", ORANGE),
        ("✅ 정상", "재고 충분\n자재", GREEN),
        ("🟣 키팅", "키팅 적용\n자재 수", PURPLE),
        ("⚠ 미매칭", "BOM 없는\nP/N 수", MID_GRAY),
    ]
    cx = Inches(0.35)
    for (title, desc, color) in summary_cards:
        box(s, cx, Inches(1.75), Inches(1.78), Inches(1.15),
            fill_color=WHITE, line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1))
        txt(s, title, cx+Inches(0.08), Inches(1.82),
            Inches(1.65), Inches(0.38),
            size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, desc, cx+Inches(0.08), Inches(2.2),
            Inches(1.65), Inches(0.6),
            size=Pt(10), color=MID_GRAY, align=PP_ALIGN.CENTER)
        cx += Inches(1.84)

    # 탭 구조
    section_label(s, "■ 결과 탭 구조", Inches(0.4), Inches(3.1), Inches(2.2))
    tabs = [
        ("자재별 부족현황", NAVY,
         ["자품번 · 자품명 · 규격",
          "소요량 / 재고 / 키팅 / 가용",
          "부족수량 · 재고비율 · 진행바",
          "위험도 배지 · 긴급도 분류",
          "행 클릭 → 모델별 상세 펼치기"]),
        ("모델별 분석", BLUE,
         ["모델 × 생산계획 날짜별",
          "생산 가능 수량 자동 계산",
          "부족 자재 칩(chip) 표시",
          "일자별 소요량 상세 테이블",
          "부족 자재 클릭 → 파트 상세"]),
        ("일자별 소요량", RGBColor(0x16,0x7a,0x5e),
         ["날짜 × 자재별 피벗 테이블",
          "날짜 범위 필터링",
          "각 자재별 일일 소요량",
          "누계 소요량 vs 재고 비교"]),
        ("재고현황 조회", ORANGE,
         ["업로드된 재고 원본 표시",
          "자재 검색 · 필터",
          "재고 수동 수정 기능",
          "수정 이력 관리"]),
    ]
    tx = Inches(0.35)
    for (title, color, items) in tabs:
        box(s, tx, Inches(3.5), Inches(3.15), Inches(3.1),
            fill_color=WHITE, line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1))
        box(s, tx, Inches(3.5), Inches(3.15), Inches(0.38), fill_color=color)
        txt(s, title, tx+Inches(0.1), Inches(3.55),
            Inches(2.97), Inches(0.3), size=Pt(11.5), bold=True, color=WHITE)
        iy = Inches(3.97)
        for item in items:
            txt(s, "• " + item, tx+Inches(0.13), iy,
                Inches(2.94), Inches(0.38), size=Pt(10.5), color=DARK_GRAY)
            iy += Inches(0.46)
        tx += Inches(3.25)

    box(s, Inches(0.35), Inches(6.4), Inches(12.6), Inches(0.76), fill_color=NAVY)
    txt(s, "🎨 UI 특징:  위험도별 색상 코딩 · 컬럼 리사이즈 · 드래그앤드롭 업로드 · 스티키 헤더 · 모바일 반응형 · Excel 내보내기",
        Inches(0.55), Inches(6.52), Inches(12.1), Inches(0.52),
        size=Pt(11.5), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 9 : 키팅 자동화 (Python)
# ═══════════════════════════════════════════════════════════
def slide_automation(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "08. 키팅 자동화 (Python)", "sMES 자동 로그인 → 키팅 파일 다운로드 → 웹 자동 업로드")

    # 자동화 단계
    auto_steps = [
        ("Step 1", "sMES 실행",
         "subprocess.Popen으로\nsMES.exe 실행\n4초 로딩 대기",
         NAVY),
        ("Step 2", "자동 로그인",
         "pywinauto로 창 탐색\nID/PW 입력 자동화\nEnter 키 전송",
         BLUE),
        ("Step 3", "키팅 화면 이동",
         "메뉴 자동 클릭\n오늘 날짜 자동 입력\n조회 버튼 실행",
         RGBColor(0x16,0x7a,0x5e)),
        ("Step 4", "Excel 저장",
         "출력(Print) 버튼 클릭\n저장 경로 자동 입력\n파일명 날짜 포함",
         ORANGE),
        ("Step 5", "웹 업로드",
         "Playwright로 브라우저 제어\nSupabase 로그인\n파일 업로드 자동 실행",
         PURPLE),
    ]

    ax = Inches(0.3)
    for i, (step, title, desc, color) in enumerate(auto_steps):
        box(s, ax, Inches(1.35), Inches(2.45), Inches(3.45),
            fill_color=WHITE, line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1))
        box(s, ax, Inches(1.35), Inches(2.45), Inches(0.4), fill_color=color)
        txt(s, step, ax+Inches(0.1), Inches(1.4),
            Inches(0.8), Inches(0.3), size=Pt(10), bold=True, color=WHITE)
        txt(s, title, ax+Inches(0.9), Inches(1.4),
            Inches(1.5), Inches(0.3), size=Pt(11), bold=True, color=WHITE)
        txt(s, desc, ax+Inches(0.13), Inches(1.85),
            Inches(2.25), Inches(2.85), size=Pt(10.5), color=DARK_GRAY)
        if i < 4:
            txt(s, "→", ax+Inches(2.47), Inches(2.95),
                Inches(0.38), Inches(0.38),
                size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        ax += Inches(2.6)

    # 주요 라이브러리
    section_label(s, "■ 사용 라이브러리", Inches(0.4), Inches(5.05), Inches(2.2))
    libs = [
        ("pywinauto",  "Win32 GUI 자동화 (창 탐색·컨트롤 조작)"),
        ("pyautogui",  "마우스·키보드 제어 (픽셀 기반 백업 방법)"),
        ("playwright", "Chromium 브라우저 자동화 (웹 업로드)"),
        ("psutil",     "sMES.exe 프로세스 감지 및 PID 추적"),
        ("run.bat",    "더블클릭 1회로 전체 자동화 실행"),
    ]
    lx = Inches(0.35)
    for lib, desc in libs:
        box(s, lx, Inches(5.45), Inches(2.47), Inches(0.75),
            fill_color=NAVY)
        txt(s, lib, lx+Inches(0.1), Inches(5.5),
            Inches(2.3), Inches(0.28), size=Pt(11), bold=True, color=WHITE)
        box(s, lx, Inches(6.2), Inches(2.47), Inches(0.55),
            fill_color=WHITE, line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(0.8))
        txt(s, desc, lx+Inches(0.1), Inches(6.23),
            Inches(2.35), Inches(0.48), size=Pt(10), color=DARK_GRAY)
        lx += Inches(2.6)

    box(s, Inches(0.35), Inches(6.4), Inches(12.6), Inches(0.76), fill_color=PURPLE)
    txt(s, "🤖 실행 방법:  pip install pywinauto pyautogui playwright psutil  →  playwright install chromium  →  run.bat 더블클릭",
        Inches(0.55), Inches(6.52), Inches(12.1), Inches(0.52),
        size=Pt(11.5), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 10 : 파일 공유 & 실시간 동기화
# ═══════════════════════════════════════════════════════════
def slide_sync(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "09. 파일 공유 & 실시간 동기화", "Supabase Storage + upload_logs로 팀 전체 파일 공유")

    # 동기화 흐름도
    section_label(s, "■ 파일 공유 흐름", Inches(0.4), Inches(1.35), Inches(2.2))
    flow = [
        ("기기 A\n(업로더)", NAVY),
        ("IndexedDB\n(로컬 캐시)", BLUE),
        ("Supabase\nStorage", RGBColor(0x16,0x7a,0x5e)),
        ("upload_logs\n테이블", ORANGE),
        ("기기 B\n(자동 복원)", PURPLE),
    ]
    fx = Inches(0.35)
    for i, (label, color) in enumerate(flow):
        box(s, fx, Inches(1.75), Inches(2.3), Inches(1.1),
            fill_color=color, line_color=color)
        txt(s, label, fx+Inches(0.1), Inches(1.98),
            Inches(2.12), Inches(0.72),
            size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, "→", fx+Inches(2.32), Inches(2.18),
                Inches(0.38), Inches(0.38),
                size=Pt(20), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        fx += Inches(2.6)

    # 상세 설명
    sync_details = [
        ("📤 업로드 시 처리 흐름",
         NAVY,
         ["1. FileReader로 ArrayBuffer 읽기",
          "2. IndexedDB에 raw 바이너리 저장",
          "3. Supabase Storage에 파일 업로드",
          "4. upload_logs 테이블 업데이트 (시각·담당자·파일명)",
          "5. LocalStorage에 업로드 타임스탬프 저장"]),
        ("📥 접속 시 복원 흐름",
         RGBColor(0x16,0x7a,0x5e),
         ["1. 페이지 로드 시 syncFilesFromSupabase() 호출",
          "2. upload_logs 테이블 조회 (파일 목록·시각)",
          "3. 로컬 캐시보다 최신이면 Storage에서 다운로드",
          "4. IndexedDB에 저장 후 UI 자동 업데이트",
          "5. 업로드 시각·담당자명 헤더에 표시"]),
        ("💾 로컬 vs 클라우드",
         BLUE,
         ["• IndexedDB: 브라우저 재시작 후 복원",
          "• LocalStorage: 메타데이터(시각·담당자)",
          "• Supabase Storage: 기기간 공유 원본",
          "• upload_logs: 누가 언제 올렸는지 이력",
          "• 오프라인: IndexedDB 캐시로 동작"]),
    ]
    dx = Inches(0.35)
    for (title, color, items) in sync_details:
        box(s, dx, Inches(3.15), Inches(4.15), Inches(3.05),
            fill_color=WHITE, line_color=RGBColor(0xcc,0xd6,0xe5), line_width=Pt(1))
        box(s, dx, Inches(3.15), Inches(4.15), Inches(0.4), fill_color=color)
        txt(s, title, dx+Inches(0.12), Inches(3.2),
            Inches(3.97), Inches(0.32), size=Pt(12), bold=True, color=WHITE)
        iy = Inches(3.65)
        for item in items:
            txt(s, item, dx+Inches(0.15), iy,
                Inches(3.92), Inches(0.38), size=Pt(10.5), color=DARK_GRAY)
            iy += Inches(0.44)
        dx += Inches(4.3)

    box(s, Inches(0.35), Inches(6.4), Inches(12.6), Inches(0.76), fill_color=NAVY)
    txt(s, "☁️ 결과:  A 기기에서 Excel 업로드 → B, C 기기 재접속 시 자동으로 동일 데이터 복원 → 팀 전원이 항상 최신 자재현황 공유",
        Inches(0.55), Inches(6.52), Inches(12.1), Inches(0.52),
        size=Pt(11.5), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 11 : 배포 & 운영
# ═══════════════════════════════════════════════════════════
def slide_deploy(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "10. 배포 & 운영 방법", "Vercel 웹 배포 + 로컬 자동화 실행")

    # Vercel 배포
    section_label(s, "■ 웹 배포 (Vercel)", Inches(0.4), Inches(1.35), Inches(2.4))
    vercel_steps = [
        "① vercel.json 작성: 모든 경로를 index.html로 라우팅",
        "② Vercel CLI 또는 GitHub 연동으로 자동 배포",
        "③ HTTPS URL 생성 → 팀에 링크 공유",
        "④ 파일 변경 시 push만 하면 자동 재배포",
    ]
    ty = Inches(1.78)
    for step in vercel_steps:
        txt(s, step, Inches(0.5), ty, Inches(5.8), Inches(0.38),
            size=Pt(12), color=DARK_GRAY)
        ty += Inches(0.42)

    # vercel.json 내용
    box(s, Inches(0.4), Inches(3.62), Inches(5.8), Inches(1.35),
        fill_color=RGBColor(0x1a,0x1a,0x2e), line_color=NAVY, line_width=Pt(1))
    txt(s, 'vercel.json',
        Inches(0.55), Inches(3.67), Inches(5.5), Inches(0.3),
        size=Pt(10), bold=True, color=RGBColor(0x8a,0xb4,0xf8))
    txt(s, '{\n  "rewrites": [{ "source": "/(.*)", "destination": "/index.html" }]\n}',
        Inches(0.55), Inches(3.97), Inches(5.5), Inches(0.88),
        size=Pt(11), color=RGBColor(0xa8,0xe6,0xa8))

    # Supabase 초기 설정
    section_label(s, "■ Supabase 초기 설정", Inches(6.6), Inches(1.35), Inches(2.6))
    supa_steps = [
        "① Supabase 프로젝트 생성 (무료 플랜)",
        "② supabase_setup.sql 실행 → 테이블·RLS 생성",
        "③ Storage에 'ms-files' 버킷 생성",
        "④ HTML에 SUPABASE_URL + ANON_KEY 설정",
        "⑤ 첫 계정 생성 후 SQL로 관리자 권한 부여",
    ]
    ty2 = Inches(1.78)
    for step in supa_steps:
        txt(s, step, Inches(6.7), ty2, Inches(5.9), Inches(0.38),
            size=Pt(12), color=DARK_GRAY)
        ty2 += Inches(0.42)

    # 파이썬 자동화 실행
    section_label(s, "■ Python 자동화 설치 및 실행", Inches(0.4), Inches(5.1), Inches(3.2))
    py_steps = [
        "pip install pywinauto pyautogui playwright psutil keyboard pillow",
        "playwright install chromium",
        "kitting_automation.py 상단 설정값 수정 (sMES 계정, 저장 경로 등)",
        "run.bat 더블클릭 → 전체 자동화 순차 실행",
    ]
    ty3 = Inches(5.5)
    for step in py_steps:
        box(s, Inches(0.4), ty3, Inches(12.6), Inches(0.37),
            fill_color=RGBColor(0x1a,0x1a,0x2e))
        txt(s, "  > " + step,
            Inches(0.5), ty3+Inches(0.03), Inches(12.3), Inches(0.3),
            size=Pt(10.5), color=RGBColor(0xa8,0xe6,0xa8))
        ty3 += Inches(0.4)

    box(s, Inches(0.35), Inches(6.42), Inches(12.6), Inches(0.74), fill_color=NAVY)
    txt(s, "🌐 접속 방법:  Vercel URL 브라우저 접속 → 회원가입 신청 → 관리자 승인 → 로그인 → Excel 업로드 → 분석 실행",
        Inches(0.55), Inches(6.54), Inches(12.1), Inches(0.52),
        size=Pt(11.5), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 12 : 마무리 & 요약
# ═══════════════════════════════════════════════════════════
def slide_summary(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    # 배경 장식
    circ = s.shapes.add_shape(9, Inches(10), Inches(4), Inches(5), Inches(5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x25, 0x4e, 0x7e)
    circ.line.fill.background()

    txt(s, "구현 방법 요약",
        Inches(1), Inches(0.8), Inches(11), Inches(0.65),
        size=Pt(34), bold=True, color=RGBColor(0x8a,0xb4,0xf8))
    box(s, Inches(1), Inches(1.45), Inches(4), Pt(2), fill_color=BLUE)

    summaries = [
        ("🌐 아키텍처",   "서버 없는 순수 HTML+JS · Supabase BaaS · Vercel 배포"),
        ("📊 핵심 로직",  "BOM × 계획수량 = 소요량 → 재고+키팅 대조 → 위험도 분류"),
        ("📁 데이터 처리","SheetJS 브라우저 파싱 · 헤더 자동 감지 · 복수 시트 지원"),
        ("🔐 인증/권한",  "Supabase Auth + 관리자 승인 + RLS + 3단계 권한 관리"),
        ("☁️ 파일 공유",  "업로드 → Storage 동기화 → 다른 기기 자동 복원"),
        ("🤖 자동화",     "Python(pywinauto+playwright)으로 sMES → 웹 1-클릭 연동"),
        ("📱 UX",         "탭별 결과 · 위험도 색상 · 드래그앤드롭 · Excel 내보내기"),
    ]
    ty = Inches(1.7)
    for label, desc in summaries:
        box(s, Inches(1), ty, Inches(2.1), Inches(0.44),
            fill_color=RGBColor(0x25,0x4e,0x7e))
        txt(s, label, Inches(1.1), ty+Inches(0.06),
            Inches(1.95), Inches(0.33),
            size=Pt(12), bold=True, color=RGBColor(0x8a,0xb4,0xf8))
        txt(s, desc, Inches(3.2), ty+Inches(0.06),
            Inches(8.5), Inches(0.33),
            size=Pt(12), color=WHITE)
        ty += Inches(0.52)

    # 슬라이드 목차
    box(s, Inches(1), Inches(5.55), Inches(11), Inches(0.55),
        fill_color=RGBColor(0x18,0x30,0x50))
    slides_index = [
        "01.개요", "02.아키텍처", "03.기술스택",
        "04.알고리즘", "05.파싱", "06.인증",
        "07.UI", "08.자동화", "09.동기화", "10.배포"
    ]
    txt(s, "  |  ".join(slides_index),
        Inches(1.1), Inches(5.62), Inches(10.8), Inches(0.38),
        size=Pt(10), color=RGBColor(0x80,0xa0,0xc8), align=PP_ALIGN.CENTER)

    txt(s, "📦 생산계획 대비 자재부족현황 대시보드",
        Inches(1), Inches(6.5), Inches(11), Inches(0.55),
        size=Pt(14), bold=True, color=RGBColor(0x55,0x85,0xcc),
        align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    print("슬라이드 생성 중...")
    slide_cover(prs)       ; print("  ✅ 표지")
    slide_overview(prs)    ; print("  ✅ 01. 프로젝트 개요")
    slide_architecture(prs); print("  ✅ 02. 시스템 아키텍처")
    slide_techstack(prs)   ; print("  ✅ 03. 기술 스택")
    slide_algorithm(prs)   ; print("  ✅ 04. 핵심 알고리즘")
    slide_parsing(prs)     ; print("  ✅ 05. 데이터 파싱")
    slide_auth(prs)        ; print("  ✅ 06. 인증·권한")
    slide_ui(prs)          ; print("  ✅ 07. UI 구성")
    slide_automation(prs)  ; print("  ✅ 08. 키팅 자동화")
    slide_sync(prs)        ; print("  ✅ 09. 파일 동기화")
    slide_deploy(prs)      ; print("  ✅ 10. 배포·운영")
    slide_summary(prs)     ; print("  ✅ 마무리 요약")

    out = r"C:\Users\조립\Desktop\claude\Material Shortage Status vs. Production Plan\자재부족현황_구현방법.pptx"
    prs.save(out)
    print(f"\n✅ PPT 저장 완료: {out}")
    print(f"   총 슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    main()
