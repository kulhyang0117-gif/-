#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
자재부족현황 앱 - 자동실행 구성 설명 PPT 생성기
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ───────────────────────────────────────────
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
BLUE       = RGBColor(0x29, 0x52, 0xa3)
GREEN      = RGBColor(0x1a, 0x8a, 0x4a)
GREEN_LT   = RGBColor(0xd4, 0xed, 0xda)
RED        = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE     = RGBColor(0xe6, 0x7e, 0x22)
ORANGE_LT  = RGBColor(0xfd, 0xf3, 0xe3)
PURPLE     = RGBColor(0x8e, 0x44, 0xad)
PURPLE_LT  = RGBColor(0xf0, 0xe6, 0xf7)
TEAL       = RGBColor(0x0d, 0x86, 0x8c)
TEAL_LT    = RGBColor(0xd1, 0xf0, 0xf2)
BLUE_LT    = RGBColor(0xd6, 0xe4, 0xf7)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf4, 0xf6, 0xf9)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
BORDER     = RGBColor(0xcc, 0xd6, 0xe5)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=Pt(13), bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header_bar(slide, title, subtitle=""):
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    txt(slide, title,
        Inches(0.4), Inches(0.1), Inches(10), Inches(0.6),
        size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.42), Inches(0.68), Inches(11), Inches(0.35),
            size=Pt(12), color=RGBColor(0xb0, 0xc4, 0xde))
    txt(slide, "생산계획 대비 자재부족현황 시스템",
        Inches(9.2), Inches(0.15), Inches(3.8), Inches(0.4),
        size=Pt(10), color=RGBColor(0x80, 0xa0, 0xc8), align=PP_ALIGN.RIGHT)

def step_box(slide, l, t, w, h, num, title, lines,
             num_color=NAVY, bg_color=WHITE, border_color=BORDER):
    box(slide, l, t, w, h, fill_color=bg_color, line_color=border_color, line_width=Pt(1.2))
    # 번호 원
    circ = slide.shapes.add_shape(9, l + Inches(0.12), t + Inches(0.1),
                                   Inches(0.38), Inches(0.38))
    circ.fill.solid()
    circ.fill.fore_color.rgb = num_color
    circ.line.fill.background()
    txt(slide, str(num),
        l + Inches(0.12), t + Inches(0.1), Inches(0.38), Inches(0.38),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title,
        l + Inches(0.6), t + Inches(0.11), w - Inches(0.75), Inches(0.36),
        size=Pt(13), bold=True, color=num_color)
    # 구분선
    box(slide, l + Inches(0.12), t + Inches(0.52),
        w - Inches(0.24), Pt(0.8), fill_color=border_color)
    body_t = t + Inches(0.6)
    for line in lines:
        txt(slide, line,
            l + Inches(0.18), body_t, w - Inches(0.36), Inches(0.3),
            size=Pt(10.5), color=DARK_GRAY)
        body_t += Inches(0.28)

def arrow_right(slide, x, y):
    """→ 화살표 텍스트"""
    txt(slide, "➜", x, y, Inches(0.4), Inches(0.4),
        size=Pt(20), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def arrow_down(slide, x, y):
    """↓ 화살표 텍스트"""
    txt(slide, "↓", x, y, Inches(0.4), Inches(0.35),
        size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════
def slide_cover(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    # 장식 원
    for ox, oy, ow, col in [
        (Inches(10), Inches(-1.2), Inches(6), RGBColor(0x25, 0x4e, 0x7e)),
        (Inches(-1.5), Inches(4.8), Inches(5), RGBColor(0x18, 0x30, 0x50)),
        (Inches(6),   Inches(5.5), Inches(3), RGBColor(0x1e, 0x45, 0x72)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    # 상단 액센트 바
    box(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H,
        fill_color=RGBColor(0x29, 0x52, 0xa3))

    txt(s, "⚙ 자동실행 구성 설명서",
        Inches(0.55), Inches(1.8), Inches(12), Inches(1.0),
        size=Pt(40), bold=True, color=WHITE)
    txt(s, "생산계획 대비 자재부족현황 시스템",
        Inches(0.6), Inches(2.85), Inches(10), Inches(0.55),
        size=Pt(20), color=RGBColor(0x90, 0xb8, 0xe8))
    txt(s, "자동 파일 동기화 · 자동 분석 실행 · 다중 사용자 공유 흐름",
        Inches(0.6), Inches(3.45), Inches(10), Inches(0.45),
        size=Pt(15), color=RGBColor(0xb0, 0xc8, 0xe0), italic=True)

    # 구분선
    box(s, Inches(0.6), Inches(4.0), Inches(5), Pt(1.5),
        fill_color=BLUE)

    # 목차 미리보기
    items = [
        "① 전체 자동실행 흐름 개요",
        "② 로그인 & 초기 파일 복원",
        "③ Supabase 파일 자동 동기화 (3분 주기)",
        "④ 파일 업로드 → 자동 분석 트리거",
        "⑤ 생산계획 DB 저장 방식 (신규)",
        "⑥ 요약 & 핵심 포인트",
    ]
    ty = Inches(4.2)
    for item in items:
        txt(s, item, Inches(0.7), ty, Inches(7), Inches(0.35),
            size=Pt(12), color=RGBColor(0xc0, 0xd8, 0xf0))
        ty += Inches(0.36)


# ════════════════════════════════════════════════════════════
# 슬라이드 2: 전체 자동실행 흐름 개요
# ════════════════════════════════════════════════════════════
def slide_overview(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "① 전체 자동실행 흐름 개요",
               "앱 실행부터 분석 완료까지 자동으로 처리되는 단계")

    # 메인 흐름 - 가로 5단계
    stages = [
        ("🔐", "로그인",        BLUE,   BLUE_LT,  ["Firebase/Supabase", "인증 처리", "세션 복원"]),
        ("💾", "파일 복원",     TEAL,   TEAL_LT,  ["IndexedDB에서", "로컬 캐시 로드", "(BOM·계획·재고)"]),
        ("☁", "서버 동기화",   PURPLE, PURPLE_LT, ["Supabase DB에서", "최신 파일 확인", "타임스탬프 비교"]),
        ("📊", "자동 분석",     GREEN,  GREEN_LT,  ["3개 파일 준비 시", "즉시 분석 실행", "결과 화면 표시"]),
        ("🔄", "주기 감시",     ORANGE, ORANGE_LT, ["3분마다 반복", "변경 감지 시", "재동기화·재분석"]),
    ]

    box_w = Inches(2.3)
    box_h = Inches(3.6)
    start_x = Inches(0.3)
    top_y   = Inches(1.4)
    gap     = Inches(0.22)

    for i, (icon, title, c_color, bg_color, lines) in enumerate(stages):
        lx = start_x + i * (box_w + gap)
        box(s, lx, top_y, box_w, box_h,
            fill_color=bg_color, line_color=c_color, line_width=Pt(1.5))
        # 아이콘+타이틀 영역
        box(s, lx, top_y, box_w, Inches(0.9), fill_color=c_color)
        txt(s, icon, lx, top_y + Inches(0.05), box_w, Inches(0.45),
            size=Pt(22), color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, lx, top_y + Inches(0.5), box_w, Inches(0.38),
            size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 내용
        body_y = top_y + Inches(1.0)
        for line in lines:
            txt(s, "• " + line, lx + Inches(0.15), body_y,
                box_w - Inches(0.3), Inches(0.35),
                size=Pt(11), color=DARK_GRAY)
            body_y += Inches(0.36)
        # 화살표 (마지막 제외)
        if i < len(stages) - 1:
            ax = lx + box_w + Inches(0.02)
            txt(s, "→", ax, top_y + Inches(1.5), gap, Inches(0.45),
                size=Pt(20), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 하단 설명
    box(s, Inches(0.3), Inches(5.2), Inches(12.7), Inches(1.0),
        fill_color=WHITE, line_color=BORDER, line_width=Pt(1))
    txt(s, "💡  모든 단계는 사용자 조작 없이 자동으로 실행됩니다. "
           "파일 업로드 후 분석 버튼을 누르지 않아도 자동으로 결과가 갱신되며, "
           "3분마다 서버를 확인해 다른 사용자가 올린 최신 파일을 자동 반영합니다.",
        Inches(0.5), Inches(5.28), Inches(12.3), Inches(0.8),
        size=Pt(11.5), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════
# 슬라이드 3: 로그인 & 초기 파일 복원
# ════════════════════════════════════════════════════════════
def slide_login(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "② 로그인 & 초기 파일 복원",
               "앱 최초 진입 시 자동으로 실행되는 인증 및 파일 캐시 복원 흐름")

    # 왼쪽: 로그인 흐름
    box(s, Inches(0.3), Inches(1.25), Inches(5.8), Inches(5.3),
        fill_color=WHITE, line_color=BLUE, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(5.8), Inches(0.45), fill_color=BLUE)
    txt(s, "🔐  로그인 처리 흐름", Inches(0.45), Inches(1.27),
        Inches(5.5), Inches(0.4), size=Pt(14), bold=True, color=WHITE)

    login_steps = [
        ("1", "앱 열기", "window.load 이벤트 발생 → initAuth() 자동 실행"),
        ("2", "세션 확인", "sb.auth.getSession() — 브라우저에 저장된 로그인 세션 확인"),
        ("3", "자동 로그인", "유효 세션 있으면 로그인 화면 건너뜀 (자동 복원)"),
        ("4", "권한 적용", "profiles 테이블에서 권한 조회 → upload / modify 권한 설정"),
        ("5", "파일 동기화", "syncFilesFromSupabase() 자동 호출 → 최신 파일 확인 시작"),
    ]
    sy = Inches(1.85)
    for num, title, desc in login_steps:
        circ = s.shapes.add_shape(9, Inches(0.48), sy + Inches(0.05),
                                   Inches(0.32), Inches(0.32))
        circ.fill.solid(); circ.fill.fore_color.rgb = BLUE; circ.line.fill.background()
        txt(s, num, Inches(0.48), sy + Inches(0.05), Inches(0.32), Inches(0.32),
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, Inches(0.9), sy, Inches(1.3), Inches(0.32),
            size=Pt(11), bold=True, color=BLUE)
        txt(s, desc, Inches(0.9), sy + Inches(0.28), Inches(5.0), Inches(0.32),
            size=Pt(9.5), color=DARK_GRAY)
        if int(num) < 5:
            txt(s, "│", Inches(0.59), sy + Inches(0.35), Inches(0.2), Inches(0.28),
                size=Pt(10), color=BORDER, align=PP_ALIGN.CENTER)
        sy += Inches(0.72)

    # 오른쪽: IndexedDB 복원 흐름
    box(s, Inches(6.5), Inches(1.25), Inches(6.5), Inches(5.3),
        fill_color=WHITE, line_color=TEAL, line_width=Pt(1.5))
    box(s, Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.45), fill_color=TEAL)
    txt(s, "💾  IndexedDB 로컬 캐시 복원", Inches(6.65), Inches(1.27),
        Inches(6.2), Inches(0.4), size=Pt(14), bold=True, color=WHITE)

    cache_items = [
        ("BOM",    "부품 소요량 정보 (.xlsx)",         "bom"),
        ("생산계획", "조립 A/B/C/D 시트 포함 (.xlsx)", "plan"),
        ("재고현황", "현재 재고 수량 (.xlsx)",           "inv"),
        ("포장사양", "제품별 포장재 사양 (.xlsx)",       "pkg"),
        ("키팅자재", "복수 파일 목록 (배열 저장)",       "kit_list"),
    ]
    ry = Inches(1.85)
    for label, desc, key in cache_items:
        box(s, Inches(6.65), ry, Inches(6.15), Inches(0.58),
            fill_color=TEAL_LT, line_color=RGBColor(0xa0, 0xd8, 0xdc), line_width=Pt(0.8))
        txt(s, f"🔑 키: '{key}'", Inches(6.78), ry + Inches(0.03),
            Inches(2.0), Inches(0.28), size=Pt(9), color=TEAL, bold=True)
        txt(s, label, Inches(8.85), ry + Inches(0.03),
            Inches(1.2), Inches(0.28), size=Pt(10), bold=True, color=DARK_GRAY)
        txt(s, desc, Inches(6.78), ry + Inches(0.3),
            Inches(6.0), Inches(0.25), size=Pt(9.5), color=MID_GRAY)
        ry += Inches(0.68)

    txt(s, "※ 로컬 캐시가 있으면 즉시 화면에 표시 (서버 응답 대기 없이)",
        Inches(6.65), Inches(5.25), Inches(6.1), Inches(0.35),
        size=Pt(10), color=MID_GRAY, italic=True)


# ════════════════════════════════════════════════════════════
# 슬라이드 4: Supabase 파일 자동 동기화
# ════════════════════════════════════════════════════════════
def slide_sync(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "③ Supabase 파일 자동 동기화 (3분 주기)",
               "서버의 최신 파일을 자동으로 감지하고 로컬에 반영하는 메커니즘")

    # 타임스탬프 비교 로직 설명
    box(s, Inches(0.3), Inches(1.25), Inches(8.0), Inches(2.5),
        fill_color=WHITE, line_color=PURPLE, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(8.0), Inches(0.42), fill_color=PURPLE)
    txt(s, "⏱  타임스탬프 비교로 불필요한 다운로드 방지",
        Inches(0.45), Inches(1.27), Inches(7.7), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)

    txt(s, "serverTs  =  upload_logs.uploaded_at  (Supabase DB)",
        Inches(0.5), Inches(1.8), Inches(7.6), Inches(0.32),
        size=Pt(11), color=DARK_GRAY)
    txt(s, "localTs   =  localStorage.ms_uptime_{type}  (브라우저)",
        Inches(0.5), Inches(2.1), Inches(7.6), Inches(0.32),
        size=Pt(11), color=DARK_GRAY)

    # 조건 박스
    conditions = [
        (ORANGE_LT, ORANGE, "localTs  <  serverTs",  "→  서버가 더 최신  →  DB에서 파일 다운로드"),
        (GREEN_LT,  GREEN,  "localTs  ≥  serverTs",  "→  로컬이 최신  →  다운로드 생략 (캐시 사용)"),
    ]
    cx = Inches(0.5)
    for bg_c, bd_c, cond, result in conditions:
        box(s, cx, Inches(2.55), Inches(3.5), Inches(0.68),
            fill_color=bg_c, line_color=bd_c, line_width=Pt(1))
        txt(s, cond, cx + Inches(0.12), Inches(2.6),
            Inches(3.26), Inches(0.32), size=Pt(12), bold=True, color=bd_c)
        txt(s, result, cx + Inches(0.12), Inches(2.92),
            Inches(3.26), Inches(0.28), size=Pt(10), color=DARK_GRAY)
        cx += Inches(3.7)

    # 동기화 실행 시점
    box(s, Inches(0.3), Inches(3.9), Inches(8.0), Inches(2.7),
        fill_color=WHITE, line_color=BLUE, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(3.9), Inches(8.0), Inches(0.42), fill_color=BLUE)
    txt(s, "🕐  동기화 실행 시점",
        Inches(0.45), Inches(3.92), Inches(7.7), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)

    triggers = [
        ("로그인 직후",  "initAuth() 완료 → syncFilesFromSupabase() 즉시 실행"),
        ("3분마다",      "setInterval(syncFilesFromSupabase, 3분) — 탭을 열어두면 자동 반복"),
        ("수동 새로고침","페이지 reload 시 load 이벤트 → 자동 실행"),
    ]
    ty = Inches(4.42)
    for trigger, desc in triggers:
        box(s, Inches(0.5), ty, Inches(1.5), Inches(0.38),
            fill_color=BLUE, line_color=None)
        txt(s, trigger, Inches(0.52), ty + Inches(0.03),
            Inches(1.46), Inches(0.32), size=Pt(10), bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(2.2), ty + Inches(0.04),
            Inches(5.9), Inches(0.32), size=Pt(11), color=DARK_GRAY)
        ty += Inches(0.56)

    # 오른쪽: 파일별 처리 흐름
    box(s, Inches(8.6), Inches(1.25), Inches(4.4), Inches(5.35),
        fill_color=WHITE, line_color=TEAL, line_width=Pt(1.5))
    box(s, Inches(8.6), Inches(1.25), Inches(4.4), Inches(0.42), fill_color=TEAL)
    txt(s, "📁  파일별 처리 순서",
        Inches(8.75), Inches(1.27), Inches(4.1), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)

    files = [
        ("BOM",    "bom",  "부품 소요량"),
        ("생산계획", "plan", "조립 시트 포함"),
        ("재고현황", "inv",  "창고 재고"),
        ("포장사양", "pkg",  "포장재 마스터"),
    ]
    fy = Inches(1.8)
    for i, (name, key, desc) in enumerate(files):
        col = [BLUE, PURPLE, TEAL, ORANGE][i]
        box(s, Inches(8.75), fy, Inches(4.1), Inches(0.8),
            fill_color=LIGHT_GRAY, line_color=col, line_width=Pt(1))
        txt(s, f"{name}  ({key})", Inches(8.9), fy + Inches(0.04),
            Inches(3.8), Inches(0.3), size=Pt(11), bold=True, color=col)
        txt(s, f"① 서버 타임스탬프 확인  ② DB file_data 조회",
            Inches(8.9), fy + Inches(0.36), Inches(3.8), Inches(0.25),
            size=Pt(9.5), color=DARK_GRAY)
        txt(s, f"③ XLSX 파싱  ④ IndexedDB 갱신  ⑤ 화면 업데이트",
            Inches(8.9), fy + Inches(0.57), Inches(3.8), Inches(0.25),
            size=Pt(9.5), color=MID_GRAY)
        fy += Inches(0.95)

    txt(s, "↓ 4개 파일 모두 준비되면 자동 분석 실행",
        Inches(8.75), Inches(5.65), Inches(4.1), Inches(0.35),
        size=Pt(11), bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 5: 파일 업로드 → 자동 분석 트리거
# ════════════════════════════════════════════════════════════
def slide_upload_trigger(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "④ 파일 업로드 → 자동 분석 트리거",
               "사용자가 파일을 올리면 즉시 분석이 자동 실행되는 흐름")

    # 업로드 방법 2가지
    for i, (method, icon, desc) in enumerate([
        ("드래그 & 드롭", "🖱", "파일을 업로드 박스에 끌어다 놓기"),
        ("클릭 파일 선택", "📂", "박스 클릭 → 파일 탐색기에서 선택"),
    ]):
        bx = Inches(0.3) + i * Inches(3.8)
        box(s, bx, Inches(1.3), Inches(3.5), Inches(0.9),
            fill_color=BLUE_LT, line_color=BLUE, line_width=Pt(1))
        txt(s, icon + "  " + method, bx + Inches(0.15), Inches(1.38),
            Inches(3.2), Inches(0.35), size=Pt(13), bold=True, color=BLUE)
        txt(s, desc, bx + Inches(0.15), Inches(1.72),
            Inches(3.2), Inches(0.38), size=Pt(10), color=DARK_GRAY)

    txt(s, "↓  _processFile(type, file) 함수 실행",
        Inches(3.5), Inches(2.3), Inches(4.5), Inches(0.38),
        size=Pt(12), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 처리 단계 흐름
    steps = [
        (BLUE,   "① XLSX 파싱",       "FileReader로 ArrayBuffer 읽기\nXLSX.read() → state[type] 메모리 저장"),
        (TEAL,   "② 로컬 저장",       "IndexedDB에 dbPut(type, {data, name})\n브라우저 캐시 갱신"),
        (PURPLE, "③ 서버 저장",       "upload_logs 테이블에 upsert\nfile_data(base64) + 타임스탬프 저장"),
        (GREEN,  "④ UI 갱신",         "✅ 파일명 표시, 업로드 시각 갱신\n업로드자 이름 표시"),
        (ORANGE, "⑤ 자동 분석 실행", "BOM + 생산계획 + 재고현황 3개\n모두 있으면 runAnalysis() 즉시 호출"),
    ]

    sx = Inches(0.3)
    sy = Inches(2.85)
    sw = Inches(2.45)
    sh = Inches(2.5)

    for i, (col, title, desc) in enumerate(steps):
        lx = sx + i * (sw + Inches(0.12))
        box(s, lx, sy, sw, sh,
            fill_color=WHITE, line_color=col, line_width=Pt(1.5))
        box(s, lx, sy, sw, Inches(0.45), fill_color=col)
        txt(s, title, lx + Inches(0.12), sy + Inches(0.07),
            sw - Inches(0.24), Inches(0.35), size=Pt(12), bold=True, color=WHITE)
        ty2 = sy + Inches(0.58)
        for line in desc.split('\n'):
            txt(s, line, lx + Inches(0.15), ty2,
                sw - Inches(0.3), Inches(0.38), size=Pt(10.5), color=DARK_GRAY)
            ty2 += Inches(0.36)
        if i < len(steps) - 1:
            txt(s, "→", lx + sw + Inches(0.0), sy + Inches(1.1),
                Inches(0.14), Inches(0.38), size=Pt(18), bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

    # 하단 서버 실패 처리
    box(s, Inches(0.3), Inches(5.55), Inches(12.7), Inches(0.75),
        fill_color=RGBColor(0xff, 0xf3, 0xcd), line_color=ORANGE, line_width=Pt(1))
    txt(s, "⚠  서버 저장 실패 시",
        Inches(0.5), Inches(5.6), Inches(2.5), Inches(0.32),
        size=Pt(12), bold=True, color=ORANGE)
    txt(s, "로컬(IndexedDB) 저장은 완료 → 본인 기기에서는 정상 분석 가능 | "
           "다른 사용자는 이전 파일 유지 | 화면에 ⚠ 경고 표시",
        Inches(3.0), Inches(5.62), Inches(9.8), Inches(0.6),
        size=Pt(10.5), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════
# 슬라이드 6: 생산계획 DB 저장 방식 (신규)
# ════════════════════════════════════════════════════════════
def slide_db_storage(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "⑤ 생산계획 파일 공유 방식 — DB 직접 저장 (신규)",
               "Supabase Storage RLS 오류 해결 후 적용된 새로운 파일 공유 아키텍처")

    # 이전 방식 (좌)
    box(s, Inches(0.3), Inches(1.25), Inches(5.9), Inches(5.2),
        fill_color=RGBColor(0xfb, 0xf0, 0xef), line_color=RED, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(5.9), Inches(0.45),
        fill_color=RED)
    txt(s, "❌  이전 방식 (Storage 업로드)",
        Inches(0.45), Inches(1.27), Inches(5.6), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)

    old_steps = [
        "1. 파일 선택",
        "2. uploadFileToStorage('plan.xlsx') → ❌ RLS 오류",
        "   new row violates row-level security policy",
        "3. syncUploadLog() → 타임스탬프만 저장 (파일 없음)",
        "4. 다른 사용자: 다운로드 시도 → 구 테스트 파일 수신",
        "5. 결과: 최신 파일이 아닌 오래된 파일로 분석됨",
    ]
    oy = Inches(1.85)
    for step in old_steps:
        col = RED if "❌" in step or "오류" in step or "오래된" in step else DARK_GRAY
        sz  = Pt(10) if step.startswith("  ") else Pt(11)
        txt(s, step, Inches(0.48), oy, Inches(5.6), Inches(0.35),
            size=sz, color=col)
        oy += Inches(0.38)

    # 화살표
    txt(s, "⟹", Inches(6.3), Inches(3.3), Inches(0.7), Inches(0.6),
        size=Pt(30), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, "개선", Inches(6.25), Inches(3.85), Inches(0.8), Inches(0.3),
        size=Pt(11), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # 신규 방식 (우)
    box(s, Inches(7.2), Inches(1.25), Inches(5.8), Inches(5.2),
        fill_color=GREEN_LT, line_color=GREEN, line_width=Pt(2))
    box(s, Inches(7.2), Inches(1.25), Inches(5.8), Inches(0.45),
        fill_color=GREEN)
    txt(s, "✅  신규 방식 (DB 직접 저장)",
        Inches(7.35), Inches(1.27), Inches(5.5), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)

    new_steps = [
        ("1. 파일 선택", GREEN),
        ("2. ArrayBuffer → base64 변환", TEAL),
        ("3. upload_logs.file_data 컬럼에 직접 upsert", TEAL),
        ("   (기존 RLS 정상 동작 테이블 재사용)", MID_GRAY),
        ("4. 타임스탬프·파일명 동시 저장", GREEN),
        ("5. 다른 사용자: DB에서 file_data 직접 조회", GREEN),
        ("6. 결과: 항상 최신 파일로 정확한 분석 ✅", GREEN),
    ]
    ny = Inches(1.85)
    for step, col in new_steps:
        sz = Pt(10) if step.startswith("  ") else Pt(11)
        txt(s, step, Inches(7.38), ny, Inches(5.5), Inches(0.35),
            size=sz, color=col)
        ny += Inches(0.38)

    # DB 스키마 설명
    box(s, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.72),
        fill_color=NAVY, line_color=None)
    txt(s, "upload_logs 테이블 컬럼:  file_type  |  uploaded_at  |  uploader_name  |  "
           "file_name  |  file_names  |  file_data (TEXT · base64)",
        Inches(0.5), Inches(6.62), Inches(12.3), Inches(0.5),
        size=Pt(11), color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════
# 슬라이드 7: 요약 & 핵심 포인트
# ════════════════════════════════════════════════════════════
def slide_summary(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    # 장식
    for ox, oy, ow, col in [
        (Inches(10.5), Inches(-0.5), Inches(5), RGBColor(0x25, 0x4e, 0x7e)),
        (Inches(-1),   Inches(5),    Inches(4), RGBColor(0x18, 0x30, 0x50)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    header_bar(s, "⑥ 요약 & 핵심 포인트", "")

    txt(s, "자동실행 핵심 구성 요약",
        Inches(0.4), Inches(1.2), Inches(12), Inches(0.5),
        size=Pt(18), bold=True, color=WHITE)

    points = [
        (GREEN,  "✅ 로그인 후 자동 파일 복원",
                 "IndexedDB 캐시 → 즉시 화면 표시, 서버 응답 불필요"),
        (BLUE,   "✅ 타임스탬프 기반 스마트 동기화",
                 "서버 > 로컬 일 때만 다운로드 → 불필요한 트래픽 제거"),
        (PURPLE, "✅ 3분 주기 자동 감시",
                 "다른 사용자가 올린 최신 파일 자동 감지 및 반영"),
        (TEAL,   "✅ 업로드 즉시 자동 분석",
                 "BOM + 생산계획 + 재고 3개 파일 준비 완료 시 runAnalysis() 자동 실행"),
        (ORANGE, "✅ DB 직접 저장 방식 (신규)",
                 "Supabase Storage RLS 우회 → upload_logs.file_data(base64) 컬럼 활용"),
        (GREEN,  "✅ 서버 저장 실패 감지",
                 "업로드 실패 시 ⚠ 경고 표시 → 로컬 분석은 유지"),
    ]

    py = Inches(1.85)
    for i, (col, title, desc) in enumerate(points):
        lx = Inches(0.4) if i % 2 == 0 else Inches(6.8)
        box(s, lx, py if i % 2 == 0 else py, Inches(6.0), Inches(0.88),
            fill_color=RGBColor(0x24, 0x4a, 0x78),
            line_color=col, line_width=Pt(1.5))
        txt(s, title, lx + Inches(0.18),
            (py if i % 2 == 0 else py) + Inches(0.06),
            Inches(5.6), Inches(0.32), size=Pt(12), bold=True, color=col)
        txt(s, desc, lx + Inches(0.18),
            (py if i % 2 == 0 else py) + Inches(0.4),
            Inches(5.6), Inches(0.38), size=Pt(10.5), color=RGBColor(0xc0, 0xd8, 0xf0))
        if i % 2 == 1:
            py += Inches(1.02)

    # 하단 한마디
    box(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.65),
        fill_color=BLUE)
    txt(s, "💡  사용자는 파일만 업로드하면 됩니다. 분석·동기화·공유는 모두 자동으로 처리됩니다.",
        Inches(0.6), Inches(6.63), Inches(12.1), Inches(0.5),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 메인
# ════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_login(prs)
    slide_sync(prs)
    slide_upload_trigger(prs)
    slide_db_storage(prs)
    slide_summary(prs)

    out_path = r"C:\Users\조립\Desktop\claude\Material Shortage Status vs. Production Plan\자동실행_구성_설명.pptx"
    prs.save(out_path)
    import sys
    sys.stdout.buffer.write(("저장 완료: " + out_path + "\n").encode("utf-8"))


if __name__ == '__main__':
    main()
