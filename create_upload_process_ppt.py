#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
업로드 파일 진행과정 PPT 생성기
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ────────────────────────────────────────────
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
BLUE       = RGBColor(0x29, 0x52, 0xa3)
GREEN      = RGBColor(0x27, 0xae, 0x60)
RED        = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE     = RGBColor(0xe6, 0x7e, 0x22)
PURPLE     = RGBColor(0x8e, 0x44, 0xad)
TEAL       = RGBColor(0x16, 0x7a, 0x5e)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf2, 0xf5)
DARK_GRAY  = RGBColor(0x44, 0x44, 0x44)
MID_GRAY   = RGBColor(0x77, 0x77, 0x77)
LIGHT_BLUE = RGBColor(0xe8, 0xf0, 0xfa)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=""):
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), NAVY)
    txt(slide, title,
        Inches(0.4), Inches(0.12), Inches(10), Inches(0.6),
        size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.42), Inches(0.72), Inches(10), Inches(0.35),
            size=Pt(13), color=RGBColor(0xb0, 0xc4, 0xde))
    txt(slide, "업로드 파일 진행과정",
        Inches(9.5), Inches(0.18), Inches(3.5), Inches(0.35),
        size=Pt(11), color=RGBColor(0x80, 0xa0, 0xc8), align=PP_ALIGN.RIGHT)


def label(slide, text, l, t, w=Inches(3), color=NAVY):
    box(slide, l, t, w, Inches(0.32), fill_color=color)
    txt(slide, text, l + Inches(0.1), t + Inches(0.03),
        w - Inches(0.2), Inches(0.28),
        size=Pt(11), bold=True, color=WHITE)


def step_box(slide, l, t, w, h, num, title, desc, color):
    box(slide, l, t, w, h, fill_color=WHITE,
        line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1))
    box(slide, l, t, w, Inches(0.42), fill_color=color)
    txt(slide, num, l + Inches(0.1), t + Inches(0.06),
        Inches(0.5), Inches(0.3), size=Pt(13), bold=True, color=WHITE)
    txt(slide, title, l + Inches(0.6), t + Inches(0.06),
        w - Inches(0.7), Inches(0.3), size=Pt(12), bold=True, color=WHITE)
    txt(slide, desc, l + Inches(0.13), t + Inches(0.52),
        w - Inches(0.26), h - Inches(0.6),
        size=Pt(10.5), color=DARK_GRAY)


def footer(slide, text, color=NAVY):
    box(slide, Inches(0.35), Inches(6.4), Inches(12.6), Inches(0.76), fill_color=color)
    txt(slide, text, Inches(0.55), Inches(6.52), Inches(12.1), Inches(0.52),
        size=Pt(11.5), color=WHITE)


# ═══════════════════════════════════════════════════════════
# 슬라이드 1 : 표지
# ═══════════════════════════════════════════════════════════
def slide_cover(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    circ = s.shapes.add_shape(9, Inches(9.5), Inches(-1), Inches(6), Inches(6))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x25, 0x4e, 0x7e)
    circ.line.fill.background()
    circ2 = s.shapes.add_shape(9, Inches(-1.5), Inches(4.5), Inches(5), Inches(5))
    circ2.fill.solid(); circ2.fill.fore_color.rgb = RGBColor(0x18, 0x30, 0x50)
    circ2.line.fill.background()

    txt(s, "📂 업로드 파일 진행과정",
        Inches(1), Inches(1.3), Inches(11), Inches(0.9),
        size=Pt(42), bold=True, color=WHITE)
    txt(s, "파일 업로드 → 파싱 → 저장 → 동기화 → 분석 실행",
        Inches(1), Inches(2.3), Inches(11), Inches(0.6),
        size=Pt(20), color=RGBColor(0x8a, 0xb4, 0xf8))

    box(s, Inches(1), Inches(3.1), Inches(5.5), Pt(2), fill_color=BLUE)

    items = [
        ("📁 업로드 대상", "BOM · 생산계획 · 재고현황 · 포장사양 · 키팅"),
        ("🔄 처리 과정",   "FileReader → SheetJS → IndexedDB → Supabase"),
        ("⚡ 자동 실행",   "업로드 즉시 분석 재실행 & 화면 갱신"),
    ]
    ty = Inches(3.3)
    for lbl, val in items:
        txt(s, f"• {lbl}", Inches(1), ty, Inches(2.5), Inches(0.38),
            size=Pt(13), bold=True, color=RGBColor(0x8a, 0xb4, 0xf8))
        txt(s, val, Inches(3.4), ty, Inches(8.5), Inches(0.38),
            size=Pt(13), color=RGBColor(0xcc, 0xdd, 0xf5))
        ty += Inches(0.48)

    txt(s, "2026 · 조립팀 내부 도구",
        Inches(1), Inches(6.7), Inches(6), Inches(0.4),
        size=Pt(11), color=RGBColor(0x55, 0x75, 0xa0), italic=True)


# ═══════════════════════════════════════════════════════════
# 슬라이드 2 : 업로드 파일 종류 개요
# ═══════════════════════════════════════════════════════════
def slide_file_types(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "01. 업로드 파일 종류", "5가지 입력 파일과 역할")

    files = [
        ("📋 BOM", NAVY,
         "bill_of_materials.xlsx",
         ["필수 파일",
          "시트: BOM정보 / BOM / 첫번째 시트",
          "컬럼: 제품번호·자품번·소요량·레벨",
          "→ 제품별 소요 자재 구성 정의",
          "→ bomMap[제품코드] 딕셔너리 생성"]),
        ("📅 생산계획", BLUE,
         "조립_A/B/C/D.xlsx",
         ["필수 파일",
          "시트: 조립_A, 조립_B, 조립_C, 조립_D",
          "컬럼: 일자·시스템P/N·도번·수량",
          "→ 날짜별 생산 제품·수량 정의",
          "→ 도번↔P/N 매핑 자동 추출"]),
        ("🏭 재고현황", TEAL,
         "inventory.xlsx",
         ["필수 파일",
          "시트: 첫 번째 시트",
          "컬럼: 품목코드·재고수량",
          "→ 현재 창고 보유 수량",
          "→ invMap[품번] 딕셔너리 생성"]),
        ("📦 포장사양", ORANGE,
         "packaging_spec.xlsx",
         ["선택 파일",
          "시트: grdList",
          "컬럼: P/N·발포지·패드·박스·포장지",
          "→ 포장재 소요량 계산 기준",
          "→ specMap[P/N] 딕셔너리 생성"]),
        ("🔧 키팅된 자재", PURPLE,
         "kit_*.xlsx (복수 파일)",
         ["선택 파일",
          "복수 파일 동시 업로드 가능",
          "컬럼: 자재코드·수량 (다양한 양식)",
          "→ 이미 키팅된 수량 차감",
          "→ kitMap[자재코드] 합산 처리"]),
    ]

    cx = Inches(0.25)
    for (title, color, filename, items) in files:
        w = Inches(2.5)
        h = Inches(5.15)
        box(s, cx, Inches(1.25), w, h, fill_color=WHITE,
            line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1.2))
        box(s, cx, Inches(1.25), w, Inches(0.45), fill_color=color)
        txt(s, title, cx + Inches(0.1), Inches(1.3),
            w - Inches(0.2), Inches(0.36), size=Pt(13), bold=True, color=WHITE)
        # 파일명 배경
        box(s, cx + Inches(0.08), Inches(1.76), w - Inches(0.16), Inches(0.36),
            fill_color=RGBColor(0xf0, 0xf4, 0xfa),
            line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(0.5))
        txt(s, filename, cx + Inches(0.12), Inches(1.8),
            w - Inches(0.24), Inches(0.28), size=Pt(9), color=MID_GRAY, italic=True)
        iy = Inches(2.2)
        for item in items:
            clr = RED if item == "필수 파일" else (MID_GRAY if item == "선택 파일" else DARK_GRAY)
            bld = item in ("필수 파일", "선택 파일")
            txt(s, ("✅ " if bld else "• ") + item,
                cx + Inches(0.12), iy, w - Inches(0.24), Inches(0.44),
                size=Pt(10), color=clr, bold=bld)
            iy += Inches(0.44)
        cx += Inches(2.62)

    footer(s, "💡 BOM · 생산계획 · 재고현황 3개 파일이 모두 업로드되어야 [분석 실행] 버튼이 활성화됩니다.")


# ═══════════════════════════════════════════════════════════
# 슬라이드 3 : 업로드 방법
# ═══════════════════════════════════════════════════════════
def slide_how_to_upload(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "02. 파일 업로드 방법", "드래그앤드롭 또는 클릭 선택")

    # 방법 1: 드래그앤드롭
    label(s, "■ 방법 1 : 드래그 앤 드롭", Inches(0.4), Inches(1.35), Inches(3.0), NAVY)
    box(s, Inches(0.4), Inches(1.75), Inches(5.9), Inches(3.5),
        fill_color=WHITE, line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1))

    drop_zone_items = [
        "1. 탐색기에서 Excel 파일 선택",
        "2. 해당 업로드 박스(📋 BOM / 📅 생산계획 등) 위로 드래그",
        "3. 마우스 버튼 놓으면 자동 업로드 시작",
        "4. 박스 테두리가 파란색으로 강조되면 정상 인식",
    ]
    ty = Inches(2.0)
    for item in drop_zone_items:
        txt(s, item, Inches(0.6), ty, Inches(5.5), Inches(0.4),
            size=Pt(12), color=DARK_GRAY)
        ty += Inches(0.42)

    # 드롭존 시각화
    box(s, Inches(0.8), Inches(3.9), Inches(5.1), Inches(1.15),
        fill_color=RGBColor(0xeb, 0xf5, 0xff),
        line_color=BLUE, line_width=Pt(1.5))
    txt(s, "📋 BOM", Inches(1.0), Inches(4.0), Inches(2), Inches(0.35),
        size=Pt(14), bold=True, color=NAVY)
    txt(s, "✅ BOM.xlsx (저장됨)", Inches(1.0), Inches(4.38), Inches(4.5), Inches(0.35),
        size=Pt(11), color=TEAL)
    txt(s, "← 파일을 여기에 드롭", Inches(3.2), Inches(4.0), Inches(2.5), Inches(0.35),
        size=Pt(10), color=MID_GRAY, italic=True)

    # 방법 2: 클릭
    label(s, "■ 방법 2 : 클릭하여 선택", Inches(6.9), Inches(1.35), Inches(3.0), BLUE)
    box(s, Inches(6.9), Inches(1.75), Inches(5.9), Inches(3.5),
        fill_color=WHITE, line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1))

    click_items = [
        "1. 업로드 박스 영역 클릭",
        "2. 파일 탐색기가 열림",
        "3. 해당 Excel 파일 선택 후 열기",
        "4. 자동으로 업로드 처리 시작",
    ]
    ty2 = Inches(2.0)
    for item in click_items:
        txt(s, item, Inches(7.1), ty2, Inches(5.5), Inches(0.4),
            size=Pt(12), color=DARK_GRAY)
        ty2 += Inches(0.42)

    # 클릭 시각화
    box(s, Inches(7.3), Inches(3.9), Inches(5.1), Inches(1.15),
        fill_color=RGBColor(0xf5, 0xf0, 0xff),
        line_color=PURPLE, line_width=Pt(1.5))
    txt(s, "📅 생산계획", Inches(7.5), Inches(4.0), Inches(3), Inches(0.35),
        size=Pt(14), bold=True, color=PURPLE)
    txt(s, "클릭하여 파일 선택", Inches(7.5), Inches(4.38), Inches(4.5), Inches(0.35),
        size=Pt(11), color=MID_GRAY)

    # 공통 안내
    label(s, "■ 업로드 후 표시 정보", Inches(0.4), Inches(5.45), Inches(2.8), TEAL)
    info_items = [
        ("✅ 상태 표시", "파일명 + '(저장됨)' 텍스트로 성공 표시"),
        ("📅 업로드 시각", "업로드한 날짜·시간·담당자명 표시"),
        ("💾 저장 버튼", "생산계획 파일은 💾 저장 버튼으로 로컬 다운로드 가능"),
    ]
    ix = Inches(0.4)
    for (lbl, desc) in info_items:
        box(s, ix, Inches(5.85), Inches(4.1), Inches(0.38),
            fill_color=WHITE, line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(0.8))
        txt(s, lbl, ix + Inches(0.1), Inches(5.9), Inches(1.4), Inches(0.3),
            size=Pt(11), bold=True, color=TEAL)
        txt(s, desc, ix + Inches(1.5), Inches(5.9), Inches(2.5), Inches(0.3),
            size=Pt(10.5), color=DARK_GRAY)
        ix += Inches(4.3)

    footer(s, "📌 키팅 파일은 복수 선택 가능 — Ctrl+클릭으로 여러 파일 동시 업로드, 기존 파일에 추가됨")


# ═══════════════════════════════════════════════════════════
# 슬라이드 4 : 업로드 처리 흐름 (핵심)
# ═══════════════════════════════════════════════════════════
def slide_process_flow(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "03. 업로드 처리 흐름", "파일 선택 후 내부적으로 일어나는 5단계")

    steps = [
        ("①", "FileReader\n읽기", NAVY,
         ["브라우저 FileReader API",
          "파일을 ArrayBuffer로 변환",
          "바이너리 데이터 메모리 로드",
          "readAsArrayBuffer(file) 호출"]),
        ("②", "SheetJS\n파싱", BLUE,
         ["XLSX.read(arrayBuffer)",
          "워크북 객체 생성",
          "시트별 데이터 추출",
          "JSON 배열로 변환 완료"]),
        ("③", "IndexedDB\n저장", TEAL,
         ["dbPut(type, {data, name})",
          "브라우저 로컬 DB 저장",
          "세션 재시작 후 복원 가능",
          "1개 파일만 유지 (덮어쓰기)"]),
        ("④", "Supabase\n업로드", ORANGE,
         ["uploadFileToStorage()",
          "plan.xlsx / bom.xlsx 등",
          "팀 전체 공유 원본 저장",
          "upload_logs 테이블 갱신"]),
        ("⑤", "분석\n자동 실행", RED,
         ["3개 파일 모두 있으면",
          "runAnalysis() 자동 호출",
          "부족현황 즉시 재계산",
          "모든 탭 화면 자동 갱신"]),
    ]

    sx = Inches(0.3)
    for i, (num, title, color, items) in enumerate(steps):
        W, H = Inches(2.4), Inches(4.6)
        box(s, sx, Inches(1.3), W, H, fill_color=WHITE,
            line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1.2))
        box(s, sx, Inches(1.3), W, Inches(0.72), fill_color=color)
        txt(s, num, sx + Inches(0.1), Inches(1.37), Inches(0.45), Inches(0.55),
            size=Pt(20), bold=True, color=WHITE)
        txt(s, title, sx + Inches(0.55), Inches(1.37), W - Inches(0.65), Inches(0.55),
            size=Pt(13), bold=True, color=WHITE)
        iy = Inches(2.18)
        for item in items:
            txt(s, "• " + item, sx + Inches(0.13), iy, W - Inches(0.26), Inches(0.42),
                size=Pt(10.5), color=DARK_GRAY)
            iy += Inches(0.44)
        # 화살표
        if i < 4:
            txt(s, "→", sx + W + Inches(0.02), Inches(3.4), Inches(0.38), Inches(0.4),
                size=Pt(22), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        sx += Inches(2.6)

    # 코드 스니펫
    box(s, Inches(0.3), Inches(6.15), Inches(12.7), Inches(0.55),
        fill_color=RGBColor(0x1a, 0x1a, 0x2e))
    txt(s, "  state[type] = XLSX.read(ab)  →  dbPut(type, {data:ab, name:file.name})  →  uploadFileToStorage(type+'.xlsx', ab)  →  runAnalysis()",
        Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.42),
        size=Pt(10.5), color=RGBColor(0xa8, 0xe6, 0xa8))

    footer(s, "⚡ 파싱·저장·업로드·분석이 모두 단일 FileReader.onload 콜백 안에서 순차 실행됩니다.")


# ═══════════════════════════════════════════════════════════
# 슬라이드 5 : 생산계획 파일 특별 처리
# ═══════════════════════════════════════════════════════════
def slide_plan_special(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "04. 생산계획 파일 특별 처리", "조립_A/B/C/D 다중 시트 + 도번↔P/N 매핑")

    # 시트 구조
    label(s, "■ 다중 시트 구조", Inches(0.4), Inches(1.35), Inches(2.4), NAVY)
    sheets = [
        ("조립_A", NAVY,   ["시스템P/N + 도번 컬럼 모두 포함", "도번↔P/N 매핑 기준 시트"]),
        ("조립_B", BLUE,   ["P/N 기준 생산계획", "도번 컬럼 있으면 매핑 추출"]),
        ("조립_C", TEAL,   ["P/N 기준 생산계획", "도번 컬럼 있으면 매핑 추출"]),
        ("조립_D", ORANGE, ["P/N 컬럼 = 도번값", "dobnMap으로 실제 P/N 변환 필요"]),
    ]
    sx = Inches(0.4)
    for (sname, color, items) in sheets:
        box(s, sx, Inches(1.75), Inches(3.0), Inches(1.85),
            fill_color=WHITE, line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1))
        box(s, sx, Inches(1.75), Inches(3.0), Inches(0.38), fill_color=color)
        txt(s, sname, sx + Inches(0.1), Inches(1.8), Inches(2.8), Inches(0.3),
            size=Pt(13), bold=True, color=WHITE)
        iy = Inches(2.22)
        for item in items:
            txt(s, "• " + item, sx + Inches(0.12), iy, Inches(2.8), Inches(0.4),
                size=Pt(10.5), color=DARK_GRAY)
            iy += Inches(0.42)
        sx += Inches(3.15)

    # 도번 매핑 흐름
    label(s, "■ 도번 ↔ 시스템P/N 자동 매핑 (buildDobnMapFromPlan)", Inches(0.4), Inches(3.82), Inches(5.5), RED)
    flow_items = [
        ("조립_A/B/C/D 전체 시트 순회", NAVY),
        ("헤더 행 자동 감지\n(P/N·시스템·품번 키워드)", BLUE),
        ("도번→P/N\nMap 생성", TEAL),
        ("bomMap에\n별칭 등록", ORANGE),
        ("조립_D 부족현황\n정상 반영", RED),
    ]
    fx = Inches(0.4)
    for i, (ftext, color) in enumerate(flow_items):
        box(s, fx, Inches(4.22), Inches(2.35), Inches(1.2),
            fill_color=color)
        txt(s, ftext, fx + Inches(0.1), Inches(4.35),
            Inches(2.18), Inches(0.95),
            size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, "→", fx + Inches(2.37), Inches(4.72),
                Inches(0.38), Inches(0.38),
                size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        fx += Inches(2.6)

    # 날짜 처리
    label(s, "■ 날짜 파싱 — 2가지 형식 자동 지원", Inches(0.4), Inches(5.65), Inches(3.6), MID_GRAY)
    date_items = [
        ("Excel 날짜 시리얼 (숫자)", "46100 → XLSX.SSF.parse_date_code() → '2026-03-19'"),
        ("한글 텍스트 날짜",         "'3월 19일 (수)' → 정규식으로 월·일 추출 → 연도 자동 계산"),
    ]
    ty = Inches(6.05)
    for (form, proc) in date_items:
        box(s, Inches(0.4), ty, Inches(3.0), Inches(0.34),
            fill_color=NAVY)
        txt(s, form, Inches(0.5), ty + Inches(0.04), Inches(2.85), Inches(0.28),
            size=Pt(10), bold=True, color=WHITE)
        txt(s, "→ " + proc, Inches(3.5), ty + Inches(0.04), Inches(9.2), Inches(0.28),
            size=Pt(10), color=DARK_GRAY)
        ty += Inches(0.37)

    footer(s, "📋 조립_D는 P/N 컬럼 값이 도번(고객사PARTNO)이므로 dobnMap으로 실제 시스템P/N으로 변환 후 BOM 조회합니다.")


# ═══════════════════════════════════════════════════════════
# 슬라이드 6 : 파일 저장 & 복원
# ═══════════════════════════════════════════════════════════
def slide_restore(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "05. 파일 저장 & 자동 복원", "IndexedDB 캐시로 세션 재시작 후 자동 복원")

    # 저장 흐름
    label(s, "■ 저장 흐름 (업로드 시)", Inches(0.4), Inches(1.35), Inches(2.6), NAVY)
    save_steps = [
        ("파일 선택",    "드래그앤드롭\n또는 클릭",           NAVY),
        ("ArrayBuffer",  "FileReader로\n바이너리 변환",       BLUE),
        ("IndexedDB",    "dbPut(key, {data, name})\n로컬 저장", TEAL),
        ("Supabase\nStorage", "클라우드 업로드\n팀 공유",     ORANGE),
        ("LocalStorage", "업로드 시각·\n담당자명 저장",       PURPLE),
    ]
    sx = Inches(0.4)
    for i, (title, desc, color) in enumerate(save_steps):
        box(s, sx, Inches(1.75), Inches(2.3), Inches(1.5),
            fill_color=color)
        txt(s, title, sx + Inches(0.1), Inches(1.88),
            Inches(2.12), Inches(0.42),
            size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, sx + Inches(0.1), Inches(2.32),
            Inches(2.12), Inches(0.82),
            size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, "→", sx + Inches(2.32), Inches(2.35),
                Inches(0.38), Inches(0.38),
                size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        sx += Inches(2.6)

    # 복원 흐름
    label(s, "■ 복원 흐름 (페이지 로드 시)", Inches(0.4), Inches(3.52), Inches(3.0), TEAL)
    restore_steps = [
        "① window.load 이벤트 발생",
        "② ['bom','plan','inv','pkg'] 순서로 IndexedDB 조회",
        "③ 데이터 있으면 XLSX.read() 로 워크북 복원",
        "④ 업로드 박스에 '✅ 파일명 (저장됨)' 표시",
        "⑤ syncFilesFromSupabase() 호출 — 서버 파일이 더 최신이면 자동 덮어쓰기",
        "⑥ 3개 파일 복원 완료 시 분석 실행 버튼 자동 활성화",
    ]
    ty = Inches(3.92)
    for step in restore_steps:
        txt(s, step, Inches(0.5), ty, Inches(12.2), Inches(0.38),
            size=Pt(12), color=DARK_GRAY)
        ty += Inches(0.4)

    # 생산계획 저장 버튼
    label(s, "■ 💾 생산계획 저장 버튼", Inches(0.4), Inches(6.3), Inches(2.6), GREEN)
    txt(s, "업로드/복원 시 자동 표시 → 클릭하면 IndexedDB의 원본 xlsx 파일을 그대로 다운로드",
        Inches(3.1), Inches(6.33), Inches(9.9), Inches(0.35),
        size=Pt(11.5), color=DARK_GRAY)

    footer(s, "🔄 오프라인 상태에서도 IndexedDB 캐시로 분석 가능 — Supabase 연결 실패 시 로컬 데이터 사용")


# ═══════════════════════════════════════════════════════════
# 슬라이드 7 : 기기간 파일 동기화
# ═══════════════════════════════════════════════════════════
def slide_sync(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "06. 기기간 파일 동기화", "A 기기 업로드 → B·C 기기 자동 복원")

    # 동기화 다이어그램
    label(s, "■ 동기화 흐름도", Inches(0.4), Inches(1.35), Inches(2.0), NAVY)

    # 기기 A
    box(s, Inches(0.4), Inches(1.8), Inches(2.5), Inches(2.6),
        fill_color=NAVY, line_color=NAVY)
    txt(s, "💻 기기 A\n(업로더)", Inches(0.5), Inches(2.5),
        Inches(2.32), Inches(0.8), size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Excel 파일 업로드", Inches(0.5), Inches(1.95),
        Inches(2.32), Inches(0.35), size=Pt(10), color=RGBColor(0xb0, 0xc8, 0xff), align=PP_ALIGN.CENTER)

    # 중간 처리
    mid = [
        (Inches(3.15), "📦 IndexedDB",   BLUE,   "로컬 캐시\n(즉시 저장)"),
        (Inches(5.55), "☁️ Supabase\nStorage", TEAL, "클라우드\n원본 저장"),
        (Inches(7.95), "📋 upload_logs", ORANGE, "메타데이터\n(시각·담당자)"),
    ]
    for (mx, mt, color, desc) in mid:
        box(s, mx, Inches(2.0), Inches(2.2), Inches(2.2),
            fill_color=color)
        txt(s, mt, mx + Inches(0.1), Inches(2.15),
            Inches(2.0), Inches(0.55), size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, mx + Inches(0.1), Inches(2.78),
            Inches(2.0), Inches(0.6), size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

    # 기기 B
    box(s, Inches(10.4), Inches(1.8), Inches(2.55), Inches(2.6),
        fill_color=RED)
    txt(s, "💻 기기 B/C\n(자동 복원)", Inches(10.5), Inches(2.5),
        Inches(2.37), Inches(0.8), size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "재접속 시 자동 동기화", Inches(10.5), Inches(1.95),
        Inches(2.37), Inches(0.35), size=Pt(10), color=RGBColor(0xff, 0xcc, 0xcc), align=PP_ALIGN.CENTER)

    # 화살표
    for ax in [Inches(2.92), Inches(5.32), Inches(7.72), Inches(10.17)]:
        txt(s, "→", ax, Inches(2.62), Inches(0.38), Inches(0.38),
            size=Pt(20), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 동기화 상세
    label(s, "■ syncFilesFromSupabase() 동작 방식", Inches(0.4), Inches(4.62), Inches(3.8), TEAL)
    sync_items = [
        "① 페이지 로드 시 upload_logs 테이블에서 각 파일의 최신 업로드 시각 조회",
        "② 로컬(LocalStorage)의 업로드 시각과 비교",
        "③ 서버 파일이 더 최신이면 Supabase Storage에서 다운로드",
        "④ 다운로드한 파일을 IndexedDB에 저장 + UI 자동 갱신",
        "⑤ 3개 파일 복원 완료 → runAnalysis() 자동 실행",
    ]
    ty = Inches(5.02)
    for item in sync_items:
        txt(s, item, Inches(0.5), ty, Inches(12.2), Inches(0.36),
            size=Pt(11.5), color=DARK_GRAY)
        ty += Inches(0.37)

    footer(s, "☁️ 결과: A 기기에서 생산계획 업로드 → B 기기 새로고침 → 자동으로 동일 데이터 복원 → 팀 전원 최신 현황 공유")


# ═══════════════════════════════════════════════════════════
# 슬라이드 8 : 업로드 후 분석 실행
# ═══════════════════════════════════════════════════════════
def slide_analysis(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_GRAY)
    header_bar(s, "07. 업로드 후 분석 실행", "3개 파일 업로드 완료 → 자동 분석 → 결과 표시")

    # 분석 트리거 조건
    label(s, "■ 분석 자동 실행 조건", Inches(0.4), Inches(1.35), Inches(2.6), NAVY)
    box(s, Inches(0.4), Inches(1.75), Inches(5.9), Inches(1.65),
        fill_color=WHITE, line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1))
    cond_items = [
        "• BOM 파일 업로드/복원 완료     (state.bom ✅)",
        "• 생산계획 파일 업로드/복원 완료  (state.plan ✅)",
        "• 재고현황 파일 업로드/복원 완료  (state.inv ✅)",
        "→ 3개 모두 충족 시 runAnalysis() 자동 호출",
    ]
    ty = Inches(1.88)
    for item in cond_items:
        clr = GREEN if "→" in item else DARK_GRAY
        txt(s, item, Inches(0.6), ty, Inches(5.5), Inches(0.36),
            size=Pt(11.5), color=clr, bold=("→" in item))
        ty += Inches(0.37)

    # 분석 결과물
    label(s, "■ 분석 결과 탭", Inches(6.6), Inches(1.35), Inches(2.0), BLUE)
    tabs = [
        ("자재별 부족현황", NAVY,
         "BOM 소요량 × 계획수량 → 재고+키팅 대조\n위험도 🚨🔴🟡✅ 분류 후 테이블 표시"),
        ("모델별 분석", BLUE,
         "모델별 생산 가능 수량 자동 계산\n부족 자재 칩(chip) 표시"),
        ("포장재 소요량", ORANGE,
         "포장사양 × 생산계획 → 발포지/패드/박스 수량\n라인별·일자별·유형별 집계"),
        ("재고현황 조회", TEAL,
         "업로드된 재고 원본 표 표시\n자재 검색 및 수동 수정 가능"),
    ]
    tx = Inches(6.6)
    tw = Inches(3.2)
    for (title, color, desc) in tabs:
        box(s, tx, Inches(1.75), tw, Inches(1.55),
            fill_color=WHITE, line_color=RGBColor(0xcc, 0xd6, 0xe5), line_width=Pt(1))
        box(s, tx, Inches(1.75), tw, Inches(0.36), fill_color=color)
        txt(s, title, tx + Inches(0.1), Inches(1.8),
            tw - Inches(0.2), Inches(0.28), size=Pt(11.5), bold=True, color=WHITE)
        txt(s, desc, tx + Inches(0.12), Inches(2.18),
            tw - Inches(0.24), Inches(1.0), size=Pt(10), color=DARK_GRAY)
        tx += Inches(3.37)

    # 수식
    label(s, "■ 핵심 계산 수식", Inches(0.4), Inches(3.62), Inches(2.2), RED)
    formulas = [
        ("소요량",   "BOM 소요량  ×  생산 계획 수량"),
        ("가용수량",  "재고수량  +  키팅수량"),
        ("부족수량",  "가용수량  −  소요량  (음수 = 부족)"),
        ("재고비율",  "가용수량  ÷  소요량  × 100%"),
        ("위험도",   "< 0 → 부족(적색) / < 120% → 경고(황색) / ≥ 120% → 정상(녹색)"),
    ]
    ty = Inches(4.02)
    for (fname, formula) in formulas:
        box(s, Inches(0.4), ty, Inches(1.7), Inches(0.35), fill_color=NAVY)
        txt(s, fname, Inches(0.5), ty + Inches(0.04),
            Inches(1.55), Inches(0.28), size=Pt(11), bold=True, color=WHITE)
        txt(s, "=  " + formula, Inches(2.18), ty + Inches(0.04),
            Inches(10.5), Inches(0.28), size=Pt(11.5), color=DARK_GRAY)
        ty += Inches(0.38)

    footer(s, "⚡ 파일 하나라도 새로 업로드되면 3개 파일 유무를 체크 후 자동으로 전체 분석을 재실행합니다.")


# ═══════════════════════════════════════════════════════════
# 슬라이드 9 : 전체 흐름 요약
# ═══════════════════════════════════════════════════════════
def slide_summary(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    circ = s.shapes.add_shape(9, Inches(10), Inches(4), Inches(5), Inches(5))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x25, 0x4e, 0x7e)
    circ.line.fill.background()

    txt(s, "업로드 파일 진행과정 요약",
        Inches(1), Inches(0.6), Inches(11), Inches(0.7),
        size=Pt(34), bold=True, color=RGBColor(0x8a, 0xb4, 0xf8))
    box(s, Inches(1), Inches(1.3), Inches(4), Pt(2), fill_color=BLUE)

    # 단계별 요약 플로우
    flow = [
        ("📂 파일 선택",      "드래그앤드롭\n또는 클릭", NAVY),
        ("🔍 SheetJS 파싱",   "Excel → JS\n워크북 객체", BLUE),
        ("💾 IndexedDB 저장", "로컬 캐시\n세션 복원 가능", TEAL),
        ("☁️ Supabase 업로드","팀 공유\n원본 보관", ORANGE),
        ("⚡ 분석 자동 실행", "부족현황\n즉시 갱신", RED),
    ]
    fx = Inches(1)
    for i, (title, desc, color) in enumerate(flow):
        box(s, fx, Inches(1.55), Inches(2.1), Inches(1.5), fill_color=color)
        txt(s, title, fx + Inches(0.1), Inches(1.65),
            Inches(1.92), Inches(0.42), size=Pt(11.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, fx + Inches(0.1), Inches(2.1),
            Inches(1.92), Inches(0.82), size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, "→", fx + Inches(2.12), Inches(2.22),
                Inches(0.38), Inches(0.38),
                size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        fx += Inches(2.28)

    # 핵심 요약
    summaries = [
        ("📁 업로드 대상",  "BOM(필수) · 생산계획(필수) · 재고현황(필수) · 포장사양(선택) · 키팅(선택)"),
        ("🔄 자동 처리",    "FileReader → SheetJS 파싱 → IndexedDB 저장 → Supabase 업로드 → 분석 재실행"),
        ("📅 생산계획 특이사항", "조립_A/B/C/D 다중 시트 · 도번↔P/N 자동 매핑 · 텍스트/시리얼 날짜 모두 지원"),
        ("☁️ 팀 공유",     "A 기기 업로드 → Supabase 동기화 → B·C 기기 재접속 시 자동 복원"),
        ("⚡ 분석 트리거",  "BOM + 생산계획 + 재고현황 3개 모두 있으면 분석 버튼 활성화 & 자동 실행"),
        ("💾 파일 저장",    "생산계획 💾 저장 버튼으로 IndexedDB 원본 xlsx 파일 로컬 다운로드"),
    ]
    ty = Inches(3.35)
    for (lbl, desc) in summaries:
        box(s, Inches(1), ty, Inches(2.3), Inches(0.44),
            fill_color=RGBColor(0x25, 0x4e, 0x7e))
        txt(s, lbl, Inches(1.1), ty + Inches(0.06),
            Inches(2.15), Inches(0.33),
            size=Pt(11), bold=True, color=RGBColor(0x8a, 0xb4, 0xf8))
        txt(s, desc, Inches(3.4), ty + Inches(0.06),
            Inches(9.2), Inches(0.33),
            size=Pt(11), color=WHITE)
        ty += Inches(0.5)

    txt(s, "📂 생산계획 대비 자재부족현황 · 업로드 파일 진행과정",
        Inches(1), Inches(6.7), Inches(11), Inches(0.45),
        size=Pt(13), bold=True, color=RGBColor(0x55, 0x85, 0xcc), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    print("슬라이드 생성 중...")
    slide_cover(prs)       ; print("  ✅ 표지")
    slide_file_types(prs)  ; print("  ✅ 01. 업로드 파일 종류")
    slide_how_to_upload(prs); print("  ✅ 02. 업로드 방법")
    slide_process_flow(prs); print("  ✅ 03. 업로드 처리 흐름")
    slide_plan_special(prs); print("  ✅ 04. 생산계획 특별 처리")
    slide_restore(prs)     ; print("  ✅ 05. 파일 저장 & 복원")
    slide_sync(prs)        ; print("  ✅ 06. 기기간 동기화")
    slide_analysis(prs)    ; print("  ✅ 07. 분석 실행")
    slide_summary(prs)     ; print("  ✅ 요약")

    out = r"C:\Users\조립\Desktop\claude\Material Shortage Status vs. Production Plan\업로드파일_진행과정.pptx"
    prs.save(out)
    print(f"\n✅ PPT 저장 완료: {out}")
    print(f"   총 슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    main()
