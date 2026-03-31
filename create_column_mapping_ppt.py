# -*- coding: utf-8 -*-
"""
자재부족현황 - 생산계획 컬럼 매핑 PPT 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 컬러 팔레트 ──────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE   = RGBColor(0x29, 0x52, 0xA3)
LBLUE  = RGBColor(0xDC, 0xE8, 0xFF)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED    = RGBColor(0xC0, 0x39, 0x2B)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFF, 0xF0, 0x80)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 빈 레이아웃

def add_slide():
    return prs.slides.add_slide(BLANK)

def txb(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def para(tf, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return p

def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def title_slide_bg(slide):
    r = rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
    r2 = rect(slide, 0, 6.2, 13.33, 1.3, fill=BLUE)

def section_header(slide, text, y=0.3):
    r = rect(slide, 0.4, y, 12.53, 0.55, fill=NAVY)
    tb = txb(slide, 0.5, y+0.05, 12, 0.45)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

def page_num(slide, n, total):
    tb = txb(slide, 12.0, 7.1, 1.2, 0.3)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = f"{n} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MGRAY
    p.alignment = PP_ALIGN.RIGHT

def info_box(slide, x, y, w, h, title, body_lines, title_color=NAVY, bg=LBLUE, border=BLUE):
    r = rect(slide, x, y, w, h, fill=bg, line=border)
    tb = txb(slide, x+0.1, y+0.05, w-0.2, 0.35)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color

    tb2 = txb(slide, x+0.12, y+0.42, w-0.24, h-0.5)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    first = True
    for line in body_lines:
        p2 = tf2.add_paragraph() if not first else tf2.paragraphs[0]
        first = False
        p2.text = line
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = DGRAY

TOTAL = 7

# ══════════════════════════════════════════════════════════════════════
# 슬라이드 1 ─ 표지
# ══════════════════════════════════════════════════════════════════════
s1 = add_slide()
title_slide_bg(s1)

tb = txb(s1, 1.2, 1.5, 10.9, 1.2)
tf = tb.text_frame
tf.word_wrap = False
p = tf.add_paragraph()
p.text = "자재부족현황"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

tb2 = txb(s1, 1.2, 2.8, 10.9, 0.8)
tf2 = tb2.text_frame
tf2.word_wrap = False
p2 = tf2.add_paragraph()
p2.text = "생산계획 파일 컬럼 매핑 시스템"
p2.font.size = Pt(28)
p2.font.bold = False
p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)
p2.alignment = PP_ALIGN.CENTER

tb3 = txb(s1, 1.2, 3.7, 10.9, 0.5)
tf3 = tb3.text_frame
tf3.word_wrap = False
p3 = tf3.add_paragraph()
p3.text = "컬럼 자동감지 · 우선순위 매핑 · 시트별 특수처리 설계 문서"
p3.font.size = Pt(15)
p3.font.color.rgb = RGBColor(0x88, 0xAA, 0xCC)
p3.alignment = PP_ALIGN.CENTER

# 목차 박스
box = rect(s1, 2.5, 4.4, 8.3, 1.5, fill=RGBColor(0x24, 0x4E, 0x88), line=RGBColor(0x44, 0x88, 0xCC))
toc_lines = [
    "  ① 컬럼 매핑 개요        ② PLAN_COL_KW 키워드 정의",
    "  ③ 감지 우선순위 로직    ④ 시트별 적용 현황",
    "  ⑤ 조립_D 특수처리       ⑥ UI 표시 패널",
]
tb4 = txb(s1, 2.6, 4.5, 8.1, 1.3)
tf4 = tb4.text_frame
tf4.word_wrap = True
first = True
for line in toc_lines:
    p4 = tf4.paragraphs[0] if first else tf4.add_paragraph()
    first = False
    p4.text = line
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

tb5 = txb(s1, 1.2, 6.3, 10.9, 0.4)
tf5 = tb5.text_frame
tf5.word_wrap = False
p5 = tf5.add_paragraph()
p5.text = "Material Shortage Status vs. Production Plan  |  2026"
p5.font.size = Pt(11)
p5.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
p5.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════════════
# 슬라이드 2 ─ 컬럼 매핑 개요
# ══════════════════════════════════════════════════════════════════════
s2 = add_slide()
section_header(s2, "① 컬럼 매핑 개요  —  왜 필요한가?")
page_num(s2, 2, TOTAL)

# 배경 설명
tb = txb(s2, 0.5, 1.05, 12.3, 0.45)
tf = tb.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "생산계획 엑셀 파일은 라인(조립_A / B / C / D)마다 컬럼명이 달라 고정 인덱스로는 데이터를 읽을 수 없습니다. 자동감지 매핑이 이를 해결합니다."
p.font.size = Pt(12)
p.font.color.rgb = DGRAY

# 문제 → 해결 화살표
problems = [
    ("조립_A", '"시스템 P/N"  / "도번"'),
    ("조립_B", '"시스템 P/N"  / "도번"'),
    ("조립_C", '"시스템 P/N"  / "도번"'),
    ("조립_D", '"품목코드"  / "P/N"'),
]
cols = [
    ("date",  "일자/날짜"),
    ("pn",    "시스템 P/N"),
    ("qty",   "계획수량"),
    ("model", "MODEL"),
    ("spec",  "사양"),
    ("lhrh",  "LH/RH → 사양"),
]
col_colors = [
    RGBColor(0xDC, 0xE8, 0xFF),
    RGBColor(0xD5, 0xF5, 0xE3),
    RGBColor(0xFD, 0xED, 0xD0),
    RGBColor(0xE8, 0xD5, 0xF5),
    RGBColor(0xFF, 0xF0, 0x80),
    RGBColor(0xFF, 0xDD, 0xDD),
]

# 왼쪽: 시트별 원본 컬럼명
rect(s2, 0.4, 1.65, 5.0, 0.4, fill=NAVY)
tb_h = txb(s2, 0.4, 1.67, 5.0, 0.36)
tf_h = tb_h.text_frame
p_h = tf_h.add_paragraph()
p_h.text = "  라인별 원본 컬럼명 (다양)"
p_h.font.size = Pt(12)
p_h.font.bold = True
p_h.font.color.rgb = WHITE

for idx, (sheet, pncols) in enumerate(problems):
    yy = 2.15 + idx * 0.6
    bg = LBLUE if idx % 2 == 0 else RGBColor(0xF0, 0xF4, 0xFF)
    r = rect(s2, 0.4, yy, 5.0, 0.52, fill=bg, line=BLUE)
    tb = txb(s2, 0.5, yy+0.05, 4.8, 0.42)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = f"  {sheet}  →  {pncols}"
    p.font.size = Pt(11.5)
    p.font.color.rgb = NAVY

# 화살표 (중앙)
arr = txb(s2, 5.5, 2.7, 1.5, 0.8)
tf_a = arr.text_frame
p_a = tf_a.add_paragraph()
p_a.text = "자동\n감지"
p_a.font.size = Pt(18)
p_a.font.bold = True
p_a.font.color.rgb = ACCENT
p_a.alignment = PP_ALIGN.CENTER

arr2 = txb(s2, 5.55, 3.4, 1.4, 0.5)
tf_a2 = arr2.text_frame
p_a2 = tf_a2.add_paragraph()
p_a2.text = "▶▶"
p_a2.font.size = Pt(24)
p_a2.font.color.rgb = ACCENT
p_a2.alignment = PP_ALIGN.CENTER

# 오른쪽: 정규화된 컬럼
rect(s2, 7.1, 1.65, 5.8, 0.4, fill=NAVY)
tb_r = txb(s2, 7.1, 1.67, 5.8, 0.36)
tf_r = tb_r.text_frame
p_r = tf_r.add_paragraph()
p_r.text = "  정규화된 컬럼명 (통일)"
p_r.font.size = Pt(12)
p_r.font.bold = True
p_r.font.color.rgb = WHITE

for idx, (key, label) in enumerate(cols):
    yy = 2.15 + idx * 0.55
    r = rect(s2, 7.1, yy, 5.8, 0.48, fill=col_colors[idx], line=BLUE)
    tb = txb(s2, 7.2, yy+0.06, 5.6, 0.36)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = f"  {label}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY


# ══════════════════════════════════════════════════════════════════════
# 슬라이드 3 ─ PLAN_COL_KW 키워드 정의
# ══════════════════════════════════════════════════════════════════════
s3 = add_slide()
section_header(s3, "② PLAN_COL_KW  —  컬럼별 인식 키워드 목록")
page_num(s3, 3, TOTAL)

kw_data = [
    ("📅 일자 (date)",   "일자, 날짜, DATE, 일, 생산일",                         col_colors[0]),
    ("🔑 시스템P/N (pn)","품목코드, 시스템, 제품번호, 제품코드, 품번, PART, PN, P/N", col_colors[1]),
    ("📦 계획수량 (qty)", "계획수량, 조립수량, 계획Qty, Qty, QTY, QUANTITY, 수량\n   (누계/실적/키팅/출고/재고 포함 셀 제외)", col_colors[2]),
    ("🏷️ 모델 (model)",  "MODEL, 모델, 모델명, MODELNO, MODEL NO",               col_colors[3]),
    ("📝 사양 (spec)",   "사양, SPEC, 규격, 제품사양, 차종사양, 사양명, 형식, 타입, TYPE", col_colors[4]),
    ("↔️ LH/RH (lhrh)", "LH/RH, LH-RH, LH_RH, LHRH, L/R, LH, RH, 좌우, 좌/우, Left/Right", col_colors[5]),
]

rect(s3, 0.4, 1.05, 12.53, 0.42, fill=NAVY)
headers = ["컬럼 (canonical)", "인식 키워드 목록 (우선순위 순)"]
widths = [3.0, 9.0]
xpos = 0.45
for h, w in zip(headers, widths):
    tb = txb(s3, xpos, 1.07, w, 0.38)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = h
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    xpos += w + 0.05

for idx, (name, kws, bg) in enumerate(kw_data):
    yy = 1.55 + idx * 0.82
    r = rect(s3, 0.4, yy, 3.0, 0.74, fill=bg, line=BLUE)
    tb = txb(s3, 0.5, yy+0.1, 2.8, 0.55)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = name
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY

    r2 = rect(s3, 3.45, yy, 9.48, 0.74, fill=LGRAY if idx%2==0 else WHITE, line=MGRAY)
    tb2 = txb(s3, 3.55, yy+0.1, 9.3, 0.55)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.add_paragraph()
    p2.text = kws
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = DGRAY


# ══════════════════════════════════════════════════════════════════════
# 슬라이드 4 ─ 감지 우선순위 로직
# ══════════════════════════════════════════════════════════════════════
s4 = add_slide()
section_header(s4, "③ 감지 우선순위 로직  —  detectPlanColIdx()")
page_num(s4, 4, TOTAL)

tb = txb(s4, 0.5, 1.05, 12.3, 0.4)
tf = tb.text_frame
p = tf.add_paragraph()
p.text = "헤더 행의 각 셀을 키워드와 비교할 때 아래 2단계 우선순위를 적용합니다."
p.font.size = Pt(12)
p.font.color.rgb = DGRAY

steps = [
    ("1순위", "완전 일치  (대소문자 무시)",
     '"일자" 키워드 → 셀값 "일자" ✅  /  셀값 "생산일자" ✗',
     GREEN, RGBColor(0xD5, 0xF5, 0xE3)),
    ("2순위", "포함 일치  (Contains)",
     '"시스템" 키워드 → 셀값 "시스템 P/N" ✅  /  완전일치 없을 때만 사용',
     ORANGE, RGBColor(0xFD, 0xED, 0xD0)),
]
for idx, (rank, title, desc, col, bg) in enumerate(steps):
    yy = 1.6 + idx * 1.55
    # 번호 뱃지
    badge = rect(s4, 0.4, yy, 1.2, 1.3, fill=col)
    tb_b = txb(s4, 0.4, yy+0.25, 1.2, 0.8)
    tf_b = tb_b.text_frame
    p_b = tf_b.add_paragraph()
    p_b.text = rank
    p_b.font.size = Pt(20)
    p_b.font.bold = True
    p_b.font.color.rgb = WHITE
    p_b.alignment = PP_ALIGN.CENTER

    # 내용
    body = rect(s4, 1.7, yy, 11.23, 1.3, fill=bg, line=col)
    tb_c = txb(s4, 1.85, yy+0.1, 11.0, 0.45)
    tf_c = tb_c.text_frame
    tf_c.word_wrap = False
    p_c = tf_c.add_paragraph()
    p_c.text = title
    p_c.font.size = Pt(16)
    p_c.font.bold = True
    p_c.font.color.rgb = col

    tb_d = txb(s4, 1.85, yy+0.6, 11.0, 0.6)
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    p_d = tf_d.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = DGRAY

# 수량 예외 박스
rect(s4, 0.4, 4.85, 12.53, 1.05, fill=RGBColor(0xFF, 0xF8, 0xE1), line=ORANGE)
tb_e = txb(s4, 0.55, 4.9, 12.3, 0.35)
tf_e = tb_e.text_frame
p_e = tf_e.add_paragraph()
p_e.text = "⚠️  계획수량 컬럼 특수 처리  —  오탐 방지 제외 키워드"
p_e.font.size = Pt(13)
p_e.font.bold = True
p_e.font.color.rgb = ORANGE

tb_f = txb(s4, 0.55, 5.28, 12.3, 0.55)
tf_f = tb_f.text_frame
p_f = tf_f.add_paragraph()
p_f.text = '셀 값에  "누계", "실적", "키팅", "출고", "재고"  중 하나라도 포함되면 계획수량 컬럼에서 제외  →  실제 계획수량 컬럼만 정확히 감지'
p_f.font.size = Pt(11.5)
p_f.font.color.rgb = DGRAY


# ══════════════════════════════════════════════════════════════════════
# 슬라이드 5 ─ 시트별 적용 현황
# ══════════════════════════════════════════════════════════════════════
s5 = add_slide()
section_header(s5, "④ 시트별 적용 현황  —  조립_A / B / C / D")
page_num(s5, 5, TOTAL)

tb = txb(s5, 0.5, 1.05, 12.3, 0.4)
tf = tb.text_frame
p = tf.add_paragraph()
p.text = '대상 시트 패턴:  /^조립[_ ][A-Za-z](\\s*\\(\\d+\\))?\\s*$/  — 접두어 "조립_" 또는 "조립 "(공백)  모두 지원'
p.font.size = Pt(12)
p.font.color.rgb = DGRAY

# 표 헤더
hcols = ["시트명", "일자 컬럼", "시스템P/N 컬럼", "계획수량", "MODEL", "사양", "LH/RH"]
hwidths = [1.4, 1.7, 2.1, 1.5, 1.4, 1.4, 1.4]
rect(s5, 0.4, 1.55, 12.53, 0.42, fill=NAVY)
xp = 0.45
for h, w in zip(hcols, hwidths):
    tb = txb(s5, xp, 1.57, w, 0.38)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.add_paragraph()
    p.text = h
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    xp += w + 0.03

rows_data = [
    ("조립_A", "일자", "시스템 P/N", "계획수량", "MODEL", "사양", "LH/RH"),
    ("조립_B", "일자", "시스템 P/N", "계획수량", "MODEL", "사양", "LH/RH"),
    ("조립_C", "일자", "시스템 P/N", "계획수량", "MODEL", "사양", "LH/RH"),
    ("조립_D\n(조립 D)", "일자", "품목코드\n→ 시스템P/N", "계획수량", "MODEL", "사양", "LH/RH"),
]
row_colors = [LBLUE, RGBColor(0xF0,0xF6,0xFF), LBLUE, RGBColor(0xFF,0xF0,0xCC)]
for ri, row in enumerate(rows_data):
    yy = 2.05 + ri * 1.05
    rh = 0.98
    xp = 0.4
    for ci, (cell, w) in enumerate(zip(row, hwidths)):
        bg = row_colors[ri]
        if ri == 3 and ci == 2:
            bg = RGBColor(0xFF, 0xE8, 0xA0)  # 강조
        rect(s5, xp, yy, w, rh, fill=bg, line=MGRAY)
        tb = txb(s5, xp+0.05, yy+0.15, w-0.1, rh-0.2)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = cell
        p.font.size = Pt(10.5 if ri != 3 else 10)
        p.font.bold = (ci == 0)
        p.font.color.rgb = NAVY if ci == 0 else (RED if (ri==3 and ci==2) else DGRAY)
        p.alignment = PP_ALIGN.CENTER
        xp += w + 0.03

# 하단 노트
rect(s5, 0.4, 6.3, 12.53, 0.65, fill=RGBColor(0xFF,0xF0,0xCC), line=ORANGE)
tb_n = txb(s5, 0.55, 6.35, 12.3, 0.55)
tf_n = tb_n.text_frame
tf_n.word_wrap = True
p_n = tf_n.add_paragraph()
p_n.text = '★  조립_D만 "품목코드" = 시스템P/N,  "P/N" = 도번(고객사 PARTNO) 으로 컬럼명이 반전됩니다.  자세한 내용은 다음 슬라이드 참조.'
p_n.font.size = Pt(11)
p_n.font.bold = True
p_n.font.color.rgb = RGBColor(0xA0, 0x50, 0x00)


# ══════════════════════════════════════════════════════════════════════
# 슬라이드 6 ─ 조립_D 특수처리
# ══════════════════════════════════════════════════════════════════════
s6 = add_slide()
section_header(s6, "⑤ 조립_D 특수처리  —  품목코드 / P/N 컬럼 반전")
page_num(s6, 6, TOTAL)

# 좌측: 문제 설명
rect(s6, 0.4, 1.05, 5.8, 5.4, fill=RGBColor(0xFF,0xEE,0xEE), line=RED)
tb_l = txb(s6, 0.5, 1.1, 5.6, 0.4)
tf_l = tb_l.text_frame
p_l = tf_l.add_paragraph()
p_l.text = "⛔  수정 전 (버그)"
p_l.font.size = Pt(14)
p_l.font.bold = True
p_l.font.color.rgb = RED

lines_before = [
    "탐색 코드:",
    'const b = r.findIndex(v =>',
    '  v.includes("P/N") ||',
    '  v.includes("품번") ||',
    '  v.includes("시스템"));',
    "",
    "결과:",
    '  → "P/N" 컬럼(=도번) 을 먼저 발견',
    '  → cP = "P/N" 컬럼 인덱스 저장',
    '  → 시스템P/N이 아닌 도번으로 BOM 조회',
    '  → 전체 조립_D 제품 "BOM없음" 오류',
]
tb_lb = txb(s6, 0.55, 1.6, 5.6, 4.7)
tf_lb = tb_lb.text_frame
tf_lb.word_wrap = True
first = True
for line in lines_before:
    p2 = tf_lb.paragraphs[0] if first else tf_lb.add_paragraph()
    first = False
    p2.text = line
    p2.font.size = Pt(10.5)
    is_code = line.startswith('const') or line.startswith('  v.') or line.startswith('  →')
    p2.font.color.rgb = RGBColor(0xCC,0x00,0x00) if is_code else DGRAY

# 중앙 화살표
arr = txb(s6, 6.3, 3.0, 0.8, 0.6)
tf_a = arr.text_frame
p_a = tf_a.add_paragraph()
p_a.text = "▶"
p_a.font.size = Pt(32)
p_a.font.color.rgb = GREEN
p_a.alignment = PP_ALIGN.CENTER

fix = txb(s6, 6.1, 3.55, 1.2, 0.5)
tf_f2 = fix.text_frame
p_f2 = tf_f2.add_paragraph()
p_f2.text = "수정"
p_f2.font.size = Pt(16)
p_f2.font.bold = True
p_f2.font.color.rgb = GREEN
p_f2.alignment = PP_ALIGN.CENTER

# 우측: 수정 후
rect(s6, 7.4, 1.05, 5.5, 5.4, fill=RGBColor(0xEE,0xFF,0xEE), line=GREEN)
tb_r2 = txb(s6, 7.5, 1.1, 5.3, 0.4)
tf_r2 = tb_r2.text_frame
p_r2 = tf_r2.add_paragraph()
p_r2.text = "✅  수정 후"
p_r2.font.size = Pt(14)
p_r2.font.bold = True
p_r2.font.color.rgb = GREEN

lines_after = [
    "탐색 코드:",
    'let b = r.findIndex(v => v === "품목코드");',
    'if (b < 0) b = r.findIndex(v =>',
    '  v.includes("시스템"));',
    'if (b < 0) b = r.findIndex(v =>',
    '  v.includes("P/N") || v.includes("품번"));',
    "",
    "결과:",
    '  → "품목코드" 완전일치 먼저 탐색',
    '  → cP = 품목코드 컬럼 인덱스 저장',
    '  → 올바른 시스템P/N으로 BOM 조회',
    '  → 조립_D 제품 정상 분석 ✅',
]
tb_rb = txb(s6, 7.55, 1.6, 5.25, 4.7)
tf_rb = tb_rb.text_frame
tf_rb.word_wrap = True
first = True
for line in lines_after:
    p3 = tf_rb.paragraphs[0] if first else tf_rb.add_paragraph()
    first = False
    p3.text = line
    p3.font.size = Pt(10.5)
    is_code = line.startswith('let') or line.startswith('if') or line.startswith('  v.') or line.startswith('  →')
    p3.font.color.rgb = RGBColor(0x00, 0x88, 0x00) if is_code else DGRAY

# 하단 원칙
rect(s6, 0.4, 6.55, 12.53, 0.5, fill=LBLUE, line=BLUE)
tb_p = txb(s6, 0.55, 6.6, 12.3, 0.4)
tf_p = tb_p.text_frame
p_p = tf_p.add_paragraph()
p_p.text = "적용 원칙:  buildDobnMapFromPlan()  /  buildPkgData()  두 함수 모두 동일한 우선순위 로직 적용  — 완전일치 → 시스템 포함 → P/N 포함"
p_p.font.size = Pt(11)
p_p.font.bold = True
p_p.font.color.rgb = NAVY


# ══════════════════════════════════════════════════════════════════════
# 슬라이드 7 ─ UI 표시 패널
# ══════════════════════════════════════════════════════════════════════
s7 = add_slide()
section_header(s7, "⑥ UI 표시 패널  —  showPlanColMap() 감지 결과")
page_num(s7, 7, TOTAL)

tb = txb(s7, 0.5, 1.05, 12.3, 0.4)
tf = tb.text_frame
p = tf.add_paragraph()
p.text = "생산계획 파일 업로드 후 각 시트의 컬럼 감지 결과를 아래 형태로 화면에 표시합니다."
p.font.size = Pt(12)
p.font.color.rgb = DGRAY

# UI 패널 모형
panel_sheets = ["조립_A", "조립_B", "조립_C", "조립_D"]
panel_cols = [
    ("📅 일자/날짜", "일자",   col_colors[0]),
    ("🏷️ 모델/MODEL", "MODEL",  col_colors[3]),
    ("📦 계획수량/Qty","계획수량",col_colors[2]),
    ("📝 사양",       "사양",   col_colors[4]),
    ("↔️ LH/RH",     "사양",   col_colors[5]),
]

rect(s7, 0.4, 1.55, 12.53, 5.5, fill=RGBColor(0xF8,0xF9,0xFA), line=MGRAY)
tb_title = txb(s7, 0.55, 1.6, 8.0, 0.4)
tf_title = tb_title.text_frame
p_title = tf_title.add_paragraph()
p_title.text = "📊  생산계획 컬럼 자동감지 결과"
p_title.font.size = Pt(13)
p_title.font.bold = True
p_title.font.color.rgb = NAVY

for si, sname in enumerate(panel_sheets):
    yy = 2.1 + si * 1.2
    tb_sn = txb(s7, 0.55, yy, 2.0, 0.35)
    tf_sn = tb_sn.text_frame
    p_sn = tf_sn.add_paragraph()
    p_sn.text = sname
    p_sn.font.size = Pt(11)
    p_sn.font.bold = True
    p_sn.font.color.rgb = NAVY

    for ci, (label_text, canonical, bg) in enumerate(panel_cols):
        xp = 0.55 + ci * 2.45
        yp = yy + 0.38
        r = rect(s7, xp, yp, 2.35, 0.68, fill=bg, line=BLUE)
        tb_c = txb(s7, xp+0.06, yp+0.05, 2.23, 0.58)
        tf_c = tb_c.text_frame
        tf_c.word_wrap = False
        p_c = tf_c.add_paragraph()
        p_c.text = label_text
        p_c.font.size = Pt(9.5)
        p_c.font.bold = True
        p_c.font.color.rgb = NAVY

        p_c2 = tf_c.add_paragraph()
        p_c2.text = f"→ {canonical}"
        p_c2.font.size = Pt(10)
        p_c2.font.color.rgb = BLUE

# 하단 설명
rect(s7, 0.4, 7.1, 12.53, 0.35, fill=LBLUE, line=BLUE)
tb_bot = txb(s7, 0.55, 7.13, 12.3, 0.28)
tf_bot = tb_bot.text_frame
p_bot = tf_bot.add_paragraph()
p_bot.text = "미감지 시 회색(미감지) 표시  /  감지 성공 시 파란 배경으로 canonical 컬럼명 표시  /  조립_D 포함 4개 시트 모두 동일하게 표시"
p_bot.font.size = Pt(10)
p_bot.font.color.rgb = NAVY

# ── 저장 ────────────────────────────────────────────────────────────
OUT = r"C:\Users\조립\Desktop\claude\Material Shortage Status vs. Production Plan\컬럼매핑_설계문서.pptx"
prs.save(OUT)
print(f"저장 완료: {OUT}")
