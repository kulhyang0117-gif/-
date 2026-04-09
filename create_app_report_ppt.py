#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
생산계획대비 자재부족현황 앱 개발 완료보고서 PPT 생성기
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ───────────────────────────────────────────
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
BLUE       = RGBColor(0x29, 0x52, 0xa3)
BLUE2      = RGBColor(0x29, 0x80, 0xb9)
GREEN      = RGBColor(0x1a, 0x8a, 0x4a)
GREEN_LT   = RGBColor(0xd4, 0xed, 0xda)
RED        = RGBColor(0xc0, 0x39, 0x2b)
RED_LT     = RGBColor(0xfb, 0xea, 0xe8)
ORANGE     = RGBColor(0xe6, 0x7e, 0x22)
ORANGE_LT  = RGBColor(0xfd, 0xf3, 0xe3)
PURPLE     = RGBColor(0x6c, 0x35, 0x9e)
PURPLE_LT  = RGBColor(0xed, 0xe3, 0xf7)
TEAL       = RGBColor(0x0d, 0x86, 0x8c)
TEAL_LT    = RGBColor(0xd1, 0xf0, 0xf2)
BLUE_LT    = RGBColor(0xd6, 0xe4, 0xf7)
GRAY_BG    = RGBColor(0xf4, 0xf6, 0xf9)
WHITE      = RGBColor(0xff, 0xff, 0xff)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
BORDER     = RGBColor(0xcc, 0xd6, 0xe5)
GOLD       = RGBColor(0xf0, 0xb4, 0x29)
CRIMSON    = RGBColor(0x7b, 0x00, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fc=None, lc=None, lw=Pt(0)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:  sh.fill.background()
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else:  sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h,
        size=Pt(12), bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = size; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return txb

def hdr(slide, title, sub=""):
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fc=NAVY)
    txt(slide, title, Inches(0.4), Inches(0.1), Inches(10), Inches(0.58),
        size=Pt(24), bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.42), Inches(0.67), Inches(11.5), Inches(0.35),
            size=Pt(11), color=RGBColor(0xb0, 0xc4, 0xde))
    txt(slide, "생산계획대비 자재부족현황 앱 개발 완료보고서",
        Inches(9.0), Inches(0.12), Inches(4.0), Inches(0.4),
        size=Pt(9.5), color=RGBColor(0x80, 0xa0, 0xc8), align=PP_ALIGN.RIGHT)

def chip(slide, text, l, t, col):
    w = Inches(len(text)*0.105 + 0.35)
    box(slide, l, t, w, Inches(0.31), fc=col)
    txt(slide, text, l+Inches(0.1), t+Inches(0.03),
        w-Inches(0.12), Inches(0.27), size=Pt(9.5), bold=True, color=WHITE)
    return w


# ════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════
def slide_cover(prs):
    s = blank(prs)
    bg(s, NAVY)
    for ox, oy, ow, col in [
        (Inches(9.8),  Inches(-1.3), Inches(7.5), RGBColor(0x25,0x4e,0x7e)),
        (Inches(-2.0), Inches(4.6),  Inches(6.0), RGBColor(0x18,0x30,0x50)),
        (Inches(5.2),  Inches(5.8),  Inches(3.8), RGBColor(0x1e,0x45,0x72)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    # 좌측 바
    box(s, Inches(0), Inches(0), Inches(0.22), SLIDE_H, fc=BLUE)

    # 완료 배지
    box(s, Inches(0.55), Inches(1.05), Inches(2.4), Inches(0.44), fc=GREEN)
    txt(s, "✅  개발 완료 보고", Inches(0.63), Inches(1.08),
        Inches(2.2), Inches(0.38), size=Pt(13), bold=True, color=WHITE)

    txt(s, "생산계획대비 자재부족현황",
        Inches(0.55), Inches(1.68), Inches(12), Inches(0.88),
        size=Pt(40), bold=True, color=WHITE)
    txt(s, "웹 애플리케이션 개발 완료보고서",
        Inches(0.6),  Inches(2.6),  Inches(11), Inches(0.68),
        size=Pt(28), color=GOLD)
    txt(s, "Material Shortage Status vs. Production Plan  —  Web Application",
        Inches(0.6), Inches(3.38), Inches(11), Inches(0.42),
        size=Pt(13.5), color=RGBColor(0xb0,0xc8,0xe0), italic=True)

    box(s, Inches(0.6), Inches(3.92), Inches(6.5), Pt(2), fc=GOLD)

    meta = [
        ("보고 일자",   "2026년 04월 01일"),
        ("개발 유형",   "순수 프론트엔드 SPA (Single-Page Application)"),
        ("주요 기술",   "HTML5 · CSS3 · Vanilla JS · xlsx.js · Supabase · Vercel"),
        ("핵심 기능",   "BOM × 생산계획 → 재고·키팅 비교 분석 & 자재부족 판정"),
        ("연동 시스템", "sMES (조립 자재 Kitting) · Supabase DB/Storage"),
    ]
    ty = Inches(4.1)
    for label, val in meta:
        txt(s, label, Inches(0.65), ty, Inches(2.4), Inches(0.35),
            size=Pt(11), color=GOLD)
        txt(s, val,   Inches(3.0),  ty, Inches(9.5), Inches(0.35),
            size=Pt(11.5), color=RGBColor(0xd0,0xe4,0xf8))
        ty += Inches(0.42)


# ════════════════════════════════════════════════════════════
# 슬라이드 2: 프로젝트 개요 & 목적
# ════════════════════════════════════════════════════════════
def slide_overview(prs):
    s = blank(prs)
    bg(s, GRAY_BG)
    hdr(s, "프로젝트 개요 및 개발 목적",
        "기존 수작업 분석 한계 극복 → 실시간 자재부족 현황 자동 분석 웹앱 구축")

    # 좌: 배경 & 목적
    box(s, Inches(0.3), Inches(1.25), Inches(5.9), Inches(5.55),
        fc=WHITE, lc=NAVY, lw=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(5.9), Inches(0.44), fc=NAVY)
    txt(s, "📋  도입 배경 및 목적",
        Inches(0.46), Inches(1.27), Inches(5.6), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    items = [
        ("기존 문제점",
         "• 생산계획과 BOM·재고를 수동으로 Excel 대조\n"
         "• 품목별 자재부족 계산 반복 수작업 (일 30분+)\n"
         "• 데이터 불일치·누락 오류 발생 빈번"),
        ("개발 목표",
         "• 파일 업로드만으로 자동 분석 — 클릭 1회\n"
         "• BOM 다단계 전개(Lv1/Lv2) 자동 계산\n"
         "• 키팅 자재 반영 후 실질 부족 수량 산출"),
        ("기대 효과",
         "• 분석 시간 30분 → 10초 이내로 단축\n"
         "• 누락·오계산 없는 정확한 부족 현황 파악\n"
         "• 팀원 공유 — Supabase 클라우드 동기화"),
    ]
    ty = Inches(1.82)
    for title, content in items:
        box(s, Inches(0.46), ty, Inches(0.07), Inches(0.3), fc=NAVY)
        txt(s, title, Inches(0.62), ty, Inches(5.4), Inches(0.3),
            size=Pt(11.5), bold=True, color=NAVY)
        txt(s, content, Inches(0.55), ty+Inches(0.32), Inches(5.5), Inches(0.8),
            size=Pt(10.5), color=DARK_GRAY)
        ty += Inches(1.25)

    # 우: 앱 범위 6단계
    box(s, Inches(6.55), Inches(1.25), Inches(6.45), Inches(5.55),
        fc=WHITE, lc=BLUE, lw=Pt(1.5))
    box(s, Inches(6.55), Inches(1.25), Inches(6.45), Inches(0.44), fc=BLUE)
    txt(s, "🎯  앱 기능 범위",
        Inches(6.7), Inches(1.27), Inches(6.2), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    scopes = [
        (BLUE,   "입력",    "4종 + 키팅 다중",
         "BOM · 생산계획 · 재고현황 · 포장사양 · 키팅 xlsx (복수)"),
        (GREEN,  "인증",    "회원제 접근",
         "로그인 / 회원가입 / 관리자 승인 / 권한별 분리"),
        (TEAL,   "분석",    "BOM 전개 계산",
         "BOM×계획 → 필요량 산출 → (재고+키팅) 차감 → 부족 판정"),
        (ORANGE, "결과",    "6개 탭 뷰",
         "자재별·모델별·일자별·재고·모델재고·포장재 소요"),
        (PURPLE, "공유",    "클라우드 동기화",
         "Supabase Storage — 장치 무관 최신 파일 자동 복원"),
        (RED,    "자동화",  "sMES 키팅 연동",
         "키팅 자동화 버튼 → Python 스크립트 → 자동 업로드"),
    ]
    ty = Inches(1.82)
    for col, tag, subtitle, desc in scopes:
        box(s, Inches(6.7), ty, Inches(0.8), Inches(0.56), fc=col)
        txt(s, tag, Inches(6.7), ty+Inches(0.12), Inches(0.8), Inches(0.32),
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, subtitle, Inches(7.6), ty+Inches(0.02), Inches(2.2), Inches(0.28),
            size=Pt(11), bold=True, color=col)
        txt(s, desc, Inches(7.6), ty+Inches(0.28), Inches(5.2), Inches(0.28),
            size=Pt(9.5), color=DARK_GRAY)
        ty += Inches(0.78)


# ════════════════════════════════════════════════════════════
# 슬라이드 3: 시스템 아키텍처
# ════════════════════════════════════════════════════════════
def slide_arch(prs):
    s = blank(prs)
    bg(s, GRAY_BG)
    hdr(s, "시스템 아키텍처",
        "순수 프론트엔드 SPA — 백엔드 서버 불필요 · Supabase BaaS 연동")

    # 레이어 3단
    layers = [
        (NAVY,   "프레젠테이션 레이어",
         "HTML5 · CSS3 · Vanilla JavaScript",
         "단일 파일(자재부족현황.html) SPA — 서버 빌드 없음"),
        (BLUE,   "비즈니스 로직 레이어",
         "xlsx.js (SheetJS) · 인메모리 분석 엔진",
         "BOM 다단계 전개 · 생산계획 파싱 · 자재 부족 계산"),
        (TEAL,   "데이터 레이어",
         "Supabase (PostgreSQL + Storage + Auth)",
         "회원 인증 · 파일 공유 · 업로드 이력 실시간 동기화"),
    ]
    lh = Inches(1.1)
    for i, (col, title, tech, desc) in enumerate(layers):
        ty = Inches(1.25) + i*(lh + Inches(0.12))
        box(s, Inches(0.3), ty, Inches(12.73), lh, fc=WHITE, lc=col, lw=Pt(1.5))
        box(s, Inches(0.3), ty, Inches(2.4), lh, fc=col)
        txt(s, title, Inches(0.38), ty+Inches(0.22),
            Inches(2.24), Inches(0.65), size=Pt(12), bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        box(s, Inches(2.82), ty+Inches(0.14), Inches(3.5), Inches(0.32), fc=RGBColor(0xf0,0xf0,0xf0))
        txt(s, tech, Inches(2.9), ty+Inches(0.15), Inches(3.4), Inches(0.3),
            size=Pt(10.5), bold=True, color=col)
        txt(s, desc, Inches(6.6), ty+Inches(0.25), Inches(6.3), Inches(0.6),
            size=Pt(10.5), color=DARK_GRAY)

    # 컴포넌트 상세
    box(s, Inches(0.3), Inches(4.75), Inches(12.73), Inches(2.5),
        fc=WHITE, lc=BORDER, lw=Pt(1))
    box(s, Inches(0.3), Inches(4.75), Inches(12.73), Inches(0.42), fc=NAVY)
    txt(s, "주요 컴포넌트 & 흐름",
        Inches(0.5), Inches(4.77), Inches(8), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)

    components = [
        (BLUE,   "파일 업로드\n& 파싱",
                 "xlsx.js → ArrayBuffer\nsheet별 데이터 추출"),
        (TEAL,   "컬럼 매핑\n자동 감지",
                 "생산계획 조립_A~D\n헤더 자동 인식"),
        (PURPLE, "BOM 전개\n계산 엔진",
                 "Lv1/Lv2 재귀 전개\n소요량 누적 합산"),
        (ORANGE, "재고·키팅\n비교 분석",
                 "재고+키팅 → 부족량\n긴급도(1/3일) 판정"),
        (GREEN,  "결과 렌더링\n& 필터",
                 "6탭 동적 테이블\n실시간 검색·필터"),
        (RED,    "Supabase\n클라우드 동기",
                 "파일 Storage 업로드\n장치 간 자동 복원"),
    ]

    cw = Inches(2.0)
    cx = Inches(0.42)
    for i, (col, title, desc) in enumerate(components):
        box(s, cx, Inches(5.28), cw, Inches(1.78),
            fc=WHITE, lc=col, lw=Pt(1.2))
        box(s, cx, Inches(5.28), cw, Inches(0.42), fc=col)
        for j, ln in enumerate(title.split('\n')):
            txt(s, ln, cx, Inches(5.32+j*0.24), cw, Inches(0.24),
                size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, ln in enumerate(desc.split('\n')):
            txt(s, ln, cx+Inches(0.08), Inches(5.8+j*0.28),
                cw-Inches(0.16), Inches(0.26),
                size=Pt(9.5), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < len(components)-1:
            txt(s, "→", cx+cw, Inches(5.95), Inches(0.1), Inches(0.35),
                size=Pt(14), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        cx += cw + Inches(0.12)


# ════════════════════════════════════════════════════════════
# 슬라이드 4: 데이터 입력 & 인증
# ════════════════════════════════════════════════════════════
def slide_input_auth(prs):
    s = blank(prs)
    bg(s, GRAY_BG)
    hdr(s, "주요 기능 (1) — 데이터 입력 & 사용자 인증",
        "5종 Excel 파일 업로드 · 드래그앤드롭 · Supabase 회원제 인증")

    # 입력 파일 (좌상)
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(3.0),
        fc=WHITE, lc=NAVY, lw=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(0.44), fc=NAVY)
    txt(s, "📂  업로드 파일 구성",
        Inches(0.46), Inches(1.27), Inches(7.2), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    files = [
        (NAVY,   "🗂️", "BOM",       "필수", "부품 소요량 정보\n부품코드·부품명·소요량·유형(구매/사출/도장)·레벨"),
        (BLUE,   "📅", "생산계획",  "필수", "조립_A / B / C / D 시트 포함\n라인별 모델번호·수량·생산일자"),
        (TEAL,   "🏭", "재고현황",  "필수", "sMES 창고별 부품 현재고 조회 다운로드\n부품코드·재고수량"),
        (ORANGE, "📋", "포장사양",  "선택", "제품별 포장재 사양 마스터\n포장재 소요량 계산 탭 활성화"),
        (PURPLE, "🔑", "키팅 자재", "선택·복수", "sMES 키팅 xlsx 복수 등록\n키팅 수량 분석 반영 (자동화 연동)"),
    ]
    ty = Inches(1.8)
    for col, icon, name, badge, desc in files:
        box(s, Inches(0.46), ty, Inches(0.52), Inches(0.52), fc=col)
        txt(s, icon, Inches(0.46), ty+Inches(0.1), Inches(0.52), Inches(0.35),
            size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)
        box(s, Inches(1.08), ty+Inches(0.02), Inches(0.9), Inches(0.22), fc=col)
        txt(s, name, Inches(1.1), ty+Inches(0.03), Inches(0.88), Inches(0.2),
            size=Pt(9.5), bold=True, color=WHITE)
        box(s, Inches(2.06), ty+Inches(0.02), Inches(0.88), Inches(0.22),
            fc=GREEN if badge=="필수" else ORANGE if badge=="선택" else PURPLE)
        txt(s, badge, Inches(2.08), ty+Inches(0.03), Inches(0.86), Inches(0.2),
            size=Pt(8.5), bold=True, color=WHITE)
        for j, ln in enumerate(desc.split('\n')):
            txt(s, ln, Inches(1.08), ty+Inches(0.26+j*0.2),
                Inches(6.5), Inches(0.2), size=Pt(9), color=DARK_GRAY)
        ty += Inches(0.54)

    # 업로드 특징 (우상)
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(3.0),
        fc=WHITE, lc=TEAL, lw=Pt(1.5))
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(0.44), fc=TEAL)
    txt(s, "⚡  업로드 UX 특징",
        Inches(8.26), Inches(1.27), Inches(4.6), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    ux = [
        ("드래그 & 드롭",     "파일을 업로드 영역에 끌어다 놓기 지원"),
        ("클라우드 자동 복원", "Supabase Storage — 타 장치·재접속 시\n이전 업로드 파일 자동 복원"),
        ("업로드 이력 표시",   "upload_logs 테이블에 파일명·시각·담당자 저장"),
        ("섹션 접기/펼치기",  "업로드 패널 클릭 토글 (작업 후 화면 절약)"),
        ("키팅 복수 등록",    "xlsx 여러 파일 동시 선택 → chip 태그 표시"),
    ]
    ty = Inches(1.82)
    for title, desc in ux:
        box(s, Inches(8.22), ty, Inches(0.06), Inches(0.28), fc=TEAL)
        txt(s, title, Inches(8.36), ty, Inches(4.5), Inches(0.26),
            size=Pt(10.5), bold=True, color=TEAL)
        txt(s, desc, Inches(8.36), ty+Inches(0.26), Inches(4.5), Inches(0.35),
            size=Pt(9.5), color=DARK_GRAY)
        ty += Inches(0.53)

    # 인증 시스템 (하단)
    box(s, Inches(0.3), Inches(4.42), Inches(12.73), Inches(2.83),
        fc=WHITE, lc=PURPLE, lw=Pt(1.5))
    box(s, Inches(0.3), Inches(4.42), Inches(12.73), Inches(0.44), fc=PURPLE)
    txt(s, "🔐  Supabase 회원 인증 & 권한 관리",
        Inches(0.46), Inches(4.44), Inches(12), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    auth_steps = [
        (PURPLE, "회원가입 신청",
                 "이름·이메일·PW 입력\n→ Supabase Auth 생성\n→ profiles 테이블 등록\n(status=pending)"),
        (ORANGE, "관리자 승인",
                 "관리자 패널에서\n회원 목록 확인\n→ 승인/거절 처리\n→ 권한(읽기/업로드/수정) 부여"),
        (BLUE,   "로그인 & 접근",
                 "이메일+PW 로그인\n→ status=approved 확인\n→ 권한별 UI 활성화\n→ 헤더에 이름 표시"),
        (GREEN,  "관리자 패널",
                 "is_admin=true 시\n회원관리 버튼 표시\n→ 전체 회원 조회\n→ 승인·거절·권한 편집"),
    ]
    aw = Inches(2.9)
    ax = Inches(0.42)
    for col, title, desc in auth_steps:
        box(s, ax, Inches(4.98), aw, Inches(2.12),
            fc=WHITE, lc=col, lw=Pt(1))
        box(s, ax, Inches(4.98), aw, Inches(0.36), fc=col)
        txt(s, title, ax+Inches(0.1), Inches(5.0), aw-Inches(0.2), Inches(0.32),
            size=Pt(11), bold=True, color=WHITE)
        for j, ln in enumerate(desc.split('\n')):
            txt(s, ln, ax+Inches(0.12), Inches(5.4+j*0.3),
                aw-Inches(0.24), Inches(0.28), size=Pt(9.5), color=DARK_GRAY)
        ax += aw + Inches(0.37)


# ════════════════════════════════════════════════════════════
# 슬라이드 5: 분석 로직 & 결과 화면
# ════════════════════════════════════════════════════════════
def slide_analysis(prs):
    s = blank(prs)
    bg(s, GRAY_BG)
    hdr(s, "주요 기능 (2) — 분석 로직 & 결과 화면",
        "BOM 전개 → 필요량 산출 → 재고·키팅 차감 → 긴급도 판정 → 6탭 결과")

    # 분석 파이프라인
    pipeline = [
        (NAVY,   "①", "BOM\n전개",
                 "Lv1/Lv2\n선택 전개\n부품별 소요\n계수 합산"),
        (BLUE,   "②", "생산계획\n파싱",
                 "조립_A~D\n시트 통합\n기준일 이후\n수량 집계"),
        (TEAL,   "③", "필요량\n산출",
                 "BOM계수\n× 생산수량\n= 총 필요\n부품량"),
        (ORANGE, "④", "재고·키팅\n차감",
                 "재고+키팅\n→ 필요량\n차감 → 부족\n수량 산출"),
        (RED,    "⑤", "긴급도\n판정",
                 "일자별 확인\n1일이내 🚨\n3일이내 🔴\n3일이후 📋"),
        (PURPLE, "⑥", "결과\n렌더링",
                 "6탭 테이블\n필터·검색\nExcel 내보내기"),
    ]
    pw = Inches(2.05)
    ph = Inches(2.15)
    px = Inches(0.3)
    for i, (col, num, title, desc) in enumerate(pipeline):
        box(s, px, Inches(1.25), pw, ph, fc=WHITE, lc=col, lw=Pt(1.8))
        box(s, px, Inches(1.25), pw, Inches(0.5), fc=col)
        txt(s, num, px, Inches(1.27), Inches(0.45), Inches(0.45),
            size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, ln in enumerate(title.split('\n')):
            txt(s, ln, px+Inches(0.45), Inches(1.3+j*0.24),
                pw-Inches(0.5), Inches(0.24),
                size=Pt(12), bold=True, color=WHITE)
        for j, ln in enumerate(desc.split('\n')):
            txt(s, ln, px+Inches(0.1), Inches(1.82+j*0.3),
                pw-Inches(0.2), Inches(0.28),
                size=Pt(9.5), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < len(pipeline)-1:
            txt(s, "→", px+pw, Inches(2.05), Inches(0.1), Inches(0.38),
                size=Pt(14), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        px += pw + Inches(0.12)

    # 6탭 결과 (하단 좌)
    box(s, Inches(0.3), Inches(3.55), Inches(7.3), Inches(3.7),
        fc=WHITE, lc=BLUE, lw=Pt(1.5))
    box(s, Inches(0.3), Inches(3.55), Inches(7.3), Inches(0.44), fc=BLUE)
    txt(s, "📊  결과 화면 6개 탭",
        Inches(0.46), Inches(3.57), Inches(7.0), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    tabs = [
        (BLUE,   "📊 자재별 부족현황",
                 "부품별 재고·필요량·부족량·긴급도·키팅수량 통합 표시\n"
                 "부품코드/명 검색 · 상태/라인/긴급도 필터"),
        (NAVY,   "🏭 모델별 부족현황",
                 "모델별 자재 부족 집계 + 상세 펼치기(▶)\n"
                 "일자별 소요량 테이블 / 부족 자재 chip 표시"),
        (TEAL,   "📅 일자별 부족현황",
                 "날짜 기준 생산계획 대비 일별 자재 소요 추이"),
        (ORANGE, "📦 재고현황",
                 "전체 부품 재고 수량 조회 및 검색"),
        (PURPLE, "🗂️ 모델별 재고현황",
                 "모델×부품 매트릭스 재고 조회"),
        (GREEN,  "📦 포장재 소요량",
                 "포장사양 파일 기반 라인별 포장재 필요량 산출"),
    ]
    ty = Inches(4.1)
    for col, title, desc in tabs:
        box(s, Inches(0.44), ty, Inches(0.08), Inches(0.26), fc=col)
        txt(s, title, Inches(0.6), ty, Inches(2.5), Inches(0.26),
            size=Pt(10.5), bold=True, color=col)
        txt(s, desc, Inches(0.6), ty+Inches(0.26), Inches(6.8), Inches(0.32),
            size=Pt(9.5), color=DARK_GRAY)
        ty += Inches(0.5)

    # 요약 카드 & 필터 (하단 우)
    box(s, Inches(7.9), Inches(3.55), Inches(5.1), Inches(3.7),
        fc=WHITE, lc=ORANGE, lw=Pt(1.5))
    box(s, Inches(7.9), Inches(3.55), Inches(5.1), Inches(0.44), fc=ORANGE)
    txt(s, "📈  요약 카드 & 분석 설정",
        Inches(8.06), Inches(3.57), Inches(4.8), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    cards = [
        (CRIMSON, "🚨 즉시부족",  "1일이내 생산에 직접 영향"),
        (RED,     "🔴 긴급부족",  "영업일 3일이내 부족 자재"),
        (MID_GRAY,"📋 여유부족",  "3일이후 부족 — 발주 필요"),
        (ORANGE,  "⚠️ 위험자재",  "충족률 120% 미만 — 주의"),
        (GREEN,   "✅ 충분",      "현재고로 계획 대응 가능"),
        (PURPLE,  "🟣 키팅적용",  "키팅 자재 반영 품목 수"),
        (BLUE2,   "📦 총 부품",   "분석 대상 전체 부품 수"),
    ]
    ty = Inches(4.1)
    for col, name, desc in cards:
        box(s, Inches(8.06), ty, Inches(0.22), Inches(0.22), fc=col)
        txt(s, name, Inches(8.36), ty, Inches(2.2), Inches(0.22),
            size=Pt(10), bold=True, color=col)
        txt(s, desc, Inches(8.36), ty+Inches(0.22), Inches(4.5), Inches(0.2),
            size=Pt(9), color=DARK_GRAY)
        ty += Inches(0.46)

    txt(s, "설정 옵션:  기준일자 · BOM 레벨(전체/Lv1/Lv1+2) · 부품유형(구매/사출/도장)",
        Inches(8.06), Inches(7.1), Inches(4.9), Inches(0.28),
        size=Pt(9.5), color=TEAL, italic=True)


# ════════════════════════════════════════════════════════════
# 슬라이드 6: DB 설계 & 배포 구성
# ════════════════════════════════════════════════════════════
def slide_db_deploy(prs):
    s = blank(prs)
    bg(s, GRAY_BG)
    hdr(s, "DB 설계 & 배포 구성",
        "Supabase PostgreSQL 스키마 · RLS 보안 · Vercel 배포")

    # DB 스키마 (좌)
    box(s, Inches(0.3), Inches(1.25), Inches(7.3), Inches(5.55),
        fc=WHITE, lc=TEAL, lw=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(7.3), Inches(0.44), fc=TEAL)
    txt(s, "🗄️  Supabase DB 스키마",
        Inches(0.46), Inches(1.27), Inches(7.0), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    tables = [
        ("profiles", NAVY, [
            ("id",           "UUID",       "PK · auth.users 참조"),
            ("email",        "TEXT",       "로그인 이메일"),
            ("display_name", "TEXT",       "표시 이름"),
            ("status",       "TEXT",       "pending / approved / rejected"),
            ("is_admin",     "BOOLEAN",    "관리자 여부"),
            ("permissions",  "JSONB",      "{read, upload, modify}"),
            ("created_at",   "TIMESTAMPTZ","가입 일시"),
            ("approved_at",  "TIMESTAMPTZ","승인 일시"),
        ]),
        ("upload_logs", BLUE, [
            ("file_type",    "TEXT",    "PK · bom/plan/inv/pkg/kit"),
            ("uploaded_at",  "TEXT",    "Unix ms timestamp"),
            ("uploader_name","TEXT",    "업로더 이름"),
            ("file_name",    "TEXT",    "단일 파일명"),
            ("file_names",   "TEXT[]",  "복수 파일명 (키팅)"),
        ]),
    ]

    ty = Inches(1.82)
    for tname, col, cols in tables:
        box(s, Inches(0.44), ty, Inches(7.0), Inches(0.36), fc=col)
        txt(s, f"TABLE: {tname}", Inches(0.54), ty+Inches(0.04),
            Inches(6.8), Inches(0.28), size=Pt(11), bold=True, color=WHITE)
        ty += Inches(0.38)
        for cname, ctype, cdesc in cols:
            box(s, Inches(0.54), ty, Inches(1.5), Inches(0.28),
                fc=RGBColor(0xf0,0xf0,0xf0))
            txt(s, cname, Inches(0.56), ty+Inches(0.03),
                Inches(1.48), Inches(0.24), size=Pt(9.5), bold=True, color=col)
            box(s, Inches(2.1), ty, Inches(1.0), Inches(0.28), fc=ORANGE_LT)
            txt(s, ctype, Inches(2.12), ty+Inches(0.03),
                Inches(0.98), Inches(0.24), size=Pt(9), color=ORANGE)
            txt(s, cdesc, Inches(3.18), ty+Inches(0.03),
                Inches(4.0), Inches(0.24), size=Pt(9.5), color=DARK_GRAY)
            ty += Inches(0.3)
        ty += Inches(0.18)

    # RLS 정책
    box(s, Inches(0.44), ty, Inches(7.0), Inches(0.3), fc=GREEN_LT, lc=GREEN, lw=Pt(1))
    txt(s, "🔒 RLS:  인증된 사용자만 profiles / upload_logs 읽기·쓰기 허용",
        Inches(0.54), ty+Inches(0.03), Inches(6.8), Inches(0.26),
        size=Pt(9.5), color=GREEN)

    # 배포 구성 (우)
    box(s, Inches(7.9), Inches(1.25), Inches(5.1), Inches(5.55),
        fc=WHITE, lc=PURPLE, lw=Pt(1.5))
    box(s, Inches(7.9), Inches(1.25), Inches(5.1), Inches(0.44), fc=PURPLE)
    txt(s, "🚀  배포 & 인프라 구성",
        Inches(8.06), Inches(1.27), Inches(4.8), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    infra = [
        (PURPLE, "Vercel 배포",
                 "vercel.json — SPA 리다이렉트 설정\n"
                 "모든 경로 → index.html (싱글 라우팅)\n"
                 "GitHub Push → 자동 CI/CD 배포"),
        (TEAL,   "Supabase Auth",
                 "이메일+PW 인증 (이메일 확인 OFF)\n"
                 "profiles 트리거로 가입 시 자동 등록\n"
                 "로그인 유지 (localStorage 세션)"),
        (BLUE,   "Supabase Storage",
                 "버킷: ms-files\n"
                 "파일 타입별 덮어쓰기 (upsert)\n"
                 "RLS: 인증 회원만 읽기·쓰기"),
        (ORANGE, "로컬 실행",
                 "자재부족현황.html 직접 열기\n"
                 "Playwright 키팅 자동화 연동\n"
                 "인터넷 없이도 분석 가능"),
    ]

    ty = Inches(1.82)
    for col, title, desc in infra:
        box(s, Inches(8.06), ty, Inches(4.8), Inches(1.18),
            fc=GRAY_BG, lc=col, lw=Pt(1.2))
        box(s, Inches(8.06), ty, Inches(4.8), Inches(0.32), fc=col)
        txt(s, title, Inches(8.16), ty+Inches(0.03),
            Inches(4.6), Inches(0.28), size=Pt(11), bold=True, color=WHITE)
        for j, ln in enumerate(desc.split('\n')):
            txt(s, f"• {ln}", Inches(8.16), ty+Inches(0.38+j*0.26),
                Inches(4.6), Inches(0.24), size=Pt(9.5), color=DARK_GRAY)
        ty += Inches(1.28)


# ════════════════════════════════════════════════════════════
# 슬라이드 7: 화면 구성 (UI/UX)
# ════════════════════════════════════════════════════════════
def slide_ui(prs):
    s = blank(prs)
    bg(s, GRAY_BG)
    hdr(s, "화면 구성 (UI/UX)",
        "반응형 레이아웃 · 모달 팝업 · 색상 코드 시각화 · 키팅 자동화 버튼")

    screens = [
        (NAVY,   "🔐 로그인 화면",
                 "• 이메일/PW 로그인\n"
                 "• 회원가입 신청 모달\n"
                 "• 관리자 승인 대기 안내\n"
                 "• 아이디·PW 저장 체크박스"),
        (BLUE,   "📂 파일 업로드 패널",
                 "• 4종 업로드 그리드 (2×2)\n"
                 "• 드래그앤드롭 + 클릭 선택\n"
                 "• 키팅 복수 파일 chip 등록\n"
                 "• 섹션 접기/펼치기 토글"),
        (TEAL,   "⚙️ 설정 & 실행 바",
                 "• 기준일자 달력 선택\n"
                 "• BOM 레벨 드롭다운\n"
                 "• 부품유형 필터\n"
                 "• 분석실행·Excel 내보내기 버튼"),
        (ORANGE, "📊 요약 카드 (7종)",
                 "• 즉시부족/긴급/여유/위험\n"
                 "• 충분/키팅적용/총 부품 수\n"
                 "• 색상 코드 직관적 표시\n"
                 "• 긴급도 기준 날짜 배너"),
        (PURPLE, "📋 결과 테이블",
                 "• 6탭 전환 (클릭)\n"
                 "• 컬럼 너비 드래그 조절\n"
                 "• 모델별 상세 펼치기(▶)\n"
                 "• 부족 자재 chip 시각화"),
        (RED,    "🤖 키팅 자동화 모달",
                 "• 5단계 안내 표시\n"
                 "• 자동화 시작 버튼\n"
                 "• 최초 설정 bat 다운로드\n"
                 "• 긴급정지 버튼 (헤더)"),
        (GREEN,  "👥 관리자 패널",
                 "• 전체 회원 카드 목록\n"
                 "• 상태 배지 (대기/승인/거절)\n"
                 "• 승인·거절 버튼\n"
                 "• 권한(읽기/업로드/수정) 체크박스"),
        (CRIMSON,"📱 반응형 레이아웃",
                 "• 1200px 이하 업로드 2열\n"
                 "• 900px 이하 1열 + 카드 2열\n"
                 "• 테이블 가로 스크롤\n"
                 "• 560px 최대 높이 스크롤"),
    ]

    cols = 4
    sw = Inches(3.08)
    sh = Inches(2.75)
    sx = Inches(0.3)
    sy = Inches(1.25)

    for i, (col, title, desc) in enumerate(screens):
        row = i // cols
        ci  = i % cols
        lx = sx + ci*(sw+Inches(0.12))
        ty = sy + row*(sh+Inches(0.1))
        box(s, lx, ty, sw, sh, fc=WHITE, lc=col, lw=Pt(1.5))
        box(s, lx, ty, sw, Inches(0.4), fc=col)
        txt(s, title, lx+Inches(0.12), ty+Inches(0.04),
            sw-Inches(0.24), Inches(0.34), size=Pt(11.5), bold=True, color=WHITE)
        for j, ln in enumerate(desc.split('\n')):
            txt(s, ln, lx+Inches(0.14), ty+Inches(0.48+j*0.34),
                sw-Inches(0.28), Inches(0.32), size=Pt(9.5), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════
# 슬라이드 8: 결론 & 향후 계획
# ════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    s = blank(prs)
    bg(s, NAVY)

    for ox, oy, ow, col in [
        (Inches(10.0), Inches(-1.0), Inches(6.0), RGBColor(0x25,0x4e,0x7e)),
        (Inches(-2.0), Inches(4.5),  Inches(5.0), RGBColor(0x18,0x30,0x50)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    hdr(s, "개발 성과 요약 & 향후 계획", "")

    # 성과 (좌)
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(5.55),
        fc=RGBColor(0x1e,0x30,0x52), lc=GOLD, lw=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(0.44), fc=GOLD)
    txt(s, "🏆  개발 성과 요약",
        Inches(0.46), Inches(1.27), Inches(7.2), Inches(0.4),
        size=Pt(13), bold=True, color=NAVY)

    achievements = [
        (GREEN,  "순수 프론트엔드 SPA 완성",
                 "별도 서버·빌드 없이 HTML 단일 파일로 완전한 앱 구현\n"
                 "로컬 실행 + Vercel 배포 모두 지원"),
        (TEAL,   "BOM × 생산계획 자동 분석",
                 "Lv1/Lv2 BOM 다단계 전개·소요량 계산 엔진 내재화\n"
                 "수동 Excel 작업 30분 → 자동 분석 10초 이내"),
        (BLUE,   "Supabase 클라우드 연동",
                 "회원 인증·권한 관리·파일 공유 통합 구현\n"
                 "다수 장치에서 최신 파일 자동 복원"),
        (PURPLE, "sMES 키팅 자동화 연동",
                 "Python 스크립트(pywinauto·Playwright) 연동\n"
                 "웹 버튼 클릭 → 자동화 실행 → 결과 자동 반영"),
        (ORANGE, "6탭 결과 화면 & Excel 내보내기",
                 "자재별·모델별·일자별·재고·포장재 등 다각도 분석\n"
                 "결과 Excel 다운로드로 보고서 활용 가능"),
    ]
    ty = Inches(1.82)
    for col, title, desc in achievements:
        box(s, Inches(0.46), ty, Inches(0.28), Inches(0.28), fc=col)
        txt(s, title, Inches(0.84), ty, Inches(6.8), Inches(0.28),
            size=Pt(11), bold=True, color=col)
        txt(s, desc, Inches(0.84), ty+Inches(0.3), Inches(6.8), Inches(0.45),
            size=Pt(10), color=RGBColor(0xc8,0xdc,0xf0))
        ty += Inches(0.92)

    # 향후 계획 (우)
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(5.55),
        fc=RGBColor(0x1a,0x35,0x55), lc=BLUE, lw=Pt(1))
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(0.44), fc=BLUE)
    txt(s, "🔧  향후 개선 계획",
        Inches(8.26), Inches(1.27), Inches(4.6), Inches(0.4),
        size=Pt(13), bold=True, color=WHITE)

    plans = [
        ("단기",  GOLD,   [
            "생산계획 컬럼 자동 매핑 고도화",
            "키팅 자동화 조회 버튼 완전 자동화",
            "분석 결과 이메일 자동 발송",
        ]),
        ("중기",  TEAL,   [
            "스케줄러 — 매일 자동 분석 & 알림",
            "부족 자재 발주 요청 자동 생성",
            "모바일 최적화 반응형 강화",
        ]),
        ("장기",  PURPLE, [
            "sMES API 직접 연동 (GUI 의존 제거)",
            "AI 기반 재고 소진 예측 기능",
            "대시보드 & KPI 리포트 고도화",
        ]),
    ]
    ty = Inches(1.82)
    for period, col, items in plans:
        box(s, Inches(8.26), ty, Inches(1.1), Inches(0.34), fc=col)
        txt(s, period, Inches(8.28), ty+Inches(0.04),
            Inches(1.06), Inches(0.28), size=Pt(11), bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            txt(s, f"→  {item}", Inches(9.46), ty+Inches(0.04+j*0.3),
                Inches(3.4), Inches(0.28), size=Pt(10.5),
                color=RGBColor(0xc8,0xdc,0xf0))
        ty += Inches(1.22)

    # 마무리 문구
    box(s, Inches(0.3), Inches(7.1), Inches(12.73), Inches(0.3),
        fc=RGBColor(0x15,0x28,0x44))
    txt(s, "본 앱은 생산 현장의 자재 수급 의사결정을 실시간으로 지원하기 위해 개발되었습니다.",
        Inches(0.5), Inches(7.11), Inches(12.0), Inches(0.26),
        size=Pt(10), color=RGBColor(0x90,0xb8,0xe0),
        align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════
# 메인
# ════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_arch(prs)
    slide_input_auth(prs)
    slide_analysis(prs)
    slide_db_deploy(prs)
    slide_ui(prs)
    slide_conclusion(prs)

    out = (r"C:\Users\조립\Desktop\claude"
           r"\Material Shortage Status vs. Production Plan"
           r"\자재부족현황_앱개발_완료보고서.pptx")
    prs.save(out)
    import sys
    sys.stdout.buffer.write(("저장 완료: " + out + "\n").encode("utf-8"))

if __name__ == '__main__':
    main()
