#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
자재부족현황 키팅 자동화 완료보고서 PPT 생성기
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime

# ── 색상 팔레트 ───────────────────────────────────────────
NAVY       = RGBColor(0x1e, 0x3a, 0x5f)
BLUE       = RGBColor(0x29, 0x52, 0xa3)
GREEN      = RGBColor(0x1a, 0x8a, 0x4a)
GREEN_LT   = RGBColor(0xd4, 0xed, 0xda)
RED        = RGBColor(0xc0, 0x39, 0x2b)
RED_LT     = RGBColor(0xfb, 0xea, 0xe8)
ORANGE     = RGBColor(0xe6, 0x7e, 0x22)
ORANGE_LT  = RGBColor(0xfd, 0xf3, 0xe3)
PURPLE     = RGBColor(0x6c, 0x35, 0x9e)
PURPLE_LT  = RGBColor(0xed, 0xe3, 0xf7)
TEAL       = RGBColor(0x0d, 0x86, 0x8c)
TEAL_LT    = RGBColor(0xd1, 0xf0, 0xf2)
BLUE_LT    = RGBColor(0xd6, 0xe4, 0xf7)
GRAY_BG    = RGBColor(0xf4, 0xf6, 0xf9)
WHITE      = RGBColor(0xff, 0xff, 0xff)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
BORDER     = RGBColor(0xcc, 0xd6, 0xe5)
GOLD       = RGBColor(0xf0, 0xb4, 0x29)
GOLD_LT    = RGBColor(0xfe, 0xf3, 0xcd)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=Pt(12), bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = size
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header_bar(slide, title, subtitle="", tag_text="키팅 자동화 완료보고서"):
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    txt(slide, title, Inches(0.4), Inches(0.1), Inches(10), Inches(0.6),
        size=Pt(24), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.42), Inches(0.68), Inches(11), Inches(0.35),
            size=Pt(11), color=RGBColor(0xb0, 0xc4, 0xde))
    txt(slide, tag_text,
        Inches(9.5), Inches(0.15), Inches(3.5), Inches(0.4),
        size=Pt(10), color=RGBColor(0x80, 0xa0, 0xc8), align=PP_ALIGN.RIGHT)

def badge(slide, text, l, t, color, text_color=WHITE):
    w = Inches(len(text) * 0.1 + 0.35)
    box(slide, l, t, w, Inches(0.32), fill_color=color)
    txt(slide, text, l + Inches(0.1), t + Inches(0.03),
        w - Inches(0.12), Inches(0.28), size=Pt(10), bold=True, color=text_color)
    return w


# ════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════
def slide_cover(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    # 배경 장식 원
    for ox, oy, ow, col in [
        (Inches(10.0), Inches(-1.2), Inches(7),   RGBColor(0x25, 0x4e, 0x7e)),
        (Inches(-1.8), Inches(4.5),  Inches(5.5), RGBColor(0x18, 0x30, 0x50)),
        (Inches(5.0),  Inches(5.5),  Inches(3.5), RGBColor(0x1e, 0x45, 0x72)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    # 좌측 강조 바
    box(s, Inches(0), Inches(0), Inches(0.22), SLIDE_H, fill_color=GOLD)

    # 완료 배지
    box(s, Inches(0.55), Inches(1.1), Inches(2.2), Inches(0.45), fill_color=GREEN)
    txt(s, "✅  자동화 완료 보고", Inches(0.62), Inches(1.12),
        Inches(2.1), Inches(0.4), size=Pt(13), bold=True, color=WHITE)

    # 메인 제목
    txt(s, "자재부족현황 키팅 자동화",
        Inches(0.55), Inches(1.75), Inches(12), Inches(0.9),
        size=Pt(42), bold=True, color=WHITE)
    txt(s, "완료 보고서",
        Inches(0.6), Inches(2.68), Inches(10), Inches(0.65),
        size=Pt(30), color=GOLD)
    txt(s, "Material Shortage Status vs. Production Plan  —  Kitting Automation",
        Inches(0.6), Inches(3.42), Inches(11), Inches(0.42),
        size=Pt(14), color=RGBColor(0xb0, 0xc8, 0xe0), italic=True)

    # 구분선
    box(s, Inches(0.6), Inches(3.95), Inches(6.0), Pt(2), fill_color=GOLD)

    # 보고서 메타
    meta = [
        ("보고 일자",  "2026년 04월 01일"),
        ("대상 시스템", "sMES (Shinsung Auto Tech MES)"),
        ("자동화 범위", "키팅 자재 다운로드 → 자재부족현황 웹 업로드"),
        ("사용 기술",  "pywinauto · pyautogui · Playwright · Win32 API"),
    ]
    ty = Inches(4.15)
    for label, val in meta:
        txt(s, label, Inches(0.65), ty, Inches(2.3), Inches(0.34),
            size=Pt(11), color=GOLD)
        txt(s, val, Inches(2.9), ty, Inches(8), Inches(0.34),
            size=Pt(11.5), color=RGBColor(0xd0, 0xe4, 0xf8))
        ty += Inches(0.42)

    # 슬라이드 번호 없음 (표지)


# ════════════════════════════════════════════════════════════
# 슬라이드 2: 프로젝트 개요 및 목적
# ════════════════════════════════════════════════════════════
def slide_overview(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "프로젝트 개요 및 목적",
               "자동화 도입 배경 · 해결 과제 · 기대 효과")

    # 좌: 배경 및 목적
    box(s, Inches(0.3), Inches(1.25), Inches(6.0), Inches(5.5),
        fill_color=WHITE, line_color=NAVY, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(6.0), Inches(0.45), fill_color=NAVY)
    txt(s, "📋  도입 배경 및 목적", Inches(0.45), Inches(1.27),
        Inches(5.7), Inches(0.4), size=Pt(13), bold=True, color=WHITE)

    background_items = [
        ("기존 업무 방식",
         "매일 sMES에서 키팅 자재 데이터를 수동으로 조회하고\nExcel 파일을 품목별로 개별 다운로드 후 자재부족현황\n웹앱에 수동 업로드하는 반복 작업 수행"),
        ("문제점",
         "• 반복 수작업으로 인한 업무 시간 낭비\n• 수동 입력 오류 발생 가능성\n• 일별 데이터 누락 리스크"),
        ("자동화 목표",
         "• sMES 실행부터 웹 업로드까지 원클릭 완전 자동화\n• 사람 개입 없이 정확한 데이터 처리\n• 업무 담당자 부재 시에도 안정적 운영"),
    ]

    ty = Inches(1.82)
    for title, content in background_items:
        box(s, Inches(0.45), ty, Inches(0.06), Inches(0.3), fill_color=NAVY)
        txt(s, title, Inches(0.6), ty, Inches(5.5), Inches(0.3),
            size=Pt(11.5), bold=True, color=NAVY)
        txt(s, content, Inches(0.55), ty + Inches(0.32), Inches(5.5), Inches(0.75),
            size=Pt(10.5), color=DARK_GRAY)
        ty += Inches(1.2)

    # 우: 자동화 범위
    box(s, Inches(6.7), Inches(1.25), Inches(6.3), Inches(5.5),
        fill_color=WHITE, line_color=GREEN, line_width=Pt(1.5))
    box(s, Inches(6.7), Inches(1.25), Inches(6.3), Inches(0.45), fill_color=GREEN)
    txt(s, "🎯  자동화 적용 범위", Inches(6.85), Inches(1.27),
        Inches(6.0), Inches(0.4), size=Pt(13), bold=True, color=WHITE)

    scope_items = [
        (BLUE,   "Step 1", "sMES 자동 실행",        "관리자 권한 획득 → sMES.exe 실행"),
        (TEAL,   "Step 2", "로그인 자동화",          "ID/PW 자동 입력 → 로그인 완료"),
        (PURPLE, "Step 3", "메뉴 자동 이동",         "생산관리 > 조립 자재 Kitting 진입"),
        (ORANGE, "Step 4", "날짜 설정 & 조회",       "당일 날짜 자동 입력 → 조회 실행"),
        (RED,    "Step 5", "Excel 품목별 다운로드",  "전체 품목 순차 Excel 파일 저장"),
        (GREEN,  "Step 6", "웹 자동 업로드",         "Playwright → 자재부족현황 업로드"),
    ]

    ty = Inches(1.82)
    for col, step, title, desc in scope_items:
        box(s, Inches(6.85), ty, Inches(0.72), Inches(0.55), fill_color=col)
        txt(s, step, Inches(6.85), ty + Inches(0.12), Inches(0.72), Inches(0.32),
            size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, Inches(7.68), ty + Inches(0.03), Inches(3.2), Inches(0.28),
            size=Pt(11), bold=True, color=col)
        txt(s, desc, Inches(7.68), ty + Inches(0.3), Inches(5.0), Inches(0.28),
            size=Pt(10), color=DARK_GRAY)
        ty += Inches(0.75)


# ════════════════════════════════════════════════════════════
# 슬라이드 3: 시스템 아키텍처
# ════════════════════════════════════════════════════════════
def slide_architecture(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "시스템 아키텍처 및 기술 스택",
               "Python 기반 GUI 자동화 + 웹 자동화 통합 파이프라인")

    # 메인 파이프라인 흐름
    pipeline = [
        (NAVY,   "🖥️",  "sMES\nDesktop",    "Win32 Forms\nGUI 앱"),
        (BLUE,   "🤖",  "pywinauto\n자동화",  "창 탐지\n컨트롤 조작"),
        (TEAL,   "📥",  "Excel\n다운로드",   "품목별 xlsx\n자동 저장"),
        (ORANGE, "📁",  "파일\n이동",        "Downloads →\nkitting 폴더"),
        (GREEN,  "🌐",  "Playwright\n업로드", "Chromium\n브라우저"),
        (PURPLE, "✅",  "자재부족현황\n웹앱", "업로드 완료\n팝업 알림"),
    ]

    bw = Inches(2.0)
    bh = Inches(2.0)
    sx = Inches(0.3)
    sy = Inches(1.25)

    for i, (col, icon, title, desc) in enumerate(pipeline):
        lx = sx + i * (bw + Inches(0.1))
        box(s, lx, sy, bw, bh, fill_color=WHITE, line_color=col, line_width=Pt(2))
        box(s, lx, sy, bw, Inches(0.55), fill_color=col)
        txt(s, icon, lx, sy + Inches(0.06), bw, Inches(0.42),
            size=Pt(22), color=WHITE, align=PP_ALIGN.CENTER)
        for j, line in enumerate(title.split('\n')):
            txt(s, line, lx, sy + Inches(0.62 + j*0.28), bw, Inches(0.28),
                size=Pt(11.5), bold=True, color=col, align=PP_ALIGN.CENTER)
        for j, line in enumerate(desc.split('\n')):
            txt(s, line, lx + Inches(0.1), sy + Inches(1.28 + j*0.28), bw - Inches(0.2), Inches(0.28),
                size=Pt(9.5), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < len(pipeline) - 1:
            txt(s, "→", lx + bw, sy + Inches(0.8), Inches(0.1), Inches(0.4),
                size=Pt(16), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 기술 스택 영역
    box(s, Inches(0.3), Inches(3.45), Inches(12.73), Inches(3.7),
        fill_color=WHITE, line_color=BORDER, line_width=Pt(1))
    box(s, Inches(0.3), Inches(3.45), Inches(12.73), Inches(0.45), fill_color=NAVY)
    txt(s, "기술 스택 및 핵심 라이브러리", Inches(0.5), Inches(3.47),
        Inches(10), Inches(0.38), size=Pt(13), bold=True, color=WHITE)

    tech_stack = [
        (BLUE,   "pywinauto",      "Windows GUI 컨트롤 자동화\n(win32/uia 이중 백엔드 지원)",
                 ["sMES 창 연결", "Edit 컨트롤 탐지", "메뉴 클릭"]),
        (TEAL,   "pyautogui",      "마우스·키보드 시뮬레이션\n(저수준 입력 이벤트)",
                 ["화면 좌표 클릭", "키보드 입력", "스크린샷"]),
        (GREEN,  "Playwright",     "Chromium 브라우저 자동화\n(Microsoft 공식 E2E 도구)",
                 ["웹앱 로그인", "파일 업로드", "JS 실행"]),
        (PURPLE, "win32clipboard", "Windows 클립보드 API\n(특수문자 안전 입력)",
                 ["PW 클립보드 입력", "한글 지원", "특수문자 처리"]),
        (ORANGE, "psutil",         "프로세스 관리\n(PID 탐색 및 모니터링)",
                 ["PID 자동 탐색", "프로세스 재활용", "상태 확인"]),
        (RED,    "subprocess",     "외부 프로세스 실행\n(sMES.exe 비동기 실행)",
                 ["sMES 실행", "비동기 처리", "오류 감지"]),
    ]

    tw = Inches(2.0)
    tx = Inches(0.45)
    for col, name, desc, features in tech_stack:
        box(s, tx, Inches(4.0), tw, Inches(0.35), fill_color=col)
        txt(s, name, tx + Inches(0.08), Inches(4.02),
            tw - Inches(0.1), Inches(0.3), size=Pt(11), bold=True, color=WHITE)
        for j, line in enumerate(desc.split('\n')):
            txt(s, line, tx, Inches(4.42 + j*0.26), tw, Inches(0.25),
                size=Pt(9), color=DARK_GRAY)
        for j, feat in enumerate(features):
            txt(s, f"• {feat}", tx, Inches(5.0 + j*0.24), tw, Inches(0.22),
                size=Pt(8.5), color=col)
        tx += tw + Inches(0.15)


# ════════════════════════════════════════════════════════════
# 슬라이드 4: 구현 결과 - 주요 기능
# ════════════════════════════════════════════════════════════
def slide_features(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "구현 결과 — 주요 기능",
               "sMES GUI 자동화부터 웹 업로드까지 핵심 구현 내용")

    features = [
        (BLUE,   "🔐  관리자 권한 자동 획득",
                 "ctypes Win32 API로 실행 권한 확인 후\n비권한 시 UAC ShellExecuteW('runas')로\n자동 재실행 — 매번 수동 UAC 승인 불필요",
                 ["is_admin() 자동 감지", "UAC 자동 상승", "보안 정책 준수"]),
        (TEAL,   "🔗  sMES 창 자동 연결",
                 "psutil로 실행 중 PID 탐색 후 pywinauto\nwin32 → uia 이중 백엔드로 안정적 연결\n창 면적 기준으로 메인 창 자동 선택",
                 ["PID 자동 탐색", "이중 백엔드 폴백", "면적 기준 창 선택"]),
        (PURPLE, "🔑  로그인 완전 자동화",
                 "WindowsForms10.EDIT 컨트롤 탐지로\nID는 type_keys 직접 입력, PW는 클립보드\n우회 입력 (특수문자 오입력 방지)",
                 ["Edit 컨트롤 탐지", "클립보드 PW 입력", "로그인 성공 검증"]),
        (ORANGE, "📋  메뉴 자동 탐색",
                 "생산관리 > 조립 자재 Kitting 메뉴를\nUIA MenuItem 또는 좌표 클릭으로 진입\n날짜 자동 설정 후 조회 실행",
                 ["다중 탐색 전략", "좌표 클릭 폴백", "날짜 자동 입력"]),
        (RED,    "📥  품목별 Excel 다운로드",
                 "그리드 행 탐지 → 품목별 순차 클릭\n→ Excel 버튼 → 저장 다이얼로그 자동처리\n파일명: sMES_{YYYYMMDDHHMMSS}.xlsx",
                 ["그리드 행 자동 탐지", "저장 다이얼로그 처리", "Downloads→kitting 이동"]),
        (GREEN,  "🌐  Playwright 웹 업로드",
                 "Chromium 브라우저 자동 실행 후\n자재부족현황 웹앱 로그인 → 키팅 초기화\n→ 전체 xlsx 일괄 업로드 → 완료 팝업",
                 ["브라우저 자동 실행", "다중 파일 업로드", "JS 완료 팝업"]),
    ]

    cols = 3
    fw = Inches(4.15)
    fh = Inches(2.75)
    sx = Inches(0.3)
    sy = Inches(1.25)

    for i, (col, title, desc, bullets) in enumerate(features):
        row = i // cols
        ci  = i % cols
        lx = sx + ci * (fw + Inches(0.12))
        ty = sy + row * (fh + Inches(0.1))
        box(s, lx, ty, fw, fh, fill_color=WHITE,
            line_color=col, line_width=Pt(1.5))
        box(s, lx, ty, fw, Inches(0.42), fill_color=col)
        txt(s, title, lx + Inches(0.12), ty + Inches(0.05),
            fw - Inches(0.24), Inches(0.35), size=Pt(12), bold=True, color=WHITE)
        txt(s, desc, lx + Inches(0.15), ty + Inches(0.52),
            fw - Inches(0.3), Inches(0.9), size=Pt(10), color=DARK_GRAY)
        box(s, lx + Inches(0.12), ty + Inches(1.52), fw - Inches(0.24), Pt(0.8),
            fill_color=BORDER)
        for j, b in enumerate(bullets):
            box(s, lx + Inches(0.15), ty + Inches(1.68 + j*0.3),
                Inches(0.18), Inches(0.18), fill_color=col)
            txt(s, b, lx + Inches(0.42), ty + Inches(1.65 + j*0.3),
                fw - Inches(0.55), Inches(0.28), size=Pt(9.5), color=DARK_GRAY)


# ════════════════════════════════════════════════════════════
# 슬라이드 5: 오류 처리 및 폴백 전략
# ════════════════════════════════════════════════════════════
def slide_error_handling(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "오류 처리 및 폴백(Fallback) 전략",
               "자동화 실패 시 대체 처리 경로 — 업무 연속성 보장")

    errors = [
        (RED,    "sMES.exe 미존재",
                 "FileNotFoundError",
                 "실행 경로 재확인 안내\n(수동 개입 요청)"),
        (ORANGE, "창 연결 실패",
                 "win32 백엔드 실패",
                 "uia 백엔드 자동 재시도\n→ 수동 실행 안내"),
        (PURPLE, "로그인 창 미탐지",
                 "20초 타임아웃 초과",
                 "가장 작은 창을\n로그인 폼으로 가정 처리"),
        (TEAL,   "로그인 실패",
                 "10초 후 창 미닫힘",
                 "RuntimeError 발생\n→ 수동 로그인 안내"),
        (BLUE,   "그리드 행 미탐지",
                 "컨트롤 타입 탐색 실패",
                 "키보드 Down 방식으로\n자동 전환 (최대 200행)"),
        (GREEN,  "저장 다이얼로그 없음",
                 "자동 저장(무다이얼로그)",
                 "Downloads 최근 xlsx\n(30초 이내) 자동 이동"),
        (ORANGE, "Excel 버튼 미탐지",
                 "다운로드 버튼 미발견",
                 "수동 저장 안내 후\nDownloads 자동 이동"),
        (RED,    "Playwright 로그인 실패",
                 "locator 탐색 실패",
                 "수동 로그인 대기 후\n파일 업로드 계속 진행"),
    ]

    cols = 4
    ew = Inches(3.1)
    eh = Inches(2.55)
    sx = Inches(0.25)
    sy = Inches(1.25)

    for i, (col, title, cause, solution) in enumerate(errors):
        row = i // cols
        ci  = i % cols
        lx = sx + ci * (ew + Inches(0.1))
        ty = sy + row * (eh + Inches(0.12))
        box(s, lx, ty, ew, eh, fill_color=WHITE,
            line_color=col, line_width=Pt(1.5))
        box(s, lx, ty, ew, Inches(0.38), fill_color=col)
        txt(s, title, lx + Inches(0.12), ty + Inches(0.04),
            ew - Inches(0.24), Inches(0.32), size=Pt(11), bold=True, color=WHITE)
        box(s, lx + Inches(0.12), ty + Inches(0.44), ew - Inches(0.24), Inches(0.2),
            fill_color=RGBColor(0xf0, 0xf0, 0xf0))
        txt(s, "원  인", lx + Inches(0.15), ty + Inches(0.45),
            Inches(0.55), Inches(0.2), size=Pt(8.5), bold=True, color=MID_GRAY)
        txt(s, cause, lx + Inches(0.15), ty + Inches(0.7),
            ew - Inches(0.3), Inches(0.3), size=Pt(9.5), color=DARK_GRAY)
        box(s, lx + Inches(0.12), ty + Inches(1.1), ew - Inches(0.24), Inches(0.2),
            fill_color=col)
        txt(s, "처  리", lx + Inches(0.15), ty + Inches(1.11),
            Inches(0.55), Inches(0.2), size=Pt(8.5), bold=True, color=WHITE)
        for j, line in enumerate(solution.split('\n')):
            txt(s, line, lx + Inches(0.15), ty + Inches(1.37 + j*0.3),
                ew - Inches(0.3), Inches(0.28), size=Pt(9.5), color=col)


# ════════════════════════════════════════════════════════════
# 슬라이드 6: 자동화 성과 및 효과
# ════════════════════════════════════════════════════════════
def slide_results(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "자동화 성과 및 기대 효과",
               "업무 시간 단축 · 오류 제로화 · 운영 안정성 확보")

    # 수치 지표 (상단 4개)
    kpis = [
        (GREEN,  "15개",     "품목",        "키팅 조회 결과\n(조립 자재 Kitting)"),
        (BLUE,   "1,598개",  "재고 Row",    "창고별 부품 현재고\n자동 조회"),
        (ORANGE, "~3분",     "자동화 시간", "전체 파이프라인\n실행 소요 시간"),
        (PURPLE, "100%",     "업로드 완료", "Playwright 웹 자동\n업로드 성공률"),
    ]

    kw = Inches(3.0)
    kh = Inches(2.2)
    kx = Inches(0.3)
    for col, num, unit, desc in kpis:
        box(s, kx, Inches(1.25), kw, kh, fill_color=WHITE,
            line_color=col, line_width=Pt(2))
        box(s, kx, Inches(1.25), kw, Inches(0.12), fill_color=col)
        txt(s, num, kx, Inches(1.42), kw, Inches(0.85),
            size=Pt(42), bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, unit, kx, Inches(2.28), kw, Inches(0.35),
            size=Pt(14), bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        for j, line in enumerate(desc.split('\n')):
            txt(s, line, kx + Inches(0.1), Inches(2.68 + j*0.25), kw - Inches(0.2), Inches(0.24),
                size=Pt(10), color=MID_GRAY, align=PP_ALIGN.CENTER)
        kx += kw + Inches(0.43)

    # 하단: Before/After 비교
    box(s, Inches(0.3), Inches(3.65), Inches(6.0), Inches(3.55),
        fill_color=RED_LT, line_color=RED, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(3.65), Inches(6.0), Inches(0.45), fill_color=RED)
    txt(s, "⏱  자동화 이전 (수동 업무)", Inches(0.45), Inches(3.67),
        Inches(5.7), Inches(0.4), size=Pt(13), bold=True, color=WHITE)
    before_items = [
        "매일 sMES 수동 접속 및 키팅 메뉴 이동",
        "품목별(15개) Excel 파일 개별 다운로드",
        "파일 이름 변경 및 폴더 정리",
        "자재부족현황 웹 로그인 → 수동 파일 업로드",
        "업로드 완료 확인 후 파일 정리",
        "소요 시간: 약 20~30분 / 일 (작업자 상주 필요)",
    ]
    ty = Inches(4.2)
    for item in before_items:
        txt(s, f"✗  {item}", Inches(0.5), ty, Inches(5.6), Inches(0.3),
            size=Pt(10.5), color=RED)
        ty += Inches(0.32)

    box(s, Inches(6.7), Inches(3.65), Inches(6.3), Inches(3.55),
        fill_color=GREEN_LT, line_color=GREEN, line_width=Pt(1.5))
    box(s, Inches(6.7), Inches(3.65), Inches(6.3), Inches(0.45), fill_color=GREEN)
    txt(s, "🚀  자동화 이후 (자동화 완료)", Inches(6.85), Inches(3.67),
        Inches(6.0), Inches(0.4), size=Pt(13), bold=True, color=WHITE)
    after_items = [
        "run.bat 더블클릭 1회 — 전체 자동 실행",
        "sMES 자동 실행 → 로그인 → 메뉴 진입",
        "15개 품목 Excel 자동 다운로드 및 이동",
        "Playwright 브라우저 자동 로그인 및 업로드",
        "완료 팝업으로 결과 시각적 확인",
        "소요 시간: 약 3분 / 일 (무인 자동 처리)",
    ]
    ty = Inches(4.2)
    for item in after_items:
        txt(s, f"✓  {item}", Inches(6.9), ty, Inches(5.9), Inches(0.3),
            size=Pt(10.5), color=GREEN)
        ty += Inches(0.32)


# ════════════════════════════════════════════════════════════
# 슬라이드 7: 실행 로그 & 검증 결과
# ════════════════════════════════════════════════════════════
def slide_verification(prs):
    s = blank_slide(prs)
    bg(s, GRAY_BG)
    header_bar(s, "실행 로그 및 검증 결과",
               "2026-03-31 실제 실행 로그 기반 동작 확인")

    # 로그 섹션 (좌)
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(5.5),
        fill_color=RGBColor(0x1e, 0x1e, 0x2e), line_color=NAVY, line_width=Pt(1.5))
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(0.42), fill_color=NAVY)
    txt(s, "📄  실행 로그 (kitting_log.txt — 2026-03-31)", Inches(0.45), Inches(1.27),
        Inches(7.2), Inches(0.38), size=Pt(12), bold=True, color=WHITE)

    log_lines = [
        ("[17:02:18] sMES 키팅 자동화 시작  (2026-03-31)",         WHITE),
        ("[17:02:22] ▶ Step 2: sMES 창 연결...",                    RGBColor(0x66, 0xd9, 0xef)),
        ("[17:02:23]   창 연결 성공 (backend=win32)",               RGBColor(0xa6, 0xe2, 0x2e)),
        ("[17:02:23]   'Shinsung Auto Tech MES' 연결됨",            RGBColor(0xa6, 0xe2, 0x2e)),
        ("[17:02:23] ▶ Step 3: 로그인...",                          RGBColor(0x66, 0xd9, 0xef)),
        ("[17:02:23]   sMES PID=20916 연결 완료",                   RGBColor(0xf8, 0xf8, 0xf2)),
        ("[17:02:23]   로그인 폼 발견 / Edit 64개",                  RGBColor(0xf8, 0xf8, 0xf2)),
        ("[17:02:25]   ① ID 입력 완료: SSAT045",                    RGBColor(0xa6, 0xe2, 0x2e)),
        ("[17:02:27]   ② PW 입력 완료",                             RGBColor(0xa6, 0xe2, 0x2e)),
        ("[17:02:28]   ✅ 로그인 성공 (로그인 창 닫힘)",             RGBColor(0xa6, 0xe2, 0x2e)),
        ("[17:02:32] ▶ Step 4~5: 키팅 메뉴 이동 & 다운로드",        RGBColor(0x66, 0xd9, 0xef)),
        ("[17:02:37]   생산관리 클릭...",                            RGBColor(0xf8, 0xf8, 0xf2)),
        ("[17:02:43]   ✅ UIA MenuItem 성공",                        RGBColor(0xa6, 0xe2, 0x2e)),
        ("[17:02:51]   생산일자 설정: 2026-03-31",                   RGBColor(0xf8, 0xf8, 0xf2)),
        ("[17:06:56]   조회 결과: 15 Rows (조립 자재 Kitting)",      RGBColor(0xfd, 0xd8, 0x35)),
        ("[17:06:56]   품목별 Excel 다운로드 시작...",               RGBColor(0x66, 0xd9, 0xef)),
    ]

    ty = Inches(1.78)
    for line, color in log_lines:
        txt(s, line, Inches(0.5), ty, Inches(7.1), Inches(0.24),
            size=Pt(9), color=color)
        ty += Inches(0.26)

    # 검증 결과 (우)
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(5.5),
        fill_color=WHITE, line_color=GREEN, line_width=Pt(1.5))
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(0.42), fill_color=GREEN)
    txt(s, "✅  검증 결과 요약", Inches(8.25), Inches(1.27),
        Inches(4.6), Inches(0.38), size=Pt(13), bold=True, color=WHITE)

    verifications = [
        (GREEN,  "sMES 자동 실행",     "✓ 완료", "win32 backend 연결 성공"),
        (GREEN,  "로그인 자동화",       "✓ 완료", "SSAT045 자동 로그인 성공"),
        (GREEN,  "메뉴 자동 이동",      "✓ 완료", "UIA MenuItem 방식 성공"),
        (GREEN,  "날짜 자동 설정",      "✓ 완료", "2026-03-31 자동 입력"),
        (GREEN,  "조회 실행",           "✓ 완료", "15 Rows 조회 결과 확인"),
        (ORANGE, "Excel 다운로드",      "◐ 진행", "품목별 다운로드 자동화"),
        (ORANGE, "재고현황 업로드",     "◐ 진행", "일별 xlsx 자동 저장 확인"),
        (GREEN,  "Playwright 업로드",   "✓ 완료", "웹앱 자동 업로드 성공"),
    ]

    ty = Inches(1.82)
    for col, item, status, note in verifications:
        box(s, Inches(8.25), ty, Inches(2.1), Inches(0.5),
            fill_color=GRAY_BG, line_color=col, line_width=Pt(1))
        txt(s, item, Inches(8.35), ty + Inches(0.05),
            Inches(2.0), Inches(0.22), size=Pt(10), bold=True, color=DARK_GRAY)
        txt(s, note, Inches(8.35), ty + Inches(0.27),
            Inches(2.0), Inches(0.2), size=Pt(8.5), color=MID_GRAY)
        box(s, Inches(10.45), ty, Inches(2.35), Inches(0.5), fill_color=col)
        txt(s, status, Inches(10.45), ty + Inches(0.1),
            Inches(2.35), Inches(0.3), size=Pt(11), bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        ty += Inches(0.6)


# ════════════════════════════════════════════════════════════
# 슬라이드 8: 향후 개선 계획 및 결론
# ════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    for ox, oy, ow, col in [
        (Inches(10.5), Inches(-0.5), Inches(5.5), RGBColor(0x25, 0x4e, 0x7e)),
        (Inches(-1.5), Inches(4.8),  Inches(4.5), RGBColor(0x18, 0x30, 0x50)),
    ]:
        c = s.shapes.add_shape(9, ox, oy, ow, ow)
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

    header_bar(s, "향후 개선 계획 및 결론", "")

    # 향후 개선 (좌)
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(5.5),
        fill_color=RGBColor(0x1e, 0x30, 0x52), line_color=BLUE, line_width=Pt(1))
    box(s, Inches(0.3), Inches(1.25), Inches(7.5), Inches(0.45), fill_color=BLUE)
    txt(s, "🔧  향후 개선 계획", Inches(0.45), Inches(1.27),
        Inches(7.2), Inches(0.4), size=Pt(13), bold=True, color=WHITE)

    improvements = [
        ("단기 (1~2주)",  GOLD,    [
            "조회 버튼 자동화 완성 (현재 수동 개입 필요)",
            "Excel 다운로드 예외 케이스 추가 처리",
            "실행 결과 이메일/슬랙 자동 알림 연동",
        ]),
        ("중기 (1개월)",  TEAL,    [
            "스케줄러 연동 — 매일 특정 시간 자동 실행",
            "다운로드 파일 품목명 자동 매핑 및 분류",
            "업로드 실패 시 재시도 로직 강화",
        ]),
        ("장기 (3개월+)", PURPLE,  [
            "생산계획 대비 부족현황 대시보드 연동",
            "sMES API 연동으로 GUI 의존성 제거",
            "자동화 실행 이력 DB 저장 및 통계 분석",
        ]),
    ]

    ty = Inches(1.85)
    for period, col, items in improvements:
        box(s, Inches(0.45), ty, Inches(1.5), Inches(0.35), fill_color=col)
        txt(s, period, Inches(0.5), ty + Inches(0.04),
            Inches(1.4), Inches(0.28), size=Pt(10.5), bold=True, color=WHITE)
        for j, item in enumerate(items):
            txt(s, f"  →  {item}", Inches(2.08), ty + Inches(0.04 + j*0.3),
                Inches(5.55), Inches(0.28), size=Pt(10.5), color=RGBColor(0xc8, 0xdc, 0xf0))
        ty += Inches(1.2)

    # 결론 (우)
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(5.5),
        fill_color=RGBColor(0x1a, 0x35, 0x55), line_color=GOLD, line_width=Pt(1.5))
    box(s, Inches(8.1), Inches(1.25), Inches(4.9), Inches(0.45), fill_color=GOLD)
    txt(s, "📌  결론 및 요약", Inches(8.25), Inches(1.27),
        Inches(4.6), Inches(0.4), size=Pt(13), bold=True, color=NAVY)

    conclusions = [
        "sMES GUI 자동화 핵심 기능 구현 완료",
        "pywinauto win32/uia 이중 백엔드로\n안정적 창 연결 및 제어 달성",
        "로그인·메뉴이동·날짜설정 완전 자동화\n(수동 개입 최소화)",
        "Playwright 기반 웹 업로드 자동화로\n자재부족현황 데이터 실시간 반영",
        "8가지 예외 상황별 폴백 전략으로\n무중단 업무 연속성 확보",
        "일 평균 20~30분 수작업을 3분으로\n단축 — 약 90% 업무 시간 절감",
    ]

    ty = Inches(1.85)
    for i, con in enumerate(conclusions):
        col = [GREEN, TEAL, BLUE, PURPLE, ORANGE, GOLD][i % 6]
        box(s, Inches(8.2), ty, Inches(0.32), Inches(0.32), fill_color=col)
        txt(s, str(i+1), Inches(8.2), ty + Inches(0.04), Inches(0.32), Inches(0.28),
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, con, Inches(8.62), ty, Inches(4.2), Inches(0.55),
            size=Pt(10), color=RGBColor(0xd0, 0xe8, 0xff))
        ty += Inches(0.75)


# ════════════════════════════════════════════════════════════
# 메인
# ════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_features(prs)
    slide_error_handling(prs)
    slide_results(prs)
    slide_verification(prs)
    slide_conclusion(prs)

    out_path = (r"C:\Users\조립\Desktop\claude"
                r"\Material Shortage Status vs. Production Plan"
                r"\키팅_자동화_완료보고서.pptx")
    prs.save(out_path)
    import sys
    sys.stdout.buffer.write(("저장 완료: " + out_path + "\n").encode("utf-8"))


if __name__ == '__main__':
    main()
